from __future__ import annotations

import tempfile
import unittest
import zipfile
from pathlib import Path

from office_artifact_runtime import office_backend_statuses, run_openxml_sdk_validator


class OpenXmlSdkValidatorTests(unittest.TestCase):
    def setUp(self) -> None:
        status = {item["id"]: item for item in office_backend_statuses()}["openxml-sdk"]
        if status["status"] != "ready":
            self.skipTest(status.get("detail", "Open XML SDK validator is unavailable"))

    def test_official_schema_validator_accepts_minimal_three_format_packages(self) -> None:
        import docx
        from pptx import Presentation

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            docx_path = root / "sample.docx"
            document = docx.Document()
            document.add_paragraph("Schema-valid Word")
            document.save(docx_path)
            xlsx_path = root / "sample.xlsx"
            with zipfile.ZipFile(xlsx_path, "w") as archive:
                archive.writestr("[Content_Types].xml", b'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types"><Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/><Default Extension="xml" ContentType="application/xml"/><Override PartName="/xl/workbook.xml" ContentType="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet.main+xml"/><Override PartName="/xl/worksheets/sheet1.xml" ContentType="application/vnd.openxmlformats-officedocument.spreadsheetml.worksheet+xml"/></Types>''')
                archive.writestr("_rels/.rels", b'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="xl/workbook.xml"/></Relationships>''')
                archive.writestr("xl/workbook.xml", b'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<workbook xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"><sheets><sheet name="Sheet1" sheetId="1" r:id="rId1"/></sheets></workbook>''')
                archive.writestr("xl/_rels/workbook.xml.rels", b'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/worksheet" Target="worksheets/sheet1.xml"/></Relationships>''')
                archive.writestr("xl/worksheets/sheet1.xml", b'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<worksheet xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main"><sheetData><row r="1"><c r="A1" t="inlineStr"><is><t>Schema-valid Excel</t></is></c></row></sheetData></worksheet>''')
            pptx_path = root / "sample.pptx"
            presentation = Presentation()
            presentation.slides.add_slide(presentation.slide_layouts[6])
            presentation.save(pptx_path)

            for path in (docx_path, xlsx_path, pptx_path):
                evidence = run_openxml_sdk_validator(path)
                self.assertEqual("pass", evidence["status"], evidence)
                self.assertEqual("3.5.1", evidence["engineVersion"])
                self.assertEqual("Microsoft365", evidence["fileFormatVersion"])


if __name__ == "__main__":
    unittest.main()
