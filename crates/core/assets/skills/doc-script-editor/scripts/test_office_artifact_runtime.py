from __future__ import annotations

import argparse
import contextlib
import hashlib
import io
import json
import os
import tempfile
import types
import unittest
import zipfile
from pathlib import Path
from unittest import mock
from xml.etree import ElementTree as ET

import edit_doc
import office_artifact_service
from office_artifact_runtime import (
    publish_staged_artifact,
    scan_ooxml_risks,
    staging_path,
    validate_ooxml_package,
)
from office_artifact_service import OfficeArtifactJob, execute_job

CONTENT_TYPES_NS = "http://schemas.openxmlformats.org/package/2006/content-types"
RELATIONSHIPS_NS = "http://schemas.openxmlformats.org/package/2006/relationships"


def _rewrite_zip(
    path: Path,
    replacements: dict[str, bytes],
    additions: dict[str, bytes] | None = None,
    output: Path | None = None,
) -> Path:
    rewritten = output or path.with_name(f"{path.stem}-rewritten{path.suffix}")
    with zipfile.ZipFile(path) as source, zipfile.ZipFile(rewritten, "w") as destination:
        for info in source.infolist():
            destination.writestr(info, replacements.get(info.filename, source.read(info.filename)))
        for name, data in (additions or {}).items():
            destination.writestr(name, data)
    return rewritten


class OfficeArtifactRuntimeTests(unittest.TestCase):
    def setUp(self) -> None:
        self.temp = tempfile.TemporaryDirectory()
        self.root = Path(self.temp.name).resolve()
        self.previous_cwd = Path.cwd()
        os.chdir(self.root)

    def tearDown(self) -> None:
        os.chdir(self.previous_cwd)
        self.temp.cleanup()

    def test_docx_creation_emits_real_hyperlink_and_valid_package(self) -> None:
        path = self.root / "linked.docx"
        args = argparse.Namespace(
            path=str(path),
            template=None,
            font="Calibri",
            title="Linked report",
            subtitle="",
            input_md=None,
            body="Read [Nexa](https://example.com/nexa) for details.",
            footer="",
            author="Nexa",
        )
        with contextlib.redirect_stdout(io.StringIO()):
            self.assertEqual(0, edit_doc.cmd_create_docx(args))

        with zipfile.ZipFile(path) as archive:
            document_xml = archive.read("word/document.xml").decode("utf-8")
            relationships = archive.read("word/_rels/document.xml.rels").decode("utf-8")
        self.assertIn("<w:hyperlink", document_xml)
        self.assertIn("https://example.com/nexa", relationships)
        self.assertEqual("pass", validate_ooxml_package(path).status)

    def test_docx_replace_matches_across_runs_and_header_story(self) -> None:
        import docx

        path = self.root / "report.docx"
        document = docx.Document()
        paragraph = document.add_paragraph()
        paragraph.add_run("Q")
        paragraph.add_run("3")
        paragraph.add_run(" Report")
        header = document.sections[0].header.paragraphs[0]
        header.add_run("Q3")
        header.add_run(" Report")
        document.save(path)

        with contextlib.redirect_stdout(io.StringIO()):
            self.assertEqual(0, edit_doc._replace_docx(path, "Q3 Report", "Q4 Review", False))

        reopened = docx.Document(path)
        self.assertEqual("Q4 Review", reopened.paragraphs[0].text)
        self.assertEqual("Q4 Review", reopened.sections[0].header.paragraphs[0].text)
        snapshots = list((self.root / ".nexa" / "doc-history").rglob("report.docx"))
        self.assertEqual(1, len(snapshots))

    def test_replace_preconditions_prevent_stale_or_ambiguous_edits(self) -> None:
        import docx

        path = self.root / "preconditions.docx"
        document = docx.Document()
        document.add_paragraph("Q3 Q3")
        document.save(path)
        before = hashlib.sha256(path.read_bytes()).hexdigest()

        with self.assertRaises(SystemExit):
            with contextlib.redirect_stderr(io.StringIO()):
                edit_doc.cmd_replace(argparse.Namespace(
                    path=str(path),
                    find="Q3",
                    replace="Q4",
                    dry_run=False,
                    expected_sha256="0" * 64,
                    expected_count=2,
                ))
        self.assertEqual(before, hashlib.sha256(path.read_bytes()).hexdigest())

        with self.assertRaises(SystemExit):
            with contextlib.redirect_stderr(io.StringIO()):
                edit_doc.cmd_replace(argparse.Namespace(
                    path=str(path),
                    find="Q3",
                    replace="Q4",
                    dry_run=False,
                    expected_sha256=before,
                    expected_count=1,
                ))
        self.assertEqual(before, hashlib.sha256(path.read_bytes()).hexdigest())

        with contextlib.redirect_stdout(io.StringIO()):
            self.assertEqual(0, edit_doc.cmd_replace(argparse.Namespace(
                path=str(path),
                find="Q3",
                replace="Q4",
                dry_run=False,
                expected_sha256=before,
                expected_count=2,
            )))
        self.assertEqual("Q4 Q4", docx.Document(path).paragraphs[0].text)

    def test_docx_replace_targets_story_scope_and_occurrence(self) -> None:
        import docx

        path = self.root / "scoped.docx"
        document = docx.Document()
        document.add_paragraph("Token Token Token")
        document.sections[0].header.paragraphs[0].text = "Token"
        document.save(path)

        with contextlib.redirect_stdout(io.StringIO()):
            self.assertEqual(0, edit_doc.cmd_replace(argparse.Namespace(
                path=str(path), find="Token", replace="Header", dry_run=False,
                expected_sha256=None, expected_count=1, scope="header",
                occurrence=1, allow_style_merge=False,
            )))
        reopened = docx.Document(path)
        self.assertEqual("Token Token Token", reopened.paragraphs[0].text)
        self.assertEqual("Header", reopened.sections[0].header.paragraphs[0].text)

        with contextlib.redirect_stdout(io.StringIO()):
            self.assertEqual(0, edit_doc.cmd_replace(argparse.Namespace(
                path=str(path), find="Token", replace="Selected", dry_run=False,
                expected_sha256=None, expected_count=3, scope="body",
                occurrence=2, allow_style_merge=False,
            )))
        self.assertEqual("Token Selected Token", docx.Document(path).paragraphs[0].text)

    def test_docx_replace_blocks_incompatible_run_style_boundary(self) -> None:
        import docx

        path = self.root / "style-boundary.docx"
        document = docx.Document()
        paragraph = document.add_paragraph()
        paragraph.add_run("Secret").bold = True
        paragraph.add_run("Value")
        document.save(path)
        before = hashlib.sha256(path.read_bytes()).hexdigest()

        with self.assertRaises(SystemExit):
            with contextlib.redirect_stderr(io.StringIO()):
                edit_doc.cmd_replace(argparse.Namespace(
                    path=str(path), find="SecretValue", replace="Safe", dry_run=False,
                    expected_sha256=None, expected_count=1, scope="body",
                    occurrence=1, allow_style_merge=False,
                ))
        self.assertEqual(before, hashlib.sha256(path.read_bytes()).hexdigest())

    def test_secure_docx_redaction_removes_body_header_and_metadata_text(self) -> None:
        import docx

        path = self.root / "secure.docx"
        document = docx.Document()
        paragraph = document.add_paragraph()
        paragraph.add_run("Top")
        paragraph.add_run("Secret")
        document.sections[0].header.paragraphs[0].text = "TopSecret"
        document.core_properties.author = "TopSecret"
        document.save(path)
        output = io.StringIO()

        with contextlib.redirect_stdout(output):
            self.assertEqual(0, edit_doc.cmd_secure_redact(argparse.Namespace(
                path=str(path),
                find="TopSecret",
                replace="[REDACTED]",
                expected_count=3,
                expected_sha256=hashlib.sha256(path.read_bytes()).hexdigest(),
                privacy_scrub=True,
            )))

        result = json.loads(output.getvalue())
        self.assertTrue(result["verification"]["originalTextAbsent"])
        self.assertEqual(3, result["redactedOccurrences"])
        with zipfile.ZipFile(path) as archive:
            for name in archive.namelist():
                self.assertNotIn(b"TopSecret", archive.read(name), name)
        reopened = docx.Document(path)
        self.assertEqual("[REDACTED]", reopened.paragraphs[0].text)
        self.assertEqual("[REDACTED]", reopened.sections[0].header.paragraphs[0].text)
        self.assertEqual("", reopened.core_properties.author or "")

    def test_secure_docx_redaction_fails_closed_on_uninspectable_media(self) -> None:
        import docx

        path = self.root / "media.docx"
        document = docx.Document()
        document.add_paragraph("TopSecret")
        document.save(path)
        with_media = self.root / "media-uninspectable.docx"
        _rewrite_zip(path, {}, {"word/media/uninspectable.bin": b"opaque"}, with_media)
        before = hashlib.sha256(with_media.read_bytes()).hexdigest()

        with self.assertRaises(SystemExit):
            with contextlib.redirect_stderr(io.StringIO()):
                edit_doc.cmd_secure_redact(argparse.Namespace(
                    path=str(with_media), find="TopSecret", replace="[REDACTED]",
                    expected_count=1, expected_sha256=None, privacy_scrub=False,
                ))
        self.assertEqual(before, hashlib.sha256(with_media.read_bytes()).hexdigest())

    def test_pptx_replace_matches_across_runs(self) -> None:
        from pptx import Presentation

        path = self.root / "deck.pptx"
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        box = slide.shapes.add_textbox(0, 0, 4_000_000, 1_000_000)
        paragraph = box.text_frame.paragraphs[0]
        paragraph.add_run().text = "Q"
        paragraph.add_run().text = "3"
        paragraph.add_run().text = " Results"
        presentation.save(path)

        with contextlib.redirect_stdout(io.StringIO()):
            self.assertEqual(0, edit_doc._replace_pptx(path, "Q3 Results", "Q4 Review", False))

        reopened = Presentation(path)
        self.assertEqual("Q4 Review", reopened.slides[0].shapes[-1].text)
        self.assertEqual("pass", validate_ooxml_package(path).status)

    def test_xlsx_precise_replace_preserves_sensitive_part_bytes(self) -> None:
        import openpyxl

        path = self.root / "model.xlsx"
        workbook = openpyxl.Workbook()
        workbook.active["A1"] = "Q3 Report"
        workbook.save(path)
        workbook.close()

        with zipfile.ZipFile(path) as archive:
            content_types = ET.fromstring(archive.read("[Content_Types].xml"))
        ET.SubElement(
            content_types,
            f"{{{CONTENT_TYPES_NS}}}Override",
            {
                "PartName": "/xl/pivotCache/pivotCacheDefinition1.xml",
                "ContentType": "application/vnd.openxmlformats-officedocument.spreadsheetml.pivotCacheDefinition+xml",
            },
        )
        sensitive = b"<?xml version='1.0' encoding='UTF-8'?><pivot-cache marker='preserve-me'/>"
        path = _rewrite_zip(
            path,
            {"[Content_Types].xml": ET.tostring(content_types, encoding="utf-8", xml_declaration=True)},
            {"xl/pivotCache/pivotCacheDefinition1.xml": sensitive},
            self.root / "model-sensitive.xlsx",
        )
        source_copy = self.root / "model-sensitive-source.xlsx"
        source_copy.write_bytes(path.read_bytes())
        with zipfile.ZipFile(path) as archive:
            untouched_sheet = archive.read("xl/worksheets/sheet1.xml")

        with contextlib.redirect_stdout(io.StringIO()):
            self.assertEqual(0, edit_doc._replace_xlsx(path, "Q3 Report", "Q4 Review", False))

        reopened = openpyxl.load_workbook(path, read_only=True)
        self.assertEqual("Q4 Review", reopened.active["A1"].value)
        reopened.close()
        with zipfile.ZipFile(path) as archive:
            self.assertEqual(sensitive, archive.read("xl/pivotCache/pivotCacheDefinition1.xml"))
            self.assertNotEqual(untouched_sheet, archive.read("xl/worksheets/sheet1.xml"))
        self.assertEqual("medium", scan_ooxml_risks(path)["riskLevel"])
        evidence = office_artifact_service._preservation_evidence(
            source_copy,
            path,
            scan_ooxml_risks(source_copy),
            {"xl/worksheets/sheet1.xml"},
        )
        self.assertTrue(evidence["verified"])
        self.assertGreater(evidence["sourceParts"], 0)
        self.assertEqual([], evidence["unauthorizedParts"])
        self.assertIn("pivotCaches", evidence["verifiedFeatures"])

    def test_xlsx_replace_does_not_reserialize_unmatched_editable_parts(self) -> None:
        import openpyxl

        path = self.root / "selective.xlsx"
        workbook = openpyxl.Workbook()
        workbook.active["A1"] = "Q3 Report"
        workbook.create_sheet("Untouched")["A1"] = "Stable content"
        workbook.save(path)
        workbook.close()
        with zipfile.ZipFile(path) as archive:
            untouched_sheet = archive.read("xl/worksheets/sheet2.xml")

        with contextlib.redirect_stdout(io.StringIO()):
            self.assertEqual(0, edit_doc._replace_xlsx(path, "Q3 Report", "Q4 Review", False))

        with zipfile.ZipFile(path) as archive:
            self.assertEqual(untouched_sheet, archive.read("xl/worksheets/sheet2.xml"))

    @unittest.skipUnless(
        edit_doc._find_soffice() and edit_doc._find_pdftoppm(),
        "LibreOffice and Poppler are required for native recalculation/render smoke",
    )
    def test_xlsx_recalculates_and_renders_with_native_qa_tools(self) -> None:
        import openpyxl

        path = self.root / "calculated.xlsx"
        workbook = openpyxl.Workbook()
        workbook.active["A1"] = 1
        workbook.active["A2"] = 2
        workbook.active["A3"] = "=SUM(A1:A2)"
        workbook.save(path)
        workbook.close()

        with contextlib.redirect_stdout(io.StringIO()):
            self.assertEqual(
                0,
                edit_doc.cmd_recalc_xlsx(argparse.Namespace(path=str(path), allow_risky=False)),
            )
        recalculated = openpyxl.load_workbook(path, data_only=True, read_only=True)
        self.assertEqual(3, recalculated.active["A3"].value)
        recalculated.close()

        render_dir = self.root / "rendered"
        with contextlib.redirect_stdout(io.StringIO()):
            self.assertEqual(
                0,
                edit_doc.cmd_render(
                    argparse.Namespace(path=str(path), outdir=str(render_dir), dpi=90, format="png")
                ),
            )
        self.assertTrue(list(render_dir.glob("*.png")))

    def test_xlsx_render_surface_plan_distinguishes_all_active_and_named_sheets(self) -> None:
        import openpyxl

        path = self.root / "surface-plan.xlsx"
        workbook = openpyxl.Workbook()
        workbook.active.title = "Summary"
        workbook.create_sheet("Detail")
        workbook.create_sheet("Hidden").sheet_state = "hidden"
        workbook.active = 1
        workbook.save(path)
        workbook.close()

        self.assertEqual(["Summary", "Detail"], edit_doc._xlsx_render_surfaces(path, "all"))
        self.assertEqual(["Detail"], edit_doc._xlsx_render_surfaces(path, "active"))
        self.assertEqual(["Summary"], edit_doc._xlsx_render_surfaces(path, "Summary"))

    def test_validator_rejects_missing_relationship_target(self) -> None:
        import docx

        path = self.root / "broken.docx"
        document = docx.Document()
        document.add_paragraph("Valid before corruption")
        document.save(path)
        with zipfile.ZipFile(path) as archive:
            rels = ET.fromstring(archive.read("word/_rels/document.xml.rels"))
        ET.SubElement(
            rels,
            f"{{{RELATIONSHIPS_NS}}}Relationship",
            {
                "Id": "rIdMissing",
                "Type": "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image",
                "Target": "media/missing.png",
            },
        )
        broken = self.root / "missing-target.docx"
        replacement = ET.tostring(rels, encoding="utf-8", xml_declaration=True)
        with zipfile.ZipFile(path) as source, zipfile.ZipFile(broken, "w") as destination:
            for info in source.infolist():
                data = replacement if info.filename == "word/_rels/document.xml.rels" else source.read(info.filename)
                destination.writestr(info, data)

        report = validate_ooxml_package(broken)
        self.assertEqual("fail", report.status)
        self.assertTrue(any(issue.code == "relationship.missing_target" for issue in report.errors))

    def test_failed_staged_validation_leaves_original_unchanged(self) -> None:
        import docx

        path = self.root / "safe.docx"
        document = docx.Document()
        document.add_paragraph("Original")
        document.save(path)
        before = hashlib.sha256(path.read_bytes()).hexdigest()
        staged = staging_path(path)
        staged.write_bytes(b"not-an-office-package")

        with self.assertRaises(ValueError):
            publish_staged_artifact(staged, path, self.root, validate=True)
        after = hashlib.sha256(path.read_bytes()).hexdigest()
        self.assertEqual(before, after)

    def test_transactional_job_writes_result_manifest(self) -> None:
        output = self.root / "job.docx"
        payload = {
            "jobVersion": 1,
            "format": "docx",
            "intent": "create_new",
            "output": str(output),
            "operations": [{
                "op": "create",
                "title": "Runtime job",
                "body": "A [linked result](https://example.com/result).",
            }],
            "preservationPolicy": "strict",
            "renderPolicy": "none",
        }
        job = OfficeArtifactJob.from_dict(payload, self.root)
        result, exit_code = execute_job(job, self.root)

        self.assertEqual(0, exit_code, result)
        self.assertTrue(result["ok"])
        self.assertTrue(output.exists())
        manifest = json.loads((self.root / "artifact-manifest.json").read_text(encoding="utf-8"))
        self.assertTrue(manifest["ok"])
        self.assertEqual("nexa-openxml", manifest["backend"])
        self.assertEqual("pass", validate_ooxml_package(output).status)

    def test_docx_validation_contract_is_enforced_before_publish(self) -> None:
        output = self.root / "contract-failed.docx"
        payload = {
            "jobVersion": 1,
            "format": "docx",
            "intent": "create_new",
            "output": str(output),
            "operations": [{"op": "create", "title": "Report", "body": "Draft body"}],
            "validationContract": {
                "required_text": ["Approved by Finance"],
                "min_paragraphs": 2,
            },
        }
        result, exit_code = execute_job(OfficeArtifactJob.from_dict(payload, self.root), self.root)

        self.assertEqual(1, exit_code)
        self.assertFalse(result["ok"])
        self.assertFalse(output.exists())
        self.assertIn("required_text.missing", result["error"])

    def test_pptx_validation_contract_passes_for_required_content(self) -> None:
        from pptx import Presentation

        source = self.root / "contract-source.pptx"
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = "Executive Summary"
        presentation.save(source)
        output = self.root / "contract-result.pptx"
        payload = {
            "jobVersion": 1,
            "format": "pptx",
            "intent": "edit_existing",
            "input": str(source),
            "output": str(output),
            "operations": [{"op": "replace", "find": "Executive", "replace": "Board"}],
            "validationContract": {
                "required_text": ["Board Summary"],
                "min_slides": 1,
                "max_slides": 1,
            },
        }
        result, exit_code = execute_job(OfficeArtifactJob.from_dict(payload, self.root), self.root)

        self.assertEqual(0, exit_code, result)
        contract = result["validation"]["backend"]["contract"]
        self.assertEqual("pass", contract["status"])
        self.assertTrue(output.exists())

    def test_job_rejects_manifest_path_equal_to_artifact_path(self) -> None:
        output = self.root / "conflict.docx"
        with self.assertRaisesRegex(ValueError, "manifest path must be distinct"):
            OfficeArtifactJob.from_dict(
                {
                    "jobVersion": 1,
                    "format": "docx",
                    "intent": "create_new",
                    "output": str(output),
                    "manifest": str(output),
                    "operations": [{"op": "create", "body": "Never published"}],
                },
                self.root,
            )

    def test_job_rejects_manifest_path_equal_to_input_path(self) -> None:
        source = self.root / "source.docx"
        source.write_bytes(b"placeholder")
        with self.assertRaisesRegex(ValueError, "manifest path must be distinct"):
            OfficeArtifactJob.from_dict(
                {
                    "jobVersion": 1,
                    "format": "docx",
                    "intent": "edit_existing",
                    "input": str(source),
                    "output": str(self.root / "result.docx"),
                    "manifest": str(source),
                },
                self.root,
            )

    def test_unpack_is_lossless_and_overwrite_requires_managed_directory(self) -> None:
        import docx

        source = self.root / "lossless.docx"
        document = docx.Document()
        document.add_paragraph("Preserve exact XML bytes")
        document.save(source)
        with zipfile.ZipFile(source) as archive:
            document_xml = archive.read("word/document.xml")

        outdir = self.root / "unpacked"
        with contextlib.redirect_stdout(io.StringIO()):
            self.assertEqual(
                0,
                edit_doc.cmd_unpack(
                    argparse.Namespace(path=str(source), outdir=str(outdir), overwrite=False)
                ),
            )
        self.assertEqual(document_xml, (outdir / "word" / "document.xml").read_bytes())
        self.assertTrue((outdir / edit_doc.UNPACK_MARKER).exists())

        unmanaged = self.root / "unmanaged"
        unmanaged.mkdir()
        (unmanaged / "keep.txt").write_text("user data", encoding="utf-8")
        with self.assertRaises(SystemExit):
            with contextlib.redirect_stderr(io.StringIO()):
                edit_doc.cmd_unpack(
                    argparse.Namespace(path=str(source), outdir=str(unmanaged), overwrite=True)
                )
        self.assertEqual("user data", (unmanaged / "keep.txt").read_text(encoding="utf-8"))

        # A copied/forged marker must never turn an arbitrary directory into a
        # recursive-delete target.
        (unmanaged / edit_doc.UNPACK_MARKER).write_text(
            json.dumps({"kind": "nexa-ooxml-unpack", "version": 1}),
            encoding="utf-8",
        )
        with self.assertRaises(SystemExit):
            with contextlib.redirect_stderr(io.StringIO()):
                edit_doc.cmd_unpack(
                    argparse.Namespace(path=str(source), outdir=str(unmanaged), overwrite=True)
                )
        self.assertEqual("user data", (unmanaged / "keep.txt").read_text(encoding="utf-8"))

    def test_unpack_refuses_workspace_root_even_with_overwrite(self) -> None:
        source = self.root / "minimal.docx"
        with zipfile.ZipFile(source, "w") as archive:
            archive.writestr("[Content_Types].xml", "<Types/>")
        with self.assertRaises(SystemExit):
            with contextlib.redirect_stderr(io.StringIO()):
                edit_doc.cmd_unpack(
                    argparse.Namespace(path=str(source), outdir=str(self.root), overwrite=True)
                )
        self.assertTrue(source.exists())

    def test_validator_and_unpack_reject_high_ratio_zip_before_extraction(self) -> None:
        path = self.root / "bomb.docx"
        with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as archive:
            archive.writestr("[Content_Types].xml", b"0" * (2 * 1024 * 1024))
            archive.writestr("_rels/.rels", "<Relationships/>")
            archive.writestr("word/document.xml", "<document/>")
        report = validate_ooxml_package(path)
        self.assertEqual("fail", report.status)
        self.assertTrue(any(issue.code == "zip.compression_ratio" for issue in report.errors))

        outdir = self.root / "bomb-unpacked"
        with self.assertRaises(SystemExit):
            with contextlib.redirect_stderr(io.StringIO()):
                edit_doc.cmd_unpack(
                    argparse.Namespace(path=str(path), outdir=str(outdir), overwrite=False)
                )
        self.assertFalse(outdir.exists())

    def test_validator_rejects_symlink_and_dtd_before_xml_parsing(self) -> None:
        path = self.root / "symlink.docx"
        symlink = zipfile.ZipInfo("word/media/link.png")
        symlink.create_system = 3
        symlink.external_attr = (0o120777 << 16)
        with zipfile.ZipFile(path, "w") as archive:
            archive.writestr(
                "[Content_Types].xml",
                '<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types"/>',
            )
            archive.writestr("_rels/.rels", "<Relationships/>")
            archive.writestr("word/document.xml", "<document/>")
            archive.writestr(symlink, "target.png")
        report = validate_ooxml_package(path)
        self.assertIn("zip.symlink", {issue.code for issue in report.errors})

        dtd_path = self.root / "dtd.docx"
        with zipfile.ZipFile(dtd_path, "w") as archive:
            archive.writestr(
                "[Content_Types].xml",
                '<!DOCTYPE Types [<!ENTITY x "boom">]><Types>&x;</Types>',
            )
            archive.writestr("_rels/.rels", "<Relationships/>")
            archive.writestr("word/document.xml", "<document/>")
        dtd_report = validate_ooxml_package(dtd_path)
        self.assertIn("xml.dtd_forbidden", {issue.code for issue in dtd_report.errors})

        late_dtd = self.root / "late-dtd.docx"
        with zipfile.ZipFile(late_dtd, "w") as archive:
            archive.writestr(
                "[Content_Types].xml",
                b'<?xml version="1.0"?>' + b" " * 5000
                + b'<!DOCTYPE Types [<!ENTITY x "xml">]><Types>&x;</Types>',
            )
            archive.writestr("_rels/.rels", "<Relationships/>")
            archive.writestr("word/document.xml", "<document/>")
        late_report = validate_ooxml_package(late_dtd)
        self.assertIn("xml.dtd_forbidden", {issue.code for issue in late_report.errors})

    def test_validator_rejects_windows_path_aliases_and_ads(self) -> None:
        path = self.root / "windows-collision.docx"
        with zipfile.ZipFile(path, "w") as archive:
            archive.writestr("[Content_Types].xml", "<Types/>")
            archive.writestr("_rels/.rels", "<Relationships/>")
            archive.writestr("word/document.xml", "<document/>")
            archive.writestr("word/foo.xml", "<a/>")
            archive.writestr("word/FOO.xml", "<b/>")
            archive.writestr("word/file:stream.xml", "<c/>")
        report = validate_ooxml_package(path)
        codes = {issue.code for issue in report.errors}
        self.assertIn("zip.windows_path_collision", codes)
        self.assertIn("zip.windows_path", codes)

    def test_risk_scan_decodes_formula_entities_and_blocks_external_relationships(self) -> None:
        import openpyxl

        workbook_path = self.root / "escaped-webservice.xlsx"
        workbook = openpyxl.Workbook()
        workbook.active["A1"] = '=WEBSERVICE("http://127.0.0.1:9/x")'
        workbook.save(workbook_path)
        workbook.close()
        with zipfile.ZipFile(workbook_path) as archive:
            sheet = archive.read("xl/worksheets/sheet1.xml")
        sheet = sheet.replace(
            b'WEBSERVICE("http://127.0.0.1:9/x")',
            b'WEBSERVICE&#40;&quot;http://127.0.0.1:9/x&quot;&#41;',
        )
        workbook_path = _rewrite_zip(
            workbook_path,
            {"xl/worksheets/sheet1.xml": sheet},
            output=self.root / "escaped-webservice-encoded.xlsx",
        )
        risk = scan_ooxml_risks(workbook_path)
        self.assertIn(
            "xl/worksheets/sheet1.xml",
            risk["features"]["externalFormulaFunctions"],
        )
        with self.assertRaisesRegex(RuntimeError, "externalFormulaFunctions"):
            office_artifact_service._assert_native_network_closed(workbook_path, "xlsx")

        try:
            from pptx import Presentation
        except ImportError:
            self.skipTest("python-pptx is not installed")
        deck = self.root / "external-image.pptx"
        presentation = Presentation()
        presentation.slides.add_slide(presentation.slide_layouts[6])
        presentation.save(deck)
        with zipfile.ZipFile(deck) as archive:
            rels_name = "ppt/slides/_rels/slide1.xml.rels"
            rels = ET.fromstring(archive.read(rels_name))
        ET.SubElement(rels, f"{{{RELATIONSHIPS_NS}}}Relationship", {
            "Id": "rIdExternal",
            "Type": "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image",
            "Target": "http://127.0.0.1:9/private.png",
            "TargetMode": "External",
        })
        deck = _rewrite_zip(
            deck,
            {rels_name: ET.tostring(rels, encoding="utf-8", xml_declaration=True)},
            output=self.root / "external-image-risk.pptx",
        )
        risk = scan_ooxml_risks(deck)
        self.assertIn(rels_name, risk["features"]["unsafeExternalRelationships"])
        with self.assertRaisesRegex(RuntimeError, "unsafeExternalRelationships"):
            office_artifact_service._assert_native_network_closed(deck, "pptx")

    def test_windows_automation_security_is_forced_disabled_and_restored(self) -> None:
        class FakeOfficeApplication:
            AutomationSecurity = 1

        app = FakeOfficeApplication()
        previous = office_artifact_service._force_disable_macros(app)
        self.assertEqual(1, previous)
        self.assertEqual(3, app.AutomationSecurity)
        office_artifact_service._restore_automation_security(app, previous)
        self.assertEqual(1, app.AutomationSecurity)

    def test_excel_calculation_wait_completes_and_times_out(self) -> None:
        class FakeExcel:
            def __init__(self, states):
                self.states = iter(states)
                self.last = 0

            @property
            def CalculationState(self):
                self.last = next(self.states, self.last)
                return self.last

        self.assertEqual(
            "done",
            office_artifact_service._wait_for_excel_calculation(FakeExcel([1, 2, 0]), 1.0),
        )
        with self.assertRaises(TimeoutError):
            office_artifact_service._wait_for_excel_calculation(FakeExcel([2]), 0.0)
        with self.assertRaises(TimeoutError):
            office_artifact_service._wait_for_excel_calculation(FakeExcel([1]), 0.0)

    def test_powerpoint_native_export_images_are_counted_and_normalized(self) -> None:
        raw = self.root / "powerpoint-raw"
        out = self.root / "powerpoint-normalized"
        raw.mkdir()
        for name, payload in (("Slide3.PNG", b"3"), ("Slide1.PNG", b"1"), ("Slide2.PNG", b"2")):
            (raw / name).write_bytes(payload)
        outputs = office_artifact_service._collect_powerpoint_export_images(raw, out, 3)
        self.assertEqual(
            ["slide-001.png", "slide-002.png", "slide-003.png"],
            [path.name for path in outputs],
        )
        self.assertEqual([b"1", b"2", b"3"], [path.read_bytes() for path in outputs])
        with self.assertRaisesRegex(RuntimeError, "expected 4"):
            office_artifact_service._collect_powerpoint_export_images(raw, out, 4)

    def test_powerpoint_native_render_force_disables_macros_before_open(self) -> None:
        class FakeSlides:
            Count = 2

        class FakePresentation:
            Slides = FakeSlides()

            def Export(self, output, image_format, width, height):
                self.export = (image_format, width, height)
                Path(output).mkdir(parents=True)
                Path(output, "Slide1.PNG").write_bytes(b"one")
                Path(output, "Slide2.PNG").write_bytes(b"two")

            def Close(self):
                self.closed = True

        class FakePresentations:
            def __init__(self, app):
                self.app = app

            def Open(self, path, WithWindow=False):
                self.app.security_when_opened = self.app.AutomationSecurity
                return FakePresentation()

        class FakeApplication:
            def __init__(self):
                self.AutomationSecurity = 1
                self.Version = "99.0"
                self.Presentations = FakePresentations(self)
                self.security_when_opened = None

            def Quit(self):
                self.quit = True

        app = FakeApplication()
        client = types.ModuleType("win32com.client")
        client.DispatchEx = lambda name: app
        package = types.ModuleType("win32com")
        package.client = client
        actions = []
        with mock.patch.dict(
            "sys.modules",
            {"win32com": package, "win32com.client": client},
        ), mock.patch.object(
            office_artifact_service,
            "_guard_office_process",
            return_value=None,
        ), mock.patch.object(
            office_artifact_service,
            "_windows_process_ids",
            return_value=set(),
        ), mock.patch.object(
            office_artifact_service,
            "_assert_native_network_closed",
            return_value={},
        ):
            outputs = office_artifact_service._windows_com_render_pptx(
                self.root / "input.pptx",
                self.root / "native-render",
                actions,
            )

        self.assertEqual(3, app.security_when_opened)
        self.assertEqual(1, app.AutomationSecurity)
        self.assertEqual(["slide-001.png", "slide-002.png"], [path.name for path in outputs])
        self.assertEqual("powerpoint-native", actions[0]["renderProfile"])

    def test_validate_reports_missing_xlsx_formula_cache_without_claiming_calculation(self) -> None:
        import openpyxl

        path = self.root / "uncalculated.xlsx"
        workbook = openpyxl.Workbook()
        workbook.active["A1"] = 1
        workbook.active["A2"] = "=A1+1"
        workbook.save(path)
        workbook.close()

        output = io.StringIO()
        with contextlib.redirect_stdout(output):
            self.assertEqual(0, edit_doc.cmd_validate(argparse.Namespace(path=str(path), json=True)))
        result = json.loads(output.getvalue())
        self.assertEqual("not_calculated", result["calculation"]["level"])
        self.assertEqual(1, result["calculation"]["formulaCells"])
        self.assertEqual(0, result["calculation"]["cachedFormulaCells"])

    def test_native_excel_proof_invalidates_formula_caches_before_com(self) -> None:
        import openpyxl

        path = self.root / "stale-cache.xlsx"
        workbook = openpyxl.Workbook()
        workbook.active["A1"] = 1
        workbook.active["A2"] = "=A1+1"
        workbook.save(path)
        workbook.close()
        with zipfile.ZipFile(path) as source:
            replacements = {}
            for info in source.infolist():
                data = source.read(info.filename)
                if info.filename == "xl/worksheets/sheet1.xml":
                    root = ET.fromstring(data)
                    for cell in root.iter():
                        if cell.tag.rsplit("}", 1)[-1] == "c" and cell.attrib.get("r") == "A2":
                            value = next(
                                child for child in list(cell) if child.tag.rsplit("}", 1)[-1] == "v"
                            )
                            value.text = "999"
                    replacements[info.filename] = ET.tostring(
                        root, encoding="utf-8", xml_declaration=True
                    )
        path = _rewrite_zip(path, replacements, {}, self.root / "stale-cache-filled.xlsx")
        before = office_artifact_service._all_part_hashes(path)
        self.assertEqual(1, office_artifact_service._clear_xlsx_formula_caches(path))
        after = office_artifact_service._all_part_hashes(path)
        self.assertNotEqual(before["xl/worksheets/sheet1.xml"], after["xl/worksheets/sheet1.xml"])
        renderer_dir = (
            Path(office_artifact_service.__file__).resolve().parents[2]
            / "xlsx-workbook-design"
            / "scripts"
        )
        if str(renderer_dir) not in office_artifact_service.sys.path:
            office_artifact_service.sys.path.insert(0, str(renderer_dir))
        from xlsx_model_renderer import inspect_formula_cache, inspect_formula_inventory

        cache = inspect_formula_cache(path)
        inventory = inspect_formula_inventory(path)
        self.assertEqual(0, cache["cachedFormulaCells"])
        self.assertEqual("A1+1", inventory["formulas"][0]["formula"])

    def test_xlsx_contract_v2_checks_tie_out_reconciliation_formula_pattern_and_sha_binding(self) -> None:
        import openpyxl

        path = self.root / "contract-v2.xlsx"
        workbook = openpyxl.Workbook()
        sheet = workbook.active
        sheet.title = "Model"
        sheet["A1"], sheet["A2"], sheet["A3"] = 1, 2, 3
        sheet["B1"], sheet["B2"] = 3, 6
        sheet["C1"], sheet["C2"] = "=A1*2", "=A2*2"
        workbook.save(path)
        workbook.close()
        contract_path = self.root / "contract-v2.json"
        contract_path.write_text(json.dumps({
            "contractVersion": 2,
            "required_sheets": ["Model"],
            "tie_outs": [{"left": "Model!A3", "right": "Model!B1", "tolerance": 0}],
            "reconciliations": [{"sumRange": "Model!A1:A2", "equals": "Model!B1", "tolerance": 0}],
            "formula_patterns": [{
                "sheet": "Model",
                "range": "C1:C2",
                "pattern": "^=A[12]\\*2$",
                "minMatches": 2,
                "requireConsistentRelativePattern": True,
            }],
        }), encoding="utf-8")

        result = edit_doc._validate_xlsx_contract(path, str(contract_path))

        self.assertEqual("pass", result["status"])
        self.assertEqual(64, len(result["evidence"]["artifactSha256"]))
        self.assertEqual(64, len(result["evidence"]["contractSha256"]))
        self.assertTrue(result["checks"]["tieOuts"][0]["matches"])
        self.assertTrue(result["checks"]["reconciliations"][0]["matches"])
        self.assertTrue(result["checks"]["formulaPatterns"][0]["consistentRelativePattern"])

        failing = json.loads(contract_path.read_text(encoding="utf-8"))
        failing["reconciliations"][0]["sumRange"] = "Model!A1:A1"
        contract_path.write_text(json.dumps(failing), encoding="utf-8")
        failed = edit_doc._validate_xlsx_contract(path, str(contract_path))
        self.assertEqual("fail", failed["status"])
        self.assertTrue(any(error["code"] == "reconciliation.mismatch" for error in failed["errors"]))

    def test_direct_validation_contract_rejects_unknown_fields(self) -> None:
        import docx

        path = self.root / "unknown-contract.docx"
        document = docx.Document()
        document.add_paragraph("text")
        document.save(path)
        contract = self.root / "unknown-contract.json"
        contract.write_text(json.dumps({"contractVersion": 2, "required_tex": ["text"]}), encoding="utf-8")
        with self.assertRaises(SystemExit):
            with contextlib.redirect_stderr(io.StringIO()):
                edit_doc._validate_docx_contract(path, str(contract))

    def test_job_rejects_nested_operation_paths_outside_workspace(self) -> None:
        outside_spec = self.root.parent / f"{self.root.name}-outside-spec.json"
        outside_spec.write_text("{}", encoding="utf-8")
        try:
            payload = {
                "jobVersion": 1,
                "format": "xlsx",
                "intent": "create_new",
                "output": "blocked.xlsx",
                "operations": [{"op": "create", "spec": f"../{outside_spec.name}"}],
            }
            job = OfficeArtifactJob.from_dict(payload, self.root)
            result, exit_code = execute_job(job, self.root)
        finally:
            outside_spec.unlink(missing_ok=True)

        self.assertEqual(1, exit_code)
        self.assertIn("path escapes workspace", result["error"])
        self.assertFalse((self.root / "blocked.xlsx").exists())

    def test_manifest_failure_rolls_back_newly_published_artifact(self) -> None:
        output = self.root / "rolled-back.docx"
        manifest_directory = self.root / "manifest-directory"
        manifest_directory.mkdir()
        job = OfficeArtifactJob.from_dict(
            {
                "jobVersion": 1,
                "format": "docx",
                "intent": "create_new",
                "output": str(output),
                "manifest": str(manifest_directory),
                "operations": [{"op": "create", "body": "Must roll back"}],
            },
            self.root,
        )

        result, exit_code = execute_job(job, self.root)

        self.assertEqual(1, exit_code)
        self.assertFalse(result["ok"])
        self.assertTrue(result["rollbackApplied"])
        self.assertFalse(output.exists())

    def test_auto_backend_routes_recalculation_and_finalization_by_capability(self) -> None:
        base = {
            "jobVersion": 1,
            "input": str(self.root / "source.xlsx"),
            "output": str(self.root / "result.xlsx"),
            "format": "xlsx",
            "operations": [],
        }
        source = Path(base["input"])
        source.write_bytes(b"placeholder")
        recalculate = OfficeArtifactJob.from_dict(
            {**base, "intent": "edit_existing", "operations": [{"op": "recalculate"}]},
            self.root,
        )
        self.assertEqual("libreoffice", office_artifact_service._select_backend(recalculate))

        docx = self.root / "source.docx"
        docx.write_bytes(b"placeholder")
        finalize = OfficeArtifactJob.from_dict(
            {
                "jobVersion": 1,
                "format": "docx",
                "intent": "finalize",
                "input": str(docx),
                "output": str(self.root / "result.docx"),
            },
            self.root,
        )
        self.assertEqual("windows-com", office_artifact_service._select_backend(finalize))

    def test_transactional_xlsx_job_publishes_final_qa_sidecar(self) -> None:
        output = self.root / "job-model.xlsx"
        spec = self.root / "model-spec.json"
        spec.write_text(
            json.dumps({
                "title": "Job model",
                "sheets": [{
                    "name": "Summary",
                    "start_cell": "A1",
                    "headers": ["Metric", "Value"],
                    "rows": [["Revenue", 42]],
                }],
            }),
            encoding="utf-8",
        )
        payload = {
            "jobVersion": 1,
            "format": "xlsx",
            "intent": "create_new",
            "output": str(output),
            "operations": [{"op": "create", "spec": str(spec)}],
            "renderPolicy": "none",
        }
        job = OfficeArtifactJob.from_dict(payload, self.root)
        result, exit_code = execute_job(job, self.root)

        self.assertEqual(0, exit_code, result)
        qa_path = output.with_suffix(".xlsx.qa.json")
        self.assertTrue(qa_path.exists())
        qa = json.loads(qa_path.read_text(encoding="utf-8"))
        self.assertEqual(str(output), qa["path"])
        self.assertNotIn("nexa-stage", json.dumps(result))


if __name__ == "__main__":
    unittest.main()
