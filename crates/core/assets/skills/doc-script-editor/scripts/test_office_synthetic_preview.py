from __future__ import annotations

import tempfile
import unittest
from pathlib import Path

from office_synthetic_preview import create_synthetic_preview


class OfficeSyntheticPreviewTests(unittest.TestCase):
    def test_three_formats_emit_sha_bound_non_final_structural_previews(self) -> None:
        import docx
        import openpyxl
        from pptx import Presentation

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            document_path = root / "sample.docx"
            document = docx.Document()
            document.add_heading("Decision", 1)
            document.add_paragraph("Evidence")
            document.save(document_path)
            workbook_path = root / "sample.xlsx"
            workbook = openpyxl.Workbook()
            workbook.active.append(["Metric", "Value"])
            workbook.active.append(["Revenue", 100])
            workbook.save(workbook_path)
            workbook.close()
            presentation_path = root / "sample.pptx"
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[5])
            slide.shapes.title.text = "Decision"
            presentation.save(presentation_path)

            for index, artifact in enumerate((document_path, workbook_path, presentation_path)):
                evidence = create_synthetic_preview(artifact, root / f"preview-{index}")
                self.assertFalse(evidence["isFinalRenderEvidence"])
                self.assertRegex(evidence["artifactSha256"], r"^[0-9a-f]{64}$")
                self.assertTrue(evidence["files"])
                for item in evidence["files"]:
                    self.assertTrue(Path(item["path"]).is_file())


if __name__ == "__main__":
    unittest.main()
