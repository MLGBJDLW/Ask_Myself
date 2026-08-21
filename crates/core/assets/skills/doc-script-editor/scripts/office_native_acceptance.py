#!/usr/bin/env python3
"""Run SHA-bound Microsoft Office COM acceptance for DOCX, XLSX, and PPTX."""

from __future__ import annotations

import argparse
import hashlib
import importlib.metadata
import json
import os
import platform
import re
import subprocess
import sys
import traceback
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Callable

import docx
import openpyxl
from openpyxl.chart import BarChart, Reference
from pptx import Presentation

import office_artifact_service
from office_artifact_runtime import WindowsComBackend, validate_ooxml_package


SHA_PATTERN = re.compile(r"^[0-9a-f]{40}$")


def utc_now() -> str:
    return datetime.now(timezone.utc).isoformat().replace("+00:00", "Z")


def sha256(path: Path) -> str:
    return hashlib.sha256(path.read_bytes()).hexdigest()


def git(repo_root: Path, *arguments: str) -> str:
    completed = subprocess.run(
        ["git", "-C", str(repo_root), *arguments],
        text=True,
        capture_output=True,
        check=False,
        timeout=30,
    )
    if completed.returncode != 0:
        raise RuntimeError(completed.stderr.strip() or "git command failed")
    return completed.stdout.strip()


def assert_repository_state(repo_root: Path, expected_sha: str) -> None:
    if not SHA_PATTERN.fullmatch(expected_sha):
        raise ValueError("expected SHA must be a lowercase 40-character Git object ID")
    actual_sha = git(repo_root, "rev-parse", "HEAD").lower()
    if actual_sha != expected_sha:
        raise RuntimeError(f"checked-out SHA {actual_sha} does not match {expected_sha}")
    dirty = git(repo_root, "status", "--porcelain=v1", "--untracked-files=all")
    if dirty:
        raise RuntimeError("native acceptance requires a clean checkout")


def create_docx(path: Path) -> None:
    document = docx.Document()
    document.add_heading("Nexa native Word acceptance", level=1)
    document.add_paragraph("Word must preserve this SHA-bound acceptance marker.")
    document.add_table(rows=2, cols=2).cell(0, 0).text = "Native"
    document.save(path)


def create_xlsx(path: Path) -> None:
    workbook = openpyxl.Workbook()
    worksheet = workbook.active
    worksheet.title = "Acceptance"
    worksheet.append(["Quarter", "Value"])
    worksheet.append(["Q1", 1])
    worksheet.append(["Q2", 2])
    worksheet["B4"] = "=SUM(B2:B3)"
    chart = BarChart()
    chart.title = "Native values"
    chart.add_data(Reference(worksheet, min_col=2, min_row=1, max_row=3), titles_from_data=True)
    chart.set_categories(Reference(worksheet, min_col=1, min_row=2, max_row=3))
    worksheet.add_chart(chart, "D2")
    workbook.save(path)
    workbook.close()


def create_pptx(path: Path) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Nexa native PowerPoint acceptance"
    slide.placeholders[1].text = "PowerPoint must preserve and render this slide."
    presentation.save(path)


def verify_docx(path: Path) -> None:
    document = docx.Document(path)
    text = "\n".join(paragraph.text for paragraph in document.paragraphs)
    if "Nexa native Word acceptance" not in text:
        raise RuntimeError("Word-native save did not preserve the acceptance marker")


def verify_xlsx(path: Path) -> None:
    workbook = openpyxl.load_workbook(path, data_only=True, read_only=True)
    try:
        if workbook["Acceptance"]["B4"].value != 3:
            raise RuntimeError("Excel-native calculation did not cache B4=3")
    finally:
        workbook.close()


def verify_pptx(path: Path) -> None:
    presentation = Presentation(path)
    text = "\n".join(
        shape.text
        for slide in presentation.slides
        for shape in slide.shapes
        if hasattr(shape, "text")
    )
    if "Nexa native PowerPoint acceptance" not in text:
        raise RuntimeError("PowerPoint-native save did not preserve the acceptance marker")


def file_evidence(paths: list[Path], root: Path) -> list[dict[str, Any]]:
    return [
        {
            "path": path.relative_to(root).as_posix(),
            "bytes": path.stat().st_size,
            "sha256": sha256(path),
        }
        for path in paths
    ]


def accept_artifact(
    artifact_format: str,
    path: Path,
    output_root: Path,
    verify: Callable[[Path], None],
) -> dict[str, Any]:
    before_sha = sha256(path)
    actions: list[dict[str, Any]] = []
    office_artifact_service._windows_com_finalize(path, artifact_format, actions)
    verify(path)
    validation = validate_ooxml_package(path)
    if validation.status == "fail":
        raise RuntimeError(f"{artifact_format} failed OOXML validation after native save")

    render_dir = output_root / "renders" / artifact_format
    render_dir.parent.mkdir(parents=True, exist_ok=True)
    renderer = {
        "docx": office_artifact_service._windows_com_render_docx,
        "xlsx": office_artifact_service._windows_com_render_xlsx,
        "pptx": office_artifact_service._windows_com_render_pptx,
    }[artifact_format]
    rendered = renderer(path, render_dir, actions)
    if not rendered or any(not item.is_file() or item.stat().st_size == 0 for item in rendered):
        raise RuntimeError(f"{artifact_format} native render produced incomplete evidence")

    expected_engines = {
        "docx": "microsoft-word-com",
        "xlsx": "microsoft-excel-com",
        "pptx": "microsoft-powerpoint-com",
    }
    native_actions = [item for item in actions if item.get("engine") == expected_engines[artifact_format]]
    if len(native_actions) < 2 or any(item.get("status") != "ok" for item in native_actions):
        raise RuntimeError(f"{artifact_format} native finalize/render evidence is incomplete")

    return {
        "format": artifact_format,
        "path": path.relative_to(output_root).as_posix(),
        "sourceSha256": before_sha,
        "artifactSha256": sha256(path),
        "validation": validation.to_dict(),
        "renderedSurfaces": file_evidence(rendered, output_root),
        "actions": actions,
    }


def run(repo_root: Path, output_root: Path, expected_sha: str) -> dict[str, Any]:
    if platform.system() != "Windows" or os.environ.get("NEXA_RUN_OFFICE_NATIVE_SMOKE") != "1":
        raise RuntimeError("native acceptance requires Windows and NEXA_RUN_OFFICE_NATIVE_SMOKE=1")
    repo_root = repo_root.expanduser().resolve(strict=True)
    output_root = output_root.expanduser().resolve()
    if output_root == repo_root or output_root.is_relative_to(repo_root):
        raise ValueError("native acceptance output must be outside the Git checkout")
    if output_root.exists() and any(output_root.iterdir()):
        raise FileExistsError(f"native acceptance output is not empty: {output_root}")
    output_root.mkdir(parents=True, exist_ok=True)
    assert_repository_state(repo_root, expected_sha)

    backend = WindowsComBackend().preflight()
    if backend.status != "ready":
        raise RuntimeError(f"Windows COM backend is not ready: {backend.detail}")
    artifacts_dir = output_root / "artifacts"
    artifacts_dir.mkdir()
    artifacts = {
        "docx": artifacts_dir / "native-acceptance.docx",
        "xlsx": artifacts_dir / "native-acceptance.xlsx",
        "pptx": artifacts_dir / "native-acceptance.pptx",
    }
    create_docx(artifacts["docx"])
    create_xlsx(artifacts["xlsx"])
    create_pptx(artifacts["pptx"])

    results = [
        accept_artifact("docx", artifacts["docx"], output_root, verify_docx),
        accept_artifact("xlsx", artifacts["xlsx"], output_root, verify_xlsx),
        accept_artifact("pptx", artifacts["pptx"], output_root, verify_pptx),
    ]
    assert_repository_state(repo_root, expected_sha)
    office_versions = sorted(
        {
            action["engine"]: action["engineVersion"]
            for result in results
            for action in result["actions"]
            if action.get("engine") and action.get("engineVersion")
        }.items()
    )
    return {
        "schemaVersion": 1,
        "kind": "office-native-acceptance",
        "status": "pass",
        "repositorySha": expected_sha,
        "worktreeClean": True,
        "completedAt": utc_now(),
        "runner": {
            "platform": platform.platform(),
            "machine": platform.machine(),
            "python": platform.python_version(),
            "pywin32": importlib.metadata.version("pywin32"),
        },
        "officeVersions": [
            {"engine": engine, "version": version} for engine, version in office_versions
        ],
        "artifacts": results,
    }


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--repo-root", required=True)
    parser.add_argument("--output-dir", required=True)
    parser.add_argument("--expected-sha", required=True)
    args = parser.parse_args()
    output_root = Path(args.output_dir).expanduser().resolve()
    started_at = utc_now()
    try:
        report = run(Path(args.repo_root), output_root, args.expected_sha.lower())
        report["startedAt"] = started_at
    except Exception as error:
        output_root.mkdir(parents=True, exist_ok=True)
        report = {
            "schemaVersion": 1,
            "kind": "office-native-acceptance",
            "status": "fail",
            "repositorySha": args.expected_sha.lower(),
            "startedAt": started_at,
            "completedAt": utc_now(),
            "error": {"type": type(error).__name__, "message": str(error)},
        }
        (output_root / "office-native-acceptance.json").write_text(
            json.dumps(report, ensure_ascii=False, indent=2, sort_keys=True) + "\n",
            encoding="utf-8",
        )
        traceback.print_exc()
        return 1
    (output_root / "office-native-acceptance.json").write_text(
        json.dumps(report, ensure_ascii=False, indent=2, sort_keys=True) + "\n",
        encoding="utf-8",
    )
    print(json.dumps(report, ensure_ascii=False, indent=2, sort_keys=True))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
