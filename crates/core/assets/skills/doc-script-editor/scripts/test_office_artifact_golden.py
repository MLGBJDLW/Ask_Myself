from __future__ import annotations

import json
import shutil
import tempfile
import unittest
from pathlib import Path

from office_artifact_engine import OfficeArtifactEngine


class OfficeArtifactGoldenTests(unittest.TestCase):
    def test_public_inspect_exposes_pptx_slide_and_shape_ids(self) -> None:
        request = {
            "requestVersion": 2,
            "format": "pptx",
            "intent": "create",
            "destination": str(self.root / "inspectable.pptx"),
            "operations": [{"op": "create", "spec": str(self._copy_spec("pptx-spec.json"))}],
            "guarantees": {"quality": "standard", "render": "none"},
            "validation": {"contractVersion": 2, "min_slides": 2},
        }
        candidate = self.engine.execute(request)
        inspected = self.engine.inspect(candidate["candidatePath"], "pptx")
        slides = inspected["profile"]["slide_details"]
        self.assertEqual(["256", "257"], [slide["slide_id"] for slide in slides])
        text_shapes = [
            shape
            for slide in slides
            for shape in slide["shape_details"]
            if shape["text"]
        ]
        self.assertTrue(all(shape["shapeId"] for shape in text_shapes))
        self.assertTrue(all(shape["shapeName"] for shape in text_shapes))

    def setUp(self) -> None:
        self.temp = tempfile.TemporaryDirectory()
        self.root = Path(self.temp.name).resolve()
        self.engine = OfficeArtifactEngine(self.root)
        self.golden = Path(__file__).resolve().parents[1] / "tests" / "golden"
        self.expectations = json.loads((self.golden / "expectations.json").read_text(encoding="utf-8"))

    def tearDown(self) -> None:
        self.temp.cleanup()

    def _copy_spec(self, name: str) -> Path:
        destination = self.root / name
        shutil.copy2(self.golden / name, destination)
        return destination

    def test_cross_format_golden_candidates_match_semantic_oracles_and_sha_evidence(self) -> None:
        import docx
        import openpyxl
        from pptx import Presentation

        requests = [
            {
                "requestVersion": 2,
                "format": "docx",
                "intent": "create",
                "destination": str(self.root / "golden.docx"),
                "operations": [{"op": "create", "spec": str(self._copy_spec("docx-spec.json"))}],
                "guarantees": {"quality": "standard", "render": "none"},
                "validation": {
                    "contractVersion": 2,
                    "required_text": [self.expectations["docx"]["requiredText"]],
                    "min_tables": 1,
                    "require_table_header_rows": True,
                },
            },
            {
                "requestVersion": 2,
                "format": "xlsx",
                "intent": "create",
                "destination": str(self.root / "golden.xlsx"),
                "operations": [{"op": "create", "spec": str(self._copy_spec("xlsx-spec.json"))}],
                "guarantees": {"quality": "standard", "calculation": "static", "render": "none"},
                "validation": {"contractVersion": 2, "required_sheets": ["Summary"]},
            },
            {
                "requestVersion": 2,
                "format": "pptx",
                "intent": "create",
                "destination": str(self.root / "golden.pptx"),
                "operations": [{"op": "create", "spec": str(self._copy_spec("pptx-spec.json"))}],
                "guarantees": {"quality": "standard", "render": "none"},
                "validation": {
                    "contractVersion": 2,
                    "min_slides": 2,
                    "max_slides": 2,
                    "required_text": [self.expectations["pptx"]["requiredText"]],
                },
            },
        ]
        outcomes = {request["format"]: self.engine.execute(request) for request in requests}
        for outcome in outcomes.values():
            self.assertEqual("candidate", outcome["status"])
            self.assertEqual(outcome["sha256"], outcome["renderEvidence"]["artifactSha256"])
            self.assertTrue(outcome["renderEvidence"]["complete"])
            self.assertEqual(
                outcome["sha256"],
                outcome["validation"]["backend"]["contract"]["evidence"]["artifactSha256"],
            )
            self.assertFalse(Path(outcome["destination"]).exists())

        document = docx.Document(outcomes["docx"]["candidatePath"])
        self.assertGreaterEqual(len(document.paragraphs), self.expectations["docx"]["paragraphsAtLeast"])
        self.assertEqual(self.expectations["docx"]["tables"], len(document.tables))

        workbook = openpyxl.load_workbook(outcomes["xlsx"]["candidatePath"], data_only=False)
        try:
            self.assertEqual(self.expectations["xlsx"]["sheets"], workbook.sheetnames)
            self.assertEqual("=B2-B3", workbook["Summary"]["B4"].value)
        finally:
            workbook.close()

        presentation = Presentation(outcomes["pptx"]["candidatePath"])
        self.assertEqual(self.expectations["pptx"]["slides"], len(presentation.slides))
        text = " ".join(shape.text for slide in presentation.slides for shape in slide.shapes if hasattr(shape, "text"))
        self.assertIn(self.expectations["pptx"]["requiredText"], text)

    def test_contract_fault_injection_never_publishes_destination(self) -> None:
        spec = self._copy_spec("docx-spec.json")
        destination = self.root / "must-not-publish.docx"
        request = {
            "requestVersion": 2,
            "format": "docx",
            "intent": "create",
            "destination": str(destination),
            "operations": [{"op": "create", "spec": str(spec)}],
            "guarantees": {"quality": "standard", "render": "none"},
            "validation": {"contractVersion": 2, "required_text": ["INJECTED MISSING TEXT"]},
            "delivery": {"mode": "publish"},
        }
        with self.assertRaises(Exception):
            self.engine.execute(request)
        self.assertFalse(destination.exists())

    def test_every_format_publishes_with_receipt_and_restores_without_residue(self) -> None:
        requests = [
            {
                "requestVersion": 2,
                "format": "docx",
                "intent": "create",
                "destination": str(self.root / "published.docx"),
                "operations": [{"op": "create", "spec": str(self._copy_spec("docx-spec.json"))}],
                "guarantees": {"quality": "standard", "render": "none"},
                "validation": {"contractVersion": 2, "required_text": ["Approve the verified plan."]},
            },
            {
                "requestVersion": 2,
                "format": "xlsx",
                "intent": "create",
                "destination": str(self.root / "published.xlsx"),
                "operations": [{"op": "create", "spec": str(self._copy_spec("xlsx-spec.json"))}],
                "guarantees": {"quality": "standard", "calculation": "static", "render": "none"},
                "validation": {"contractVersion": 2, "required_sheets": ["Summary"]},
            },
            {
                "requestVersion": 2,
                "format": "pptx",
                "intent": "create",
                "destination": str(self.root / "published.pptx"),
                "operations": [{"op": "create", "spec": str(self._copy_spec("pptx-spec.json"))}],
                "guarantees": {"quality": "standard", "render": "none"},
                "validation": {
                    "contractVersion": 2,
                    "min_slides": 2,
                    "max_slides": 2,
                    "required_slide_titles": ["Golden Decision", "Decision options"],
                },
            },
        ]
        for request in requests:
            candidate = self.engine.execute(request)
            published = self.engine.decide(candidate["candidateId"], "publish")
            destination = Path(published["path"])
            manifest = destination.with_suffix(f"{destination.suffix}.manifest.json")
            self.assertTrue(destination.exists(), request["format"])
            self.assertTrue(manifest.exists(), request["format"])
            self.assertEqual(64, len(published["sha256"]))
            restored = self.engine.restore(published["receiptId"])
            self.assertEqual("restored", restored["status"])
            self.assertFalse(destination.exists(), request["format"])
            self.assertFalse(manifest.exists(), request["format"])
            if request["format"] == "xlsx":
                self.assertFalse(destination.with_suffix(".xlsx.qa.json").exists())


if __name__ == "__main__":
    unittest.main()
