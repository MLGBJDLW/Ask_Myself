#!/usr/bin/env python3
"""edit_doc.py — Skill-bundled document editor for DOCX/PPTX/PDF/XLSX.

Invoked via the app's `run_shell` tool. Reads/writes files on disk;
never accepts document content over argv. Lazy-imports backend libs
so `check` works with nothing installed.

Exit codes:
  0 success
  1 generic error
  2 missing dependency
  3 bad input / path validation failed
"""
from __future__ import annotations

import argparse
import difflib
import hashlib
import json
import os
import re
import shutil
import subprocess
import sys
import tempfile
import zipfile
from bisect import bisect_right
from pathlib import Path
from typing import Any
from xml.etree import ElementTree as ET

from office_artifact_runtime import (
    office_backend_statuses,
    publish_staged_artifact,
    scan_ooxml_risks,
    staging_path,
    validate_ooxml_package,
    write_artifact_manifest,
)

MAX_EXTRACT_BYTES = 50 * 1024
HISTORY_DIR = ".nexa/doc-history"
EXCEL_ERRORS = ("#VALUE!", "#DIV/0!", "#REF!", "#NAME?", "#NULL!", "#NUM!", "#N/A")
UNPACK_MARKER = ".nexa-ooxml-unpack.json"
DOCX_EXTENSIONS = {"docx", "docm", "dotx", "dotm"}
XLSX_EXTENSIONS = {"xlsx", "xlsm", "xltx", "xltm"}
PPTX_EXTENSIONS = {"pptx", "pptm", "potx", "potm"}
OOXML_EXTENSIONS = DOCX_EXTENSIONS | XLSX_EXTENSIONS | PPTX_EXTENSIONS


# ---------------------------------------------------------------------------
# Utilities
# ---------------------------------------------------------------------------

def _suppress_windows_error_dialogs() -> None:
    if os.name != "nt":
        return
    try:
        import ctypes
        kernel32 = ctypes.windll.kernel32
        sem_failcriticalerrors = 0x0001
        sem_nogpfaulterrorbox = 0x0002
        sem_noopenfileerrorbox = 0x8000
        kernel32.SetErrorMode(
            sem_failcriticalerrors | sem_nogpfaulterrorbox | sem_noopenfileerrorbox
        )
    except Exception:
        pass


_suppress_windows_error_dialogs()


def _run_subprocess(cmd: list[str], **kwargs: Any) -> subprocess.CompletedProcess[str]:
    if os.name == "nt":
        kwargs["creationflags"] = kwargs.get("creationflags", 0) | getattr(
            subprocess,
            "CREATE_NO_WINDOW",
            0,
        )
    return subprocess.run(cmd, check=kwargs.pop("check", False), **kwargs)


def _die(msg: str, code: int = 1) -> None:
    print(msg, file=sys.stderr)
    sys.exit(code)


def _missing(pkg: str) -> None:
    print(f"MISSING_DEP: {pkg}", file=sys.stderr)
    print(f"Install with: python -m pip install {pkg}", file=sys.stderr)
    sys.exit(2)


def _validate_path(raw: str, must_exist: bool = True) -> Path:
    if not raw:
        _die("ERROR: --path is required", 3)
    p = Path(raw)
    if not p.is_absolute():
        _die(f"ERROR: --path must be absolute: {raw}", 3)
    try:
        resolved = p.resolve()
        cwd = Path.cwd().resolve()
        # Basic traversal guard: resolved path must live under cwd.
        resolved.relative_to(cwd)
    except ValueError:
        _die(f"ERROR: path escapes workspace: {raw}", 3)
    except OSError as e:
        _die(f"ERROR: cannot resolve path: {e}", 3)
    if must_exist and not resolved.exists():
        _die(f"ERROR: file not found: {resolved}", 3)
    return resolved


def _validate_output_path(raw: str, suffixes: set[str]) -> Path:
    p = _validate_path(raw, must_exist=False)
    if _ext(p) not in suffixes:
        _die(f"ERROR: output path must end with one of: {', '.join(sorted(suffixes))}", 3)
    p.parent.mkdir(parents=True, exist_ok=True)
    return p


def _validate_output_dir(raw: str, *, allow_existing: bool = True) -> Path:
    if not raw:
        _die("ERROR: output directory is required", 3)
    p = Path(raw)
    if not p.is_absolute():
        _die(f"ERROR: output directory must be absolute: {raw}", 3)
    try:
        resolved = p.resolve()
        workspace_root = Path.cwd().resolve()
        resolved.relative_to(workspace_root)
    except ValueError:
        _die(f"ERROR: output directory escapes workspace: {raw}", 3)
    except OSError as e:
        _die(f"ERROR: cannot resolve output directory: {e}", 3)
    if resolved == workspace_root:
        _die("ERROR: output directory cannot be the workspace root", 3)
    if resolved.exists() and not allow_existing:
        _die(f"ERROR: output directory already exists: {resolved}", 3)
    resolved.mkdir(parents=True, exist_ok=True)
    return resolved


def _ext(p: Path) -> str:
    return p.suffix.lower().lstrip(".")


def _staged_copy(path: Path) -> Path:
    staged = staging_path(path)
    shutil.copy2(path, staged)
    return staged


def _publish_edit(staged: Path, path: Path) -> tuple[Path | None, dict[str, Any] | None]:
    try:
        snapshot, report = publish_staged_artifact(staged, path, Path.cwd(), validate=True)
    except ValueError as error:
        staged.unlink(missing_ok=True)
        _die(f"VALIDATION_FAILED: staged artifact was not published\n{error}", 1)
    return snapshot, report.to_dict() if report is not None else None


def _parse_pages(spec: str | None, total: int) -> list[int]:
    if not spec:
        return list(range(total))
    out: list[int] = []
    for part in spec.split(","):
        part = part.strip()
        if not part:
            continue
        if "-" in part:
            a, b = part.split("-", 1)
            out.extend(range(int(a) - 1, int(b)))
        else:
            out.append(int(part) - 1)
    return [i for i in out if 0 <= i < total]


def _truncate(text: str) -> str:
    raw = text.encode("utf-8")
    if len(raw) <= MAX_EXTRACT_BYTES:
        return text
    cut = raw[:MAX_EXTRACT_BYTES].decode("utf-8", errors="ignore")
    return cut + f"\n\n[TRUNCATED: output exceeded {MAX_EXTRACT_BYTES} bytes]"


def _read_json(path: str) -> dict[str, Any]:
    spec_path = _validate_path(path)
    with spec_path.open("r", encoding="utf-8") as f:
        data = json.load(f)
    if not isinstance(data, dict):
        _die("ERROR: JSON spec root must be an object", 3)
    return data


def _read_text(path: str) -> str:
    text_path = _validate_path(path)
    return text_path.read_text(encoding="utf-8")


def _find_soffice() -> str | None:
    for name in ("soffice", "soffice.com", "libreoffice"):
        found = shutil.which(name)
        if found:
            return found
    for candidate in (
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files\LibreOffice\program\soffice.com",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
    ):
        if Path(candidate).exists():
            return candidate
    return None


def _find_pdftoppm() -> str | None:
    for name in ("pdftoppm", "pdftoppm.exe"):
        found = shutil.which(name)
        if found:
            return found
    return None


def _soffice_env() -> dict[str, str]:
    env = os.environ.copy()
    env.setdefault("SAL_USE_VCLPLUGIN", "svp")
    return env


def _soffice_base_cmd(soffice: str, profile_dir: Path) -> list[str]:
    return [
        soffice,
        f"-env:UserInstallation={profile_dir.resolve().as_uri()}",
        "--headless",
        "--invisible",
        "--norestore",
        "--nolockcheck",
        "--nodefault",
    ]


def _run_soffice_convert(path: Path, to: str, outdir: Path) -> subprocess.CompletedProcess[str]:
    soffice = _find_soffice()
    if not soffice:
        _die("MISSING_DEP: LibreOffice/soffice\nInstall LibreOffice and ensure soffice is on PATH.", 2)
    with tempfile.TemporaryDirectory(prefix="nexa-lo-profile-") as profile:
        cmd = [
            *_soffice_base_cmd(soffice, Path(profile)),
            "--convert-to",
            to,
            "--outdir",
            str(outdir),
            str(path),
        ]
        return _run_subprocess(
            cmd,
            text=True,
            capture_output=True,
            check=False,
            env=_soffice_env(),
        )


def _expected_converted_path(input_path: Path, outdir: Path, to: str) -> Path:
    ext = to.split(":", 1)[0].split()[0].lstrip(".")
    return outdir / f"{input_path.stem}.{ext}"


# ---------------------------------------------------------------------------
# check
# ---------------------------------------------------------------------------

def cmd_check(args: argparse.Namespace) -> int:
    backends = [
        ("python-docx", "docx", "docx"),
        ("python-pptx", "pptx", "pptx"),
        ("pypdf", "pypdf", "pdf"),
        ("openpyxl", "openpyxl", "xlsx"),
    ]
    missing_core = []
    missing_optional = []
    results: list[dict[str, Any]] = []
    if not args.json:
        print(f"python: {sys.version.split()[0]}")
    for display, mod, artifact_format in backends:
        required = args.format == "all" or args.format == artifact_format
        try:
            imported = __import__(mod)
            ver = getattr(imported, "__version__", "unknown")
            results.append({
                "id": display,
                "module": mod,
                "status": "ok",
                "version": str(ver),
                "required": required,
                "formats": [artifact_format],
            })
            if not args.json:
                print(f"  {display:<14} OK      ({ver})")
        except ImportError:
            results.append({
                "id": display,
                "module": mod,
                "status": "missing",
                "required": required,
                "formats": [artifact_format],
            })
            if not args.json:
                print(f"  {display:<14} MISSING")
            (missing_core if required else missing_optional).append(display)
        except Exception as e:  # noqa: BLE001
            # Backend present but broken (e.g. numpy ABI mismatch). Treat as missing.
            results.append({
                "id": display,
                "module": mod,
                "status": "broken",
                "required": required,
                "formats": [artifact_format],
                "detail": f"{type(e).__name__}: {e}",
            })
            if not args.json:
                print(f"  {display:<14} BROKEN  ({type(e).__name__}: {e})")
            (missing_core if required else missing_optional).append(display)
    soffice = _find_soffice()
    if soffice:
        results.append({
            "id": "LibreOffice",
            "status": "ok",
            "path": soffice,
            "required": False,
        })
        if not args.json:
            print(f"  LibreOffice    OK      ({soffice})")
    else:
        results.append({
            "id": "LibreOffice",
            "status": "missing",
            "required": False,
            "detail": "needed for convert/render QA",
        })
        if not args.json:
            print("  LibreOffice    MISSING (needed for convert/render QA)")
    pdftoppm = _find_pdftoppm()
    if pdftoppm:
        results.append({
            "id": "Poppler",
            "status": "ok",
            "path": pdftoppm,
            "required": False,
        })
        if not args.json:
            print(f"  Poppler        OK      ({pdftoppm})")
    else:
        results.append({
            "id": "Poppler",
            "status": "missing",
            "required": False,
            "detail": "needed for render QA",
        })
        if not args.json:
            print("  Poppler        MISSING (needed for render QA)")
    try:
        imported = __import__("playwright")
        ver = getattr(imported, "__version__", "unknown")
        results.append({
            "id": "Playwright",
            "module": "playwright",
            "status": "ok",
            "version": str(ver),
            "required": False,
            "detail": "needed for HTML-first PPTX screenshot QA",
        })
        if not args.json:
            print(f"  Playwright     OK      ({ver})")
    except ImportError:
        results.append({
            "id": "Playwright",
            "module": "playwright",
            "status": "missing",
            "required": False,
            "detail": "needed for HTML-first PPTX screenshot QA",
        })
        if not args.json:
            print("  Playwright     MISSING (needed for HTML-first PPTX screenshot QA)")
    except Exception as e:  # noqa: BLE001
        results.append({
            "id": "Playwright",
            "module": "playwright",
            "status": "broken",
            "required": False,
            "detail": f"{type(e).__name__}: {e}",
        })
        if not args.json:
            print(f"  Playwright     BROKEN  ({type(e).__name__}: {e})")
    office_backends = office_backend_statuses()
    if not args.json:
        print("\nOffice artifact backends:")
        for backend in office_backends:
            detail = f" ({backend['detail']})" if backend.get("detail") else ""
            print(f"  {backend['label']:<22} {backend['status'].upper()}{detail}")
    if args.json:
        payload = {
            "python": {
                "version": sys.version.split()[0],
                "executable": sys.executable,
            },
            "status": "missing" if missing_core else "degraded" if missing_optional else "ok",
            "requestedFormat": args.format,
            "dependencies": results,
            "officeBackends": office_backends,
        }
        print(json.dumps(payload, ensure_ascii=False, indent=2))
    if missing_core:
        if not args.json:
            print(
                "\nInstall missing deps with:\n"
                f"  python -m pip install {' '.join(missing_core)}"
            )
        return 2
    return 0


# ---------------------------------------------------------------------------
# extract
# ---------------------------------------------------------------------------

def _extract_docx(path: Path) -> str:
    try:
        import docx  # type: ignore
    except ImportError:
        _missing("python-docx")
    doc = docx.Document(str(path))
    return "\n".join(p.text for p in doc.paragraphs)


def _extract_pptx(path: Path, pages: str | None) -> str:
    try:
        from pptx import Presentation  # type: ignore
    except ImportError:
        _missing("python-pptx")
    prs = Presentation(str(path))
    indices = _parse_pages(pages, len(prs.slides))
    out = []
    for i in indices:
        slide = prs.slides[i]
        out.append(f"--- Slide {i + 1} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    out.append(para.text)
    return "\n".join(out)


def _extract_pdf(path: Path, pages: str | None) -> str:
    try:
        from pypdf import PdfReader  # type: ignore
    except ImportError:
        _missing("pypdf")
    reader = PdfReader(str(path))
    indices = _parse_pages(pages, len(reader.pages))
    out = []
    for i in indices:
        out.append(f"--- Page {i + 1} ---")
        out.append(reader.pages[i].extract_text() or "")
    return "\n".join(out)


def _extract_xlsx(path: Path, sheets: str | None) -> str:
    try:
        import openpyxl  # type: ignore
    except ImportError:
        _missing("openpyxl")
    wb = openpyxl.load_workbook(str(path), data_only=False, read_only=True)
    wanted = {s.strip() for s in sheets.split(",")} if sheets else None
    out: list[str] = []
    for ws in wb.worksheets:
        if wanted and ws.title not in wanted:
            continue
        out.append(f"--- Sheet: {ws.title} ---")
        for row in ws.iter_rows(values_only=True):
            values = ["" if cell is None else str(cell) for cell in row]
            if any(v.strip() for v in values):
                out.append("\t".join(values).rstrip())
    return "\n".join(out)


def cmd_extract(args: argparse.Namespace) -> int:
    path = _validate_path(args.path)
    ext = _ext(path)
    if ext in DOCX_EXTENSIONS:
        text = _extract_docx(path)
    elif ext in PPTX_EXTENSIONS:
        text = _extract_pptx(path, args.pages)
    elif ext == "pdf":
        text = _extract_pdf(path, args.pages)
    elif ext in XLSX_EXTENSIONS:
        text = _extract_xlsx(path, args.sheets)
    else:
        _die(f"ERROR: extract does not support .{ext}", 3)
    sys.stdout.write(_truncate(text))
    if not text.endswith("\n"):
        sys.stdout.write("\n")
    return 0


# ---------------------------------------------------------------------------
# replace / redact (shared core)
# ---------------------------------------------------------------------------

def _replace_across_text_nodes(
    nodes: list[Any],
    find: str,
    replace: str,
    apply: bool,
    occurrence: int | None = None,
) -> int:
    fragments = [node.text or "" for node in nodes]
    combined = "".join(fragments)
    if not combined or find not in combined:
        return 0

    starts: list[int] = []
    cursor = 0
    for fragment in fragments:
        starts.append(cursor)
        cursor += len(fragment)

    matches: list[tuple[int, int]] = []
    cursor = 0
    while True:
        start = combined.find(find, cursor)
        if start < 0:
            break
        matches.append((start, start + len(find)))
        cursor = start + len(find)

    if not apply:
        return len(matches)

    selected = matches if occurrence is None else (
        [matches[occurrence - 1]] if 1 <= occurrence <= len(matches) else []
    )

    for start, end in reversed(selected):
        start_index = max(0, bisect_right(starts, start) - 1)
        end_index = max(0, bisect_right(starts, end - 1) - 1)
        start_offset = start - starts[start_index]
        end_offset = end - starts[end_index]
        start_text = nodes[start_index].text or ""
        end_text = nodes[end_index].text or ""
        if start_index == end_index:
            nodes[start_index].text = start_text[:start_offset] + replace + start_text[end_offset:]
            continue
        nodes[start_index].text = start_text[:start_offset] + replace
        for index in range(start_index + 1, end_index):
            nodes[index].text = ""
        nodes[end_index].text = end_text[end_offset:]
    return len(matches)


def _docx_group_scope(part_name: str, paragraph: Any) -> str:
    lowered = part_name.lower()
    if "/header" in lowered:
        return "header"
    if "/footer" in lowered:
        return "footer"
    if "/comments" in lowered:
        return "comments"
    if "/footnotes" in lowered:
        return "footnotes"
    if "/endnotes" in lowered:
        return "endnotes"
    ancestors = list(paragraph.iterancestors()) if hasattr(paragraph, "iterancestors") else []
    ancestor_names = {ancestor.tag.rsplit("}", 1)[-1] for ancestor in ancestors}
    if "txbxContent" in ancestor_names:
        return "textbox"
    if "tc" in ancestor_names:
        return "table"
    return "body"


def _docx_text_groups(doc, scopes: set[str] | None = None) -> list[dict[str, Any]]:
    from docx.oxml.ns import qn  # type: ignore

    groups: list[dict[str, Any]] = []
    seen_roots: set[int] = set()
    for part in doc.part.package.parts:
        root = getattr(part, "element", None)
        if root is None:
            root = getattr(part, "_element", None)
        if root is None or id(root) in seen_roots:
            continue
        seen_roots.add(id(root))
        part_name = str(getattr(part, "partname", ""))
        for paragraph in root.iter(qn("w:p")):
            nodes = list(paragraph.iter(qn("w:t")))
            if nodes:
                scope = _docx_group_scope(part_name, paragraph)
                if scopes is None or scope in scopes:
                    groups.append({
                        "scope": scope,
                        "part": part_name,
                        "nodes": nodes,
                    })
    return groups


def _docx_run_signature(node: Any) -> tuple[str, bool]:
    run = node.getparent() if hasattr(node, "getparent") else None
    properties = ""
    if run is not None:
        for child in run:
            if child.tag.rsplit("}", 1)[-1] == "rPr":
                properties = str(getattr(child, "xml", ""))
                break
    hyperlink = False
    parent = run.getparent() if run is not None and hasattr(run, "getparent") else None
    while parent is not None:
        if parent.tag.rsplit("}", 1)[-1] == "hyperlink":
            hyperlink = True
            break
        parent = parent.getparent() if hasattr(parent, "getparent") else None
    return properties, hyperlink


def _assert_docx_container_compatible(
    nodes: list[Any],
    find: str,
    occurrence: int | None,
    allow_style_merge: bool,
) -> None:
    if allow_style_merge:
        return
    fragments = [node.text or "" for node in nodes]
    combined = "".join(fragments)
    starts: list[int] = []
    cursor = 0
    for fragment in fragments:
        starts.append(cursor)
        cursor += len(fragment)
    matches: list[tuple[int, int]] = []
    cursor = 0
    while True:
        start = combined.find(find, cursor)
        if start < 0:
            break
        matches.append((start, start + len(find)))
        cursor = start + len(find)
    selected = matches if occurrence is None else (
        [matches[occurrence - 1]] if 1 <= occurrence <= len(matches) else []
    )
    for start, end in selected:
        start_index = max(0, bisect_right(starts, start) - 1)
        end_index = max(0, bisect_right(starts, end - 1) - 1)
        signatures = {
            _docx_run_signature(nodes[index])
            for index in range(start_index, end_index + 1)
        }
        if len(signatures) > 1:
            _die(
                "CONTAINER_BOUNDARY: replacement crosses incompatible run style or hyperlink boundaries; "
                "target a narrower span or explicitly pass --allow-style-merge",
                1,
            )


def _assert_expected_match_count(count: int, expected_count: int | None) -> None:
    if expected_count is not None and count != expected_count:
        _die(
            f"PRECONDITION_FAILED: expected {expected_count} text match(es), found {count}",
            1,
        )


def _replace_docx(
    path: Path,
    find: str,
    replace: str,
    dry_run: bool,
    expected_count: int | None = None,
    scopes: set[str] | None = None,
    occurrence: int | None = None,
    allow_style_merge: bool = False,
) -> int:
    try:
        import docx  # type: ignore
    except ImportError:
        _missing("python-docx")
    working = path if dry_run else _staged_copy(path)
    doc = docx.Document(str(working))
    groups = _docx_text_groups(doc, scopes)
    before_lines = [
        f"[{group['scope']}] " + "".join(node.text or "" for node in group["nodes"])
        for group in groups
    ]
    group_counts = [
        _replace_across_text_nodes(group["nodes"], find, replace, apply=False)
        for group in groups
    ]
    count = sum(group_counts)
    if expected_count is not None and count != expected_count and not dry_run:
        working.unlink(missing_ok=True)
    _assert_expected_match_count(count, expected_count)
    if occurrence is not None and not 1 <= occurrence <= count:
        if not dry_run:
            working.unlink(missing_ok=True)
        _die(f"PRECONDITION_FAILED: occurrence {occurrence} is outside 1..{count}", 1)
    if not dry_run and count:
        global_start = 1
        for group, group_count in zip(groups, group_counts, strict=True):
            local_occurrence = None
            should_apply = occurrence is None
            if occurrence is not None and global_start <= occurrence < global_start + group_count:
                local_occurrence = occurrence - global_start + 1
                should_apply = True
            if should_apply and group_count:
                _assert_docx_container_compatible(
                    group["nodes"],
                    find,
                    local_occurrence,
                    allow_style_merge,
                )
                _replace_across_text_nodes(
                    group["nodes"],
                    find,
                    replace,
                    apply=True,
                    occurrence=local_occurrence,
                )
            global_start += group_count
    if dry_run:
        after_lines = [line.replace(find, replace) for line in before_lines]
        diff = difflib.unified_diff(
            before_lines, after_lines,
            fromfile=str(path), tofile=f"{path} (preview)", lineterm="",
        )
        sys.stdout.write("\n".join(diff) + "\n")
        print(f"\n[DRY-RUN] matches: {count}")
        return 0
    if count == 0:
        working.unlink(missing_ok=True)
        print(f"replaced 0 occurrence(s) in {path}")
        return 0
    doc.save(str(working))
    snapshot, validation = _publish_edit(working, path)
    print(json.dumps({
        "replaced": count,
        "path": str(path),
        "snapshot": str(snapshot) if snapshot else None,
        "validation": validation,
    }, ensure_ascii=False, indent=2))
    return 0


def _pptx_text_groups(prs) -> list[list[Any]]:
    drawing_ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    paragraph_tag = f"{{{drawing_ns}}}p"
    text_tag = f"{{{drawing_ns}}}t"
    groups: list[list[Any]] = []
    seen_roots: set[int] = set()
    for part in prs.part.package.iter_parts():
        root = getattr(part, "element", None)
        if root is None:
            root = getattr(part, "_element", None)
        if root is None or id(root) in seen_roots:
            continue
        seen_roots.add(id(root))
        for paragraph in root.iter(paragraph_tag):
            nodes = list(paragraph.iter(text_tag))
            if nodes:
                groups.append(nodes)
    return groups


def _replace_pptx(
    path: Path,
    find: str,
    replace: str,
    dry_run: bool,
    expected_count: int | None = None,
) -> int:
    try:
        from pptx import Presentation  # type: ignore
    except ImportError:
        _missing("python-pptx")
    working = path if dry_run else _staged_copy(path)
    prs = Presentation(str(working))
    groups = _pptx_text_groups(prs)
    before_lines = ["".join(node.text or "" for node in nodes) for nodes in groups]
    count = sum(_replace_across_text_nodes(nodes, find, replace, apply=False) for nodes in groups)
    if expected_count is not None and count != expected_count and not dry_run:
        working.unlink(missing_ok=True)
    _assert_expected_match_count(count, expected_count)
    if not dry_run and count:
        for nodes in groups:
            _replace_across_text_nodes(nodes, find, replace, apply=True)
    if dry_run:
        before = "\n".join(before_lines)
        after = before.replace(find, replace)
        diff = difflib.unified_diff(
            before.splitlines(), after.splitlines(),
            fromfile=str(path), tofile=f"{path} (preview)", lineterm="",
        )
        sys.stdout.write("\n".join(diff) + "\n")
        print(f"\n[DRY-RUN] matches: {count}")
        return 0
    if count == 0:
        working.unlink(missing_ok=True)
        print(f"replaced 0 occurrence(s) in {path}")
        return 0
    prs.save(str(working))
    snapshot, validation = _publish_edit(working, path)
    print(json.dumps({
        "replaced": count,
        "path": str(path),
        "snapshot": str(snapshot) if snapshot else None,
        "validation": validation,
    }, ensure_ascii=False, indent=2))
    return 0


def _replace_xlsx(
    path: Path,
    find: str,
    replace: str,
    dry_run: bool,
    expected_count: int | None = None,
) -> int:
    spreadsheet_ns = "http://schemas.openxmlformats.org/spreadsheetml/2006/main"
    text_tag = f"{{{spreadsheet_ns}}}t"
    container_tags = {
        f"{{{spreadsheet_ns}}}si",
        f"{{{spreadsheet_ns}}}is",
        f"{{{spreadsheet_ns}}}comment",
    }
    scalar_tags = {
        f"{{{spreadsheet_ns}}}f",
        f"{{{spreadsheet_ns}}}definedName",
    }
    editable_parts = (
        "xl/sharedStrings.xml",
        "xl/workbook.xml",
        "xl/worksheets/",
        "xl/comments",
    )
    before_lines: list[str] = []
    after_lines: list[str] = []
    count = 0
    staged = _staged_copy(path) if not dry_run else None
    source = staged or path
    repacked = staging_path(source) if not dry_run else None
    try:
        with zipfile.ZipFile(source) as archive:
            destination = zipfile.ZipFile(repacked, "w") if repacked is not None else None
            try:
                for info in archive.infolist():
                    data = archive.read(info.filename)
                    if info.filename.startswith(editable_parts) and info.filename.endswith(".xml"):
                        root = ET.fromstring(data)
                        groups: list[list[Any]] = []
                        part_matches = 0
                        for element in root.iter():
                            if element.tag in container_tags:
                                nodes = list(element.iter(text_tag))
                                if nodes:
                                    groups.append(nodes)
                            elif element.tag in scalar_tags:
                                groups.append([element])
                        for nodes in groups:
                            before = "".join(node.text or "" for node in nodes)
                            matches = _replace_across_text_nodes(
                                nodes, find, replace, apply=not dry_run
                            )
                            if matches:
                                count += matches
                                part_matches += matches
                                before_lines.append(f"{info.filename}: {before}")
                                after_lines.append(
                                    f"{info.filename}: {before.replace(find, replace)}"
                                )
                        if not dry_run and part_matches:
                            data = ET.tostring(root, encoding="utf-8", xml_declaration=True)
                    if destination is not None:
                        destination.writestr(info, data)
            finally:
                if destination is not None:
                    destination.close()
        if not dry_run and repacked is not None and staged is not None:
            os.replace(repacked, staged)
    except Exception:
        if repacked is not None:
            repacked.unlink(missing_ok=True)
        if staged is not None:
            staged.unlink(missing_ok=True)
        raise
    if dry_run:
        _assert_expected_match_count(count, expected_count)
        diff = difflib.unified_diff(
            before_lines, after_lines,
            fromfile=str(path), tofile=f"{path} (preview)", lineterm="",
        )
        sys.stdout.write("\n".join(diff) + "\n")
        print(f"\n[DRY-RUN] matches: {count}")
        return 0
    if expected_count is not None and count != expected_count:
        staged.unlink(missing_ok=True)
        _assert_expected_match_count(count, expected_count)
    if count == 0:
        staged.unlink(missing_ok=True)
        print(f"replaced 0 occurrence(s) in {path}")
        return 0
    snapshot, validation = _publish_edit(staged, path)
    print(json.dumps({
        "replaced": count,
        "path": str(path),
        "snapshot": str(snapshot) if snapshot else None,
        "preservationRisk": scan_ooxml_risks(path),
        "validation": validation,
    }, ensure_ascii=False, indent=2))
    return 0


def cmd_replace(args: argparse.Namespace) -> int:
    path = _validate_path(args.path)
    if not args.find:
        _die("ERROR: --find is required", 3)
    expected_sha256 = getattr(args, "expected_sha256", None)
    if expected_sha256:
        actual_sha256 = hashlib.sha256(path.read_bytes()).hexdigest()
        if actual_sha256.lower() != str(expected_sha256).lower():
            _die(
                "PRECONDITION_FAILED: artifact SHA-256 changed since it was inspected "
                f"(expected {expected_sha256}, actual {actual_sha256})",
                1,
            )
    expected_count = getattr(args, "expected_count", None)
    if expected_count is not None and expected_count < 0:
        _die("ERROR: --expected-count cannot be negative", 3)
    ext = _ext(path)
    if ext in DOCX_EXTENSIONS:
        raw_scope = getattr(args, "scope", None)
        scopes = {item.strip().lower() for item in raw_scope.split(",") if item.strip()} if raw_scope else None
        supported_scopes = {"body", "table", "textbox", "header", "footer", "comments", "footnotes", "endnotes"}
        if scopes is not None and not scopes <= supported_scopes:
            _die(f"ERROR: unsupported DOCX scope(s): {', '.join(sorted(scopes - supported_scopes))}", 3)
        return _replace_docx(
            path,
            args.find,
            args.replace or "",
            args.dry_run,
            expected_count,
            scopes,
            getattr(args, "occurrence", None),
            bool(getattr(args, "allow_style_merge", False)),
        )
    if ext in PPTX_EXTENSIONS:
        return _replace_pptx(path, args.find, args.replace or "", args.dry_run, expected_count)
    if ext in XLSX_EXTENSIONS:
        return _replace_xlsx(path, args.find, args.replace or "", args.dry_run, expected_count)
    _die(f"ERROR: replace supports .docx/.pptx/.xlsx only (got .{ext})", 3)
    return 1


def cmd_redact(args: argparse.Namespace) -> int:
    # Compatibility command: visible-story text replacement only. Use
    # secure_redact when the original must be proven absent from package text.
    args.replace = args.replace if args.replace is not None else "[REDACTED]"
    return cmd_replace(args)


def _replace_secret_in_xml(data: bytes, secret: str, replacement: str) -> tuple[bytes, int]:
    root = ET.fromstring(data)
    count = 0
    text_tags = {"t", "delText", "instrText"}
    for paragraph in [element for element in root.iter() if element.tag.rsplit("}", 1)[-1] == "p"]:
        nodes = [
            element for element in paragraph.iter()
            if element.tag.rsplit("}", 1)[-1] in text_tags
        ]
        if nodes:
            matches = _replace_across_text_nodes(nodes, secret, replacement, apply=True)
            count += matches
    for element in root.iter():
        local = element.tag.rsplit("}", 1)[-1]
        if local not in text_tags and element.text and secret in element.text:
            occurrences = element.text.count(secret)
            element.text = element.text.replace(secret, replacement)
            count += occurrences
        if element.tail and secret in element.tail:
            occurrences = element.tail.count(secret)
            element.tail = element.tail.replace(secret, replacement)
            count += occurrences
        for attribute, value in list(element.attrib.items()):
            if secret in value:
                occurrences = value.count(secret)
                element.set(attribute, value.replace(secret, replacement))
                count += occurrences
    if count == 0:
        return data, 0
    return ET.tostring(root, encoding="utf-8", xml_declaration=True), count


def _privacy_scrub_docx_part(name: str, data: bytes) -> bytes:
    if name == "docProps/custom.xml":
        root = ET.fromstring(data)
        for child in list(root):
            root.remove(child)
        return ET.tostring(root, encoding="utf-8", xml_declaration=True)
    if name != "docProps/core.xml":
        return data
    root = ET.fromstring(data)
    scrub_names = {"creator", "lastModifiedBy", "keywords", "description", "subject"}
    for element in root.iter():
        if element.tag.rsplit("}", 1)[-1] in scrub_names:
            element.text = ""
    return ET.tostring(root, encoding="utf-8", xml_declaration=True)


def cmd_secure_redact(args: argparse.Namespace) -> int:
    path = _validate_path(args.path)
    if _ext(path) != "docx":
        _die("ERROR: secure_redact currently requires a .docx file", 3)
    secret = str(args.find or "")
    if not secret:
        _die("ERROR: --find is required", 3)
    replacement = args.replace if args.replace is not None else "[REDACTED]"
    expected_sha256 = getattr(args, "expected_sha256", None)
    if expected_sha256:
        actual_sha256 = hashlib.sha256(path.read_bytes()).hexdigest()
        if actual_sha256.lower() != str(expected_sha256).lower():
            _die(
                "PRECONDITION_FAILED: artifact SHA-256 changed since it was inspected "
                f"(expected {expected_sha256}, actual {actual_sha256})",
                1,
            )
    expected_count = getattr(args, "expected_count", None)
    blocked_prefixes = ("word/embeddings/", "word/oleObjects/", "word/media/")
    with zipfile.ZipFile(path) as archive:
        blocked_parts = sorted(
            name for name in archive.namelist()
            if not name.endswith("/") and name.startswith(blocked_prefixes)
        )
    if blocked_parts:
        _die(
            "UNINSPECTABLE_CONTENT: secure redaction cannot prove that the secret is absent from "
            "embedded objects or media: " + ", ".join(blocked_parts),
            1,
        )
    staged = staging_path(path)
    redacted_parts: dict[str, int] = {}
    total = 0
    try:
        with zipfile.ZipFile(path) as source, zipfile.ZipFile(staged, "w") as destination:
            for info in source.infolist():
                data = source.read(info.filename)
                if info.filename.endswith((".xml", ".rels")):
                    try:
                        data, matches = _replace_secret_in_xml(
                            data,
                            secret,
                            str(replacement),
                        )
                    except ET.ParseError as error:
                        _die(f"SECURE_REDACTION_FAILED: XML part cannot be inspected: {info.filename}: {error}", 1)
                    if matches:
                        redacted_parts[info.filename] = matches
                        total += matches
                    if getattr(args, "privacy_scrub", False):
                        data = _privacy_scrub_docx_part(info.filename, data)
                destination.writestr(info, data)
        _assert_expected_match_count(total, expected_count)
        if total == 0:
            _die("PRECONDITION_FAILED: secure redaction found no matches", 1)
        secret_encodings = [secret.encode("utf-8"), secret.encode("utf-16le")]
        residual_parts: list[str] = []
        with zipfile.ZipFile(staged) as archive:
            for name in archive.namelist():
                data = archive.read(name)
                if any(needle and needle in data for needle in secret_encodings):
                    residual_parts.append(name)
        if residual_parts:
            _die(
                "SECURE_REDACTION_FAILED: original text remains in package parts: "
                + ", ".join(residual_parts),
                1,
            )
        previous_snapshot_policy = os.environ.get("NEXA_OFFICE_SKIP_SNAPSHOT")
        os.environ["NEXA_OFFICE_SKIP_SNAPSHOT"] = "1"
        try:
            snapshot, validation = _publish_edit(staged, path)
        finally:
            if previous_snapshot_policy is None:
                os.environ.pop("NEXA_OFFICE_SKIP_SNAPSHOT", None)
            else:
                os.environ["NEXA_OFFICE_SKIP_SNAPSHOT"] = previous_snapshot_policy
    finally:
        staged.unlink(missing_ok=True)
    print(json.dumps({
        "kind": "secureDocxRedaction",
        "path": str(path),
        "redactedOccurrences": total,
        "redactedParts": redacted_parts,
        "privacyScrubbed": bool(getattr(args, "privacy_scrub", False)),
        "verification": {
            "utf8ResidualParts": [],
            "utf16ResidualParts": [],
            "uninspectableParts": [],
            "originalTextAbsent": True,
            "scope": "final artifact package only",
        },
        "snapshot": str(snapshot) if snapshot else None,
        "validation": validation,
    }, ensure_ascii=False, indent=2))
    return 0


# ---------------------------------------------------------------------------
# insert_slide
# ---------------------------------------------------------------------------

def cmd_insert_slide(args: argparse.Namespace) -> int:
    path = _validate_path(args.path)
    if _ext(path) not in PPTX_EXTENSIONS:
        _die("ERROR: insert_slide requires a PowerPoint OOXML package", 3)
    try:
        from pptx import Presentation  # type: ignore
        from pptx.util import Inches  # type: ignore
    except ImportError:
        _missing("python-pptx")
    staged = _staged_copy(path)
    prs = Presentation(str(staged))
    layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    # Populate title/body if placeholders exist.
    if slide.shapes.title is not None:
        slide.shapes.title.text = args.title or ""
    if args.body:
        body_placeholder = None
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 1:
                body_placeholder = ph
                break
        if body_placeholder is None:
            left = top = Inches(1)
            width = Inches(8)
            height = Inches(5)
            body_placeholder = slide.shapes.add_textbox(left, top, width, height)
        body_placeholder.text_frame.text = args.body

    # Reorder: move new slide to position after --after.
    after = max(0, int(args.after))
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    new_el = slides[-1]
    xml_slides.remove(new_el)
    xml_slides.insert(after, new_el)

    prs.save(str(staged))
    snapshot, validation = _publish_edit(staged, path)
    print(json.dumps({
        "insertedAfter": after,
        "path": str(path),
        "snapshot": str(snapshot) if snapshot else None,
        "validation": validation,
    }, ensure_ascii=False, indent=2))
    return 0


# ---------------------------------------------------------------------------
# version
# ---------------------------------------------------------------------------

def cmd_version(args: argparse.Namespace) -> int:
    path = _validate_path(args.path)
    root = Path.cwd() / HISTORY_DIR / path.name
    root.mkdir(parents=True, exist_ok=True)
    existing = sorted(
        (p for p in root.iterdir() if p.is_dir() and p.name.startswith("v")),
        key=lambda p: int(p.name[1:]) if p.name[1:].isdigit() else 0,
    )
    next_n = 1
    if existing:
        last = existing[-1].name
        try:
            next_n = int(last[1:]) + 1
        except ValueError:
            next_n = len(existing) + 1
    dest_dir = root / f"v{next_n}"
    dest_dir.mkdir()
    dest = dest_dir / path.name
    shutil.copy2(path, dest)
    print(f"v{next_n} -> {dest}")
    return 0


# ---------------------------------------------------------------------------
# create_docx / create_xlsx / create_pptx
# ---------------------------------------------------------------------------

def _numbered_markdown_text(stripped: str) -> str | None:
    prefix, sep, rest = stripped.partition(".")
    if sep and prefix.isdigit() and rest.startswith(" "):
        return rest.strip()
    return None


def _docx_add_hyperlink(paragraph, label: str, url: str) -> None:
    from docx.opc.constants import RELATIONSHIP_TYPE as RT  # type: ignore
    from docx.oxml import OxmlElement  # type: ignore
    from docx.oxml.ns import qn  # type: ignore

    relationship_id = paragraph.part.relate_to(url, RT.HYPERLINK, is_external=True)
    hyperlink = OxmlElement("w:hyperlink")
    hyperlink.set(qn("r:id"), relationship_id)
    run = OxmlElement("w:r")
    properties = OxmlElement("w:rPr")
    color = OxmlElement("w:color")
    color.set(qn("w:val"), "0563C1")
    underline = OxmlElement("w:u")
    underline.set(qn("w:val"), "single")
    properties.extend([color, underline])
    text = OxmlElement("w:t")
    text.text = label
    run.extend([properties, text])
    hyperlink.append(run)
    paragraph._p.append(hyperlink)


def _docx_add_inline_markdown(paragraph, text: str) -> None:
    """Add common inline Markdown as Word runs instead of literal markers."""
    i = 0
    plain_start = 0

    def flush_plain(end: int) -> None:
        nonlocal plain_start
        if end > plain_start:
            paragraph.add_run(text[plain_start:end])
        plain_start = end

    def add_styled(content: str, *, bold=False, italic=False, strike=False, code=False) -> None:
        if not content:
            return
        run = paragraph.add_run(content)
        run.bold = bold or None
        run.italic = italic or None
        run.font.strike = strike or None
        if code:
            run.font.name = "Consolas"

    while i < len(text):
        if text[i] == "\\" and i + 1 < len(text):
            flush_plain(i)
            paragraph.add_run(text[i + 1])
            i += 2
            plain_start = i
            continue

        if text.startswith(("**", "__"), i):
            marker = text[i:i + 2]
            end = text.find(marker, i + 2)
            if end != -1:
                flush_plain(i)
                add_styled(text[i + 2:end], bold=True)
                i = end + 2
                plain_start = i
                continue

        if text.startswith("~~", i):
            end = text.find("~~", i + 2)
            if end != -1:
                flush_plain(i)
                add_styled(text[i + 2:end], strike=True)
                i = end + 2
                plain_start = i
                continue

        if text[i] == "`":
            end = text.find("`", i + 1)
            if end != -1:
                flush_plain(i)
                add_styled(text[i + 1:end], code=True)
                i = end + 1
                plain_start = i
                continue

        if text[i] in {"*", "_"}:
            marker = text[i]
            end = text.find(marker, i + 1)
            if end != -1:
                flush_plain(i)
                add_styled(text[i + 1:end], italic=True)
                i = end + 1
                plain_start = i
                continue

        if text[i] == "[":
            label_end = text.find("](", i + 1)
            if label_end != -1:
                url_end = text.find(")", label_end + 2)
                if url_end != -1:
                    flush_plain(i)
                    label = text[i + 1:label_end]
                    url = text[label_end + 2:url_end]
                    if url and url.lower().startswith(("https://", "http://", "mailto:")):
                        _docx_add_hyperlink(paragraph, label or url, url)
                    elif label:
                        paragraph.add_run(label)
                        if url:
                            paragraph.add_run(f" ({url})")
                    elif url:
                        paragraph.add_run(url)
                    i = url_end + 1
                    plain_start = i
                    continue

        i += 1

    flush_plain(len(text))


def _docx_add_markdown(doc, markdown: str) -> None:
    lines = markdown.splitlines()
    i = 0
    while i < len(lines):
        line = lines[i]
        stripped = line.strip()
        if not stripped:
            doc.add_paragraph()
            i += 1
            continue
        if stripped.startswith("#"):
            level = len(stripped) - len(stripped.lstrip("#"))
            if 1 <= level <= 6 and len(stripped) > level and stripped[level].isspace():
                heading = doc.add_heading(level=min(level, 4))
                _docx_add_inline_markdown(heading, stripped[level:].strip())
                i += 1
                continue
        if stripped.startswith("|") and stripped.endswith("|"):
            table_lines: list[str] = []
            while i < len(lines) and lines[i].strip().startswith("|") and lines[i].strip().endswith("|"):
                table_lines.append(lines[i].strip())
                i += 1
            rows = [[cell.strip() for cell in row.strip("|").split("|")] for row in table_lines]
            rows = [
                row for row in rows
                if not all(cell and set(cell) <= {"-", ":", " "} for cell in row)
            ]
            if rows:
                table = doc.add_table(rows=len(rows), cols=max(len(r) for r in rows))
                table.style = "Table Grid"
                for ri, row in enumerate(rows):
                    for ci, value in enumerate(row):
                        cell = table.rows[ri].cells[ci]
                        cell.text = ""
                        _docx_add_inline_markdown(cell.paragraphs[0], value)
                        if ri == 0:
                            for paragraph in table.rows[ri].cells[ci].paragraphs:
                                for run in paragraph.runs:
                                    run.bold = True
                continue
        numbered = _numbered_markdown_text(stripped)
        if stripped.startswith(("- ", "* ", "• ")):
            p = doc.add_paragraph(style="List Bullet")
            _docx_add_inline_markdown(p, stripped[2:].strip())
        elif numbered is not None:
            p = doc.add_paragraph(style="List Number")
            _docx_add_inline_markdown(p, numbered)
        elif stripped.startswith("> "):
            p = doc.add_paragraph()
            p.style = "Intense Quote"
            _docx_add_inline_markdown(p, stripped[2:].strip())
        else:
            p = doc.add_paragraph()
            _docx_add_inline_markdown(p, stripped)
        i += 1


def cmd_create_docx(args: argparse.Namespace) -> int:
    try:
        import docx  # type: ignore
    except ImportError:
        _missing("python-docx")
    path = _validate_output_path(args.path, {"docx"})
    staged = staging_path(path)
    if getattr(args, "spec", None):
        spec_path = _validate_path(args.spec)
        payload = _read_json(str(spec_path))
        skills_root = Path(__file__).resolve().parents[2]
        renderer_dir = skills_root / "docx-document-design" / "scripts"
        if str(renderer_dir) not in sys.path:
            sys.path.insert(0, str(renderer_dir))
        try:
            from docx_renderer import render_docx  # type: ignore
        except ImportError as exc:
            _die(f"ERROR: failed to load DOCX renderer: {exc}", 1)
        try:
            result = render_docx(payload, staged, Path.cwd())
            snapshot, validation = _publish_edit(staged, path)
        finally:
            staged.unlink(missing_ok=True)
        result.update({
            "path": str(path),
            "snapshot": str(snapshot) if snapshot else None,
            "validation": validation,
        })
        print(json.dumps(result, ensure_ascii=False, indent=2))
        return 0
    doc = docx.Document(str(_validate_path(args.template))) if args.template else docx.Document()
    if args.font:
        doc.styles["Normal"].font.name = args.font
    if args.title:
        doc.add_heading(args.title, 0)
    if args.subtitle:
        p = doc.add_paragraph(args.subtitle)
        p.style = "Subtitle" if "Subtitle" in [s.name for s in doc.styles] else p.style
    if args.input_md:
        _docx_add_markdown(doc, _read_text(args.input_md))
    elif args.body:
        _docx_add_markdown(doc, args.body)
    if args.footer:
        for section in doc.sections:
            section.footer.paragraphs[0].text = args.footer
    if args.author:
        doc.core_properties.author = args.author
    try:
        doc.save(str(staged))
        snapshot, validation = _publish_edit(staged, path)
    finally:
        staged.unlink(missing_ok=True)
    print(json.dumps({
        "created": str(path),
        "snapshot": str(snapshot) if snapshot else None,
        "validation": validation,
    }, ensure_ascii=False, indent=2))
    return 0


def cmd_create_xlsx(args: argparse.Namespace) -> int:
    path = _validate_output_path(args.path, {"xlsx"})
    staged = staging_path(path)
    staged_qa = staged.with_suffix(".xlsx.qa.json")
    spec_path = _validate_path(args.spec)
    create_xlsx_from_spec, _, _ = _load_xlsx_renderer()
    try:
        result = create_xlsx_from_spec(
            staged,
            spec_path,
            workspace_root=Path.cwd(),
        )
        if result["qa"]["status"] == "fail":
            staged.unlink(missing_ok=True)
            staged_qa.unlink(missing_ok=True)
            print(json.dumps(result, ensure_ascii=False, indent=2))
            return 4
        snapshot, validation = _publish_edit(staged, path)
    except BaseException:
        staged_qa.unlink(missing_ok=True)
        raise
    finally:
        staged.unlink(missing_ok=True)
    final_qa = path.with_suffix(".xlsx.qa.json")
    result["path"] = str(path)
    result["qaPath"] = str(final_qa)
    result["artifact"] = {
        "path": str(path),
        "snapshot": str(snapshot) if snapshot else None,
        "validation": validation,
    }
    if staged_qa.exists():
        write_artifact_manifest(final_qa, result, Path.cwd())
        staged_qa.unlink(missing_ok=True)
    print(json.dumps(result, ensure_ascii=False, indent=2))
    return 0


def _load_xlsx_renderer():
    skills_root = Path(__file__).resolve().parents[2]
    renderer_dir = skills_root / "xlsx-workbook-design" / "scripts"
    renderer_path = renderer_dir / "xlsx_model_renderer.py"
    if not renderer_path.exists():
        _die(f"ERROR: XLSX renderer not found: {renderer_path}", 3)
    if str(renderer_dir) not in sys.path:
        sys.path.insert(0, str(renderer_dir))
    try:
        from xlsx_model_renderer import (  # type: ignore
            audit_xlsx_formula_integrity,
            create_xlsx_from_spec,
            inspect_formula_cache,
        )
    except ImportError as exc:
        _die(f"ERROR: failed to load XLSX renderer: {exc}", 1)
    return create_xlsx_from_spec, audit_xlsx_formula_integrity, inspect_formula_cache


def _load_xlsx_structured_editor():
    skills_root = Path(__file__).resolve().parents[2]
    editor_dir = skills_root / "xlsx-workbook-design" / "scripts"
    if str(editor_dir) not in sys.path:
        sys.path.insert(0, str(editor_dir))
    try:
        from xlsx_structured_editor import patch_xlsx  # type: ignore
    except ImportError as exc:
        _die(f"ERROR: failed to load XLSX structured editor: {exc}", 1)
    return patch_xlsx


def cmd_edit_xlsx(args: argparse.Namespace) -> int:
    path = _validate_path(args.path)
    if _ext(path) not in XLSX_EXTENSIONS:
        _die("ERROR: edit_xlsx requires an Excel OOXML package", 3)
    payload = _read_json(args.spec)
    operations = payload.get("operations", payload)
    if not isinstance(operations, list) or not all(isinstance(item, dict) for item in operations):
        _die("ERROR: XLSX edit spec must contain an operations array", 3)
    staged = staging_path(path)
    patch_xlsx = _load_xlsx_structured_editor()
    try:
        result = patch_xlsx(path, staged, operations)
        snapshot, validation = _publish_edit(staged, path)
    finally:
        staged.unlink(missing_ok=True)
    result.update({
        "path": str(path),
        "snapshot": str(snapshot) if snapshot else None,
        "validation": validation,
    })
    print(json.dumps(result, ensure_ascii=False, indent=2))
    return 0


def _load_pptx_renderer():
    skills_root = Path(__file__).resolve().parents[2]
    renderer_dir = skills_root / "pptx-presentation-design" / "scripts"
    renderer_path = renderer_dir / "pptx_renderer.py"
    if not renderer_path.exists():
        _die(f"ERROR: PPTX renderer not found: {renderer_path}", 3)
    if str(renderer_dir) not in sys.path:
        sys.path.insert(0, str(renderer_dir))
    try:
        from pptx_renderer import create_pptx_from_spec  # type: ignore
    except ImportError as exc:
        _die(f"ERROR: failed to load PPTX renderer: {exc}", 1)
    return create_pptx_from_spec


def _load_pptx_structured_editor():
    skills_root = Path(__file__).resolve().parents[2]
    editor_dir = skills_root / "pptx-presentation-design" / "scripts"
    if str(editor_dir) not in sys.path:
        sys.path.insert(0, str(editor_dir))
    try:
        from pptx_structured_editor import patch_pptx  # type: ignore
    except ImportError as exc:
        _die(f"ERROR: failed to load PPTX structured editor: {exc}", 1)
    return patch_pptx


def _load_docx_review_editor():
    skills_root = Path(__file__).resolve().parents[2]
    editor_dir = skills_root / "docx-document-design" / "scripts"
    if str(editor_dir) not in sys.path:
        sys.path.insert(0, str(editor_dir))
    try:
        from docx_review_editor import extract_comments, patch_docx_reviews  # type: ignore
    except ImportError as exc:
        _die(f"ERROR: failed to load DOCX review editor: {exc}", 1)
    return patch_docx_reviews, extract_comments


def cmd_review_docx(args: argparse.Namespace) -> int:
    path = _validate_path(args.path)
    if _ext(path) not in DOCX_EXTENSIONS:
        _die("ERROR: review_docx requires a Word OOXML package", 3)
    payload = _read_json(args.spec)
    operations = payload.get("operations", payload)
    if not isinstance(operations, list) or not all(isinstance(item, dict) for item in operations):
        _die("ERROR: DOCX review spec must contain an operations array", 3)
    patch_docx_reviews, _ = _load_docx_review_editor()
    staged = staging_path(path)
    try:
        result = patch_docx_reviews(path, staged, operations)
        snapshot, validation = _publish_edit(staged, path)
    finally:
        staged.unlink(missing_ok=True)
    result.update({
        "path": str(path),
        "snapshot": str(snapshot) if snapshot else None,
        "validation": validation,
    })
    print(json.dumps(result, ensure_ascii=False, indent=2))
    return 0


def cmd_comments_docx(args: argparse.Namespace) -> int:
    path = _validate_path(args.path)
    if _ext(path) not in DOCX_EXTENSIONS:
        _die("ERROR: comments_docx requires a Word OOXML package", 3)
    _, extract_comments = _load_docx_review_editor()
    print(json.dumps(extract_comments(path), ensure_ascii=False, indent=2))
    return 0


def cmd_edit_pptx(args: argparse.Namespace) -> int:
    path = _validate_path(args.path)
    if _ext(path) not in PPTX_EXTENSIONS:
        _die("ERROR: edit_pptx requires a PowerPoint OOXML package", 3)
    payload = _read_json(args.spec)
    operations = payload.get("operations", payload)
    if not isinstance(operations, list) or not all(isinstance(item, dict) for item in operations):
        _die("ERROR: PPTX edit spec must contain an operations array", 3)
    staged = staging_path(path)
    patch_pptx = _load_pptx_structured_editor()
    try:
        result = patch_pptx(path, staged, operations)
        snapshot, validation = _publish_edit(staged, path)
    finally:
        staged.unlink(missing_ok=True)
    result.update({
        "path": str(path),
        "snapshot": str(snapshot) if snapshot else None,
        "validation": validation,
    })
    print(json.dumps(result, ensure_ascii=False, indent=2))
    return 0


def _load_html_deck_renderer():
    skills_root = Path(__file__).resolve().parents[2]
    renderer_dir = skills_root / "pptx-presentation-design" / "scripts"
    renderer_path = renderer_dir / "html_deck_renderer.py"
    if not renderer_path.exists():
        _die(f"ERROR: HTML deck renderer not found: {renderer_path}", 3)
    if str(renderer_dir) not in sys.path:
        sys.path.insert(0, str(renderer_dir))
    try:
        from html_deck_renderer import render_html_deck  # type: ignore
    except ImportError as exc:
        _die(f"ERROR: failed to load HTML deck renderer: {exc}", 1)
    return render_html_deck


def cmd_create_pptx(args: argparse.Namespace) -> int:
    path = _validate_output_path(args.path, {"pptx"})
    staged = staging_path(path)
    create_pptx_from_spec = _load_pptx_renderer()
    try:
        create_pptx_from_spec(str(staged), args.spec, args.template, Path.cwd())
        snapshot, validation = _publish_edit(staged, path)
    finally:
        staged.unlink(missing_ok=True)
    print(json.dumps({
        "created": str(path),
        "snapshot": str(snapshot) if snapshot else None,
        "validation": validation,
    }, ensure_ascii=False, indent=2))
    return 0


def cmd_create_html_pptx(args: argparse.Namespace) -> int:
    if not args.path:
        _die("ERROR: --path is required for create_html_pptx", 3)
    if args.spec != "-":
        args.spec = str(_validate_path(args.spec))
    args.outdir = str(_validate_output_dir(args.outdir))
    path = _validate_output_path(args.path, {"pptx"})
    staged = staging_path(path)
    render_html_deck = _load_html_deck_renderer()
    try:
        result = render_html_deck(
            spec_path=args.spec,
            out_dir=args.outdir,
            pptx_path=str(staged),
            mode=args.mode,
            screenshot=args.screenshot,
            workspace_root=Path.cwd(),
        )
        if result["qa"]["status"] == "fail":
            staged.unlink(missing_ok=True)
            print(json.dumps(result, ensure_ascii=False, indent=2))
            return 4
        snapshot, validation = _publish_edit(staged, path)
    finally:
        staged.unlink(missing_ok=True)
    result["artifact"] = {
        "path": str(path),
        "snapshot": str(snapshot) if snapshot else None,
        "validation": validation,
    }
    pptx_result = result.get("manifest", {}).get("pptx")
    if isinstance(pptx_result, dict):
        pptx_result["path"] = str(path)
        write_artifact_manifest(Path(result["manifestPath"]), result["manifest"], Path.cwd())
    print(json.dumps(result, ensure_ascii=False, indent=2))
    return 0


# ---------------------------------------------------------------------------
# ooxml / render / recalc / validate / convert
# ---------------------------------------------------------------------------

def cmd_unpack(args: argparse.Namespace) -> int:
    path = _validate_path(args.path)
    if _ext(path) not in OOXML_EXTENSIONS:
        _die("ERROR: unpack supports Word/Excel/PowerPoint OOXML packages only", 3)
    preflight = validate_ooxml_package(path)
    if preflight.status == "fail":
        _die(
            "VALIDATION_FAILED: package is unsafe or structurally invalid and was not unpacked\n"
            + json.dumps(preflight.to_dict(), ensure_ascii=False),
            1,
        )
    outdir = _validate_output_dir(args.outdir, allow_existing=True)
    with zipfile.ZipFile(path) as zf:
        archive_files = {
            member.filename.rstrip("/")
            for member in zf.infolist()
            if member.filename and not member.is_dir()
        }
        if any(outdir.iterdir()):
            if not args.overwrite:
                _die(f"ERROR: output directory is not empty: {outdir}. Pass --overwrite to replace it.", 3)
            marker = outdir / UNPACK_MARKER
            try:
                marker_payload = json.loads(marker.read_text(encoding="utf-8"))
            except (FileNotFoundError, OSError, json.JSONDecodeError):
                _die(
                    f"ERROR: refusing to overwrite unmanaged directory: {outdir}. "
                    f"Only directories created by this unpack command may be replaced.",
                    3,
                )
            if marker_payload.get("kind") != "nexa-ooxml-unpack" or marker_payload.get("version") not in {1, 2}:
                _die(f"ERROR: invalid managed-unpack marker in: {outdir}", 3)
            existing_files = {
                item.relative_to(outdir).as_posix()
                for item in outdir.rglob("*")
                if item.is_file() and item.name != UNPACK_MARKER
            }
            if existing_files != archive_files:
                _die(
                    "ERROR: refusing destructive managed-unpack overwrite because the existing "
                    "member set differs from the new package; choose a new empty outdir",
                    3,
                )
        for member in zf.infolist():
            target = (outdir / member.filename).resolve()
            try:
                target.relative_to(outdir)
            except ValueError:
                _die(f"ERROR: unsafe ZIP member escapes output directory: {member.filename}", 3)
            if member.is_dir():
                target.mkdir(parents=True, exist_ok=True)
                continue
            target.parent.mkdir(parents=True, exist_ok=True)
            with zf.open(member) as source, target.open("wb") as output:
                shutil.copyfileobj(source, output)
        xml_count = sum(
            1 for member in zf.infolist()
            if member.filename.lower().endswith((".xml", ".rels"))
        )
    (outdir / UNPACK_MARKER).write_text(
        json.dumps({
            "kind": "nexa-ooxml-unpack",
            "version": 2,
            "source": str(path),
            "sourceSha256": hashlib.sha256(path.read_bytes()).hexdigest(),
            "ownedEntries": sorted(archive_files),
        }, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )
    print(f"unpacked {path.name} -> {outdir} ({xml_count} XML parts preserved byte-for-byte)")
    return 0


def cmd_pack(args: argparse.Namespace) -> int:
    input_dir = _validate_path(args.input_dir)
    if not input_dir.is_dir():
        _die(f"ERROR: input directory is not a directory: {input_dir}", 3)
    if not (input_dir / "[Content_Types].xml").exists():
        _die("ERROR: input directory does not look like an unpacked Office document", 3)
    path = _validate_output_path(args.path, {"docx", "pptx", "xlsx"})
    staged = staging_path(path)
    try:
        with zipfile.ZipFile(staged, "w", zipfile.ZIP_DEFLATED) as zf:
            for item in sorted(input_dir.rglob("*")):
                if item.is_symlink():
                    _die(f"ERROR: refusing to pack symlinked OOXML content: {item}", 3)
                if item.is_file() and item.name != UNPACK_MARKER:
                    try:
                        item.resolve().relative_to(input_dir.resolve())
                    except ValueError:
                        _die(f"ERROR: OOXML pack member escapes input directory: {item}", 3)
                    zf.write(item, item.relative_to(input_dir).as_posix())
        snapshot, validation = _publish_edit(staged, path)
    finally:
        staged.unlink(missing_ok=True)
    print(json.dumps({
        "packedFrom": str(input_dir),
        "path": str(path),
        "snapshot": str(snapshot) if snapshot else None,
        "validation": validation,
    }, ensure_ascii=False, indent=2))
    return 0


def _convert_to_pdf(path: Path, outdir: Path) -> Path:
    completed = _run_soffice_convert(path, "pdf", outdir)
    if completed.returncode != 0:
        stderr = completed.stderr.strip() or completed.stdout.strip() or "LibreOffice conversion failed"
        _die(stderr, completed.returncode or 1)
    pdf = _expected_converted_path(path, outdir, "pdf")
    if not pdf.exists():
        matches = sorted(outdir.glob("*.pdf"))
        if matches:
            return matches[0]
        _die(f"ERROR: LibreOffice did not produce a PDF in {outdir}", 1)
    return pdf


def _xlsx_render_surfaces(path: Path, selection: str | None) -> list[str]:
    try:
        import openpyxl  # type: ignore
    except ImportError:
        _missing("openpyxl")
    workbook = openpyxl.load_workbook(str(path), read_only=True, data_only=False)
    try:
        visible = [sheet.title for sheet in workbook.worksheets if sheet.sheet_state == "visible"]
        if not visible:
            _die("ERROR: XLSX has no visible worksheets to render", 1)
        requested = str(selection or "all").strip()
        if requested == "all":
            return visible
        if requested == "active":
            active = workbook.active.title
            return [active] if active in visible else [visible[0]]
        names = [item.strip() for item in requested.split(",") if item.strip()]
        missing = [name for name in names if name not in visible]
        if missing:
            _die(f"ERROR: requested render sheet(s) are not visible: {', '.join(missing)}", 3)
        return names
    finally:
        workbook.close()


def _render_xlsx_surfaces(
    path: Path,
    outdir: Path,
    *,
    selection: str | None,
    image_format: str,
    dpi: int,
    pdftoppm: str,
) -> int:
    try:
        import openpyxl  # type: ignore
    except ImportError:
        _missing("openpyxl")
    surfaces = _xlsx_render_surfaces(path, selection)
    records: list[dict[str, Any]] = []
    with tempfile.TemporaryDirectory(prefix=".nexa-xlsx-render-", dir=path.parent) as tmp:
        temporary_root = Path(tmp)
        for index, sheet_name in enumerate(surfaces, start=1):
            workbook = openpyxl.load_workbook(str(path), data_only=False, keep_links=True)
            try:
                for worksheet in workbook.worksheets:
                    worksheet.sheet_state = "visible" if worksheet.title == sheet_name else "hidden"
                workbook.active = workbook.sheetnames.index(sheet_name)
                surface_path = temporary_root / f"surface-{index:03d}.xlsx"
                workbook.save(surface_path)
            finally:
                workbook.close()
            pdf_dir = temporary_root / f"pdf-{index:03d}"
            pdf_dir.mkdir()
            pdf = _convert_to_pdf(surface_path, pdf_dir)
            prefix = outdir / f"sheet-{index:03d}-page"
            command = [
                pdftoppm,
                f"-{image_format}",
                "-r",
                str(dpi),
                str(pdf),
                str(prefix),
            ]
            completed = _run_subprocess(command, text=True, capture_output=True, check=False)
            if completed.returncode != 0:
                _die(
                    completed.stderr.strip() or completed.stdout.strip() or f"render failed for sheet {sheet_name}",
                    completed.returncode or 1,
                )
            suffix = "jpg" if image_format == "jpeg" else "png"
            images = sorted(outdir.glob(f"sheet-{index:03d}-page*.{suffix}"))
            if not images:
                _die(f"ERROR: renderer produced no images for worksheet: {sheet_name}", 1)
            records.append({
                "surfaceId": f"worksheet:{sheet_name}",
                "sheet": sheet_name,
                "files": [
                    {
                        "path": str(image),
                        "sha256": hashlib.sha256(image.read_bytes()).hexdigest(),
                    }
                    for image in images
                ],
            })
    manifest = {
        "kind": "officeRenderManifest",
        "version": 1,
        "format": "xlsx",
        "renderer": "libreoffice-compatible-via-temporary-openpyxl-surface-copy",
        "artifact": str(path),
        "artifactSha256": hashlib.sha256(path.read_bytes()).hexdigest(),
        "selection": selection or "all",
        "expectedSurfaces": len(surfaces),
        "renderedSurfaces": len(records),
        "complete": len(records) == len(surfaces),
        "surfaces": records,
        "preservationNote": "The source workbook was never rewritten; temporary per-sheet copies may not preserve unsupported Excel-only features.",
    }
    write_artifact_manifest(outdir / "render-manifest.json", manifest, Path.cwd())
    print(json.dumps(manifest, ensure_ascii=False, indent=2))
    return 0


def cmd_render(args: argparse.Namespace) -> int:
    path = _validate_path(args.path)
    if _ext(path) not in OOXML_EXTENSIONS | {"pdf"}:
        _die("ERROR: render supports Office OOXML and PDF packages only", 3)
    pdftoppm = _find_pdftoppm()
    if not pdftoppm:
        _die("MISSING_DEP: Poppler/pdftoppm\nInstall Poppler and ensure pdftoppm is on PATH.", 2)
    outdir = _validate_output_dir(args.outdir, allow_existing=True)
    image_format = args.format.lower()
    if image_format not in {"png", "jpeg"}:
        _die("ERROR: --format must be png or jpeg", 3)
    for pattern in (
        "page*.png", "page*.jpg", "page*.jpeg", "page*.ppm",
        "sheet-*-page*.png", "sheet-*-page*.jpg", "sheet-*-page*.jpeg",
    ):
        for old_preview in outdir.glob(pattern):
            if old_preview.is_file():
                old_preview.unlink()
    (outdir / "render-manifest.json").unlink(missing_ok=True)
    if _ext(path) == "xlsx":
        return _render_xlsx_surfaces(
            path,
            outdir,
            selection=getattr(args, "sheets", None),
            image_format=image_format,
            dpi=args.dpi,
            pdftoppm=pdftoppm,
        )
    with tempfile.TemporaryDirectory(prefix="nexa-render-") as tmp:
        pdf = path if _ext(path) == "pdf" else _convert_to_pdf(path, Path(tmp))
        prefix = outdir / "page"
        cmd = [
            pdftoppm,
            f"-{image_format}",
            "-r",
            str(args.dpi),
            str(pdf),
            str(prefix),
        ]
        completed = _run_subprocess(cmd, text=True, capture_output=True, check=False)
    if completed.stdout:
        print(completed.stdout.strip())
    if completed.stderr:
        print(completed.stderr.strip(), file=sys.stderr)
    if completed.returncode != 0:
        return completed.returncode
    images = sorted(outdir.glob(f"page*.{'jpg' if image_format == 'jpeg' else 'png'}"))
    if not images:
        _die(f"ERROR: renderer completed without producing page images in {outdir}", 1)
    print(f"rendered {len(images)} page image(s) to {outdir}")
    for image in images[:20]:
        print(image)
    if len(images) > 20:
        print(f"... {len(images) - 20} more")
    return 0


def _scan_xlsx_formula_errors(path: Path) -> tuple[int, dict[str, list[str]]]:
    _load_xlsx_renderer()
    from xlsx_model_renderer import inspect_formula_errors  # type: ignore

    evidence = inspect_formula_errors(path)
    return int(evidence["count"]), dict(evidence["byValue"])


def _count_xlsx_formulas(path: Path) -> int:
    try:
        import openpyxl  # type: ignore
    except ImportError:
        _missing("openpyxl")
    wb = openpyxl.load_workbook(str(path), data_only=False, read_only=True)
    count = 0
    try:
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            if not hasattr(sheet, "iter_rows"):
                continue
            for row in sheet.iter_rows():
                for cell in row:
                    if isinstance(cell.value, str) and cell.value.startswith("="):
                        count += 1
    finally:
        wb.close()
    return count


def cmd_recalc_xlsx(args: argparse.Namespace) -> int:
    path = _validate_path(args.path)
    if _ext(path) not in XLSX_EXTENSIONS:
        _die("ERROR: recalc_xlsx requires an Excel OOXML package", 3)
    soffice = _find_soffice()
    if not soffice:
        _die(
            "MISSING_DEP: LibreOffice/soffice\n"
            "Use lint_xlsx for static checks or install LibreOffice for real recalculation.",
            2,
        )

    risk = scan_ooxml_risks(path)
    unsafe_features = {
        key: parts
        for key, parts in risk["features"].items()
        if parts
    }
    if unsafe_features and not args.allow_risky:
        print(json.dumps({
            "status": "blocked",
            "recalculated": False,
            "reason": "LibreOffice round-trip could damage preservation-sensitive workbook features.",
            "risk": risk,
            "hint": "Use Excel COM/precise OOXML, or pass --allow-risky only after explicit review.",
        }, ensure_ascii=False, indent=2))
        return 5

    _, audit_xlsx_formula_integrity, inspect_formula_cache = _load_xlsx_renderer()
    from xlsx_model_renderer import inspect_formula_inventory  # type: ignore

    formula_inventory_before = inspect_formula_inventory(path)
    before_hash = hashlib.sha256(path.read_bytes()).hexdigest()
    with tempfile.TemporaryDirectory(prefix=".nexa-recalc-", dir=path.parent) as tmp:
        root = Path(tmp)
        source_dir = root / "source"
        output_dir = root / "output"
        profile_dir = root / "profile"
        source_dir.mkdir()
        output_dir.mkdir()
        source = source_dir / path.name
        shutil.copy2(path, source)
        cmd = [
            *_soffice_base_cmd(soffice, profile_dir),
            "--convert-to",
            "xlsx:Calc MS Excel 2007 XML",
            "--outdir",
            str(output_dir),
            str(source),
        ]
        completed = _run_subprocess(
            cmd,
            text=True,
            capture_output=True,
            check=False,
            env=_soffice_env(),
            timeout=120,
        )
        recalculated = output_dir / path.name
        if completed.returncode != 0 or not recalculated.is_file():
            detail = completed.stderr.strip() or completed.stdout.strip() or "no output file"
            _die(f"ERROR: LibreOffice recalculation failed: {detail}", 1)

        structural = validate_ooxml_package(recalculated)
        if structural.status == "fail":
            print(json.dumps({
                "status": "fail",
                "recalculated": False,
                "reason": "LibreOffice output failed OOXML validation and was not published.",
                "validation": structural.to_dict(),
            }, ensure_ascii=False, indent=2))
            return 1
        formula_qa = audit_xlsx_formula_integrity(recalculated)
        formula_inventory_after = inspect_formula_inventory(recalculated)
        if formula_inventory_before["fingerprint"] != formula_inventory_after["fingerprint"]:
            print(json.dumps({
                "status": "fail",
                "recalculated": False,
                "reason": "LibreOffice changed formula definitions; recalculated workbook was not published.",
                "formulaFingerprintBefore": formula_inventory_before["fingerprint"],
                "formulaFingerprintAfter": formula_inventory_after["fingerprint"],
                "formulaCountBefore": formula_inventory_before["formulaCells"],
                "formulaCountAfter": formula_inventory_after["formulaCells"],
            }, ensure_ascii=False, indent=2))
            return 1
        cached_error_total, cached_errors = _scan_xlsx_formula_errors(recalculated)
        if formula_qa["status"] == "fail" or cached_error_total:
            print(json.dumps({
                "status": "fail",
                "recalculated": False,
                "reason": "Recalculated workbook failed formula QA and was not published.",
                "formula_qa": formula_qa,
                "cached_formula_errors": cached_errors,
            }, ensure_ascii=False, indent=2))
            return 1
        staged = staging_path(path)
        shutil.copy2(recalculated, staged)
        snapshot, validation = _publish_edit(staged, path)

    after_hash = hashlib.sha256(path.read_bytes()).hexdigest()
    result = {
        "status": formula_qa["status"],
        "backend": "libreoffice-compatible",
        "recalculated": True,
        "rewritten": True,
        "contentChanged": before_hash != after_hash,
        "formula_qa": formula_qa,
        "calculation": {
            **inspect_formula_cache(path),
            "engine": "libreoffice",
            "profile": "compatible",
            "excelNative": False,
        },
        "formulaFingerprintBefore": formula_inventory_before["fingerprint"],
        "formulaFingerprintAfter": formula_inventory_after["fingerprint"],
        "cached_formula_errors": cached_errors,
        "total_formulas": _count_xlsx_formulas(path),
        "risk": risk,
        "snapshot": str(snapshot) if snapshot else None,
        "validation": validation,
    }
    print(json.dumps(result, ensure_ascii=False, indent=2))
    return 0


def cmd_lint_xlsx(args: argparse.Namespace) -> int:
    path = _validate_path(args.path)
    if _ext(path) not in XLSX_EXTENSIONS:
        _die("ERROR: lint_xlsx requires an Excel OOXML package", 3)
    _, audit_xlsx_formula_integrity, inspect_formula_cache = _load_xlsx_renderer()
    formula_qa = audit_xlsx_formula_integrity(path)
    total_errors, cached_errors = _scan_xlsx_formula_errors(path)
    contract_result = _validate_xlsx_contract(path, args.contract) if args.contract else None
    result = {
        "status": "fail" if (
            formula_qa["status"] == "fail"
            or total_errors
            or (contract_result is not None and contract_result["status"] == "fail")
        ) else formula_qa["status"],
        "formula_qa": formula_qa,
        "calculation": inspect_formula_cache(path),
        "cached_formula_errors": cached_errors,
        "preservationRisk": scan_ooxml_risks(path),
        "contract": contract_result,
        "note": "No LibreOffice or Excel process was used.",
    }
    print(json.dumps(result, ensure_ascii=False, indent=2))
    return 0 if result["status"] != "fail" else 1


VALIDATION_CONTRACT_KEYS: dict[str, set[str]] = {
    "docx": {
        "contractVersion", "required_text", "forbidden_text", "min_paragraphs", "min_tables",
        "required_styles", "no_heading_level_skips", "require_alt_text",
        "require_table_header_rows", "require_fixed_table_layout", "required_language",
        "min_comments", "require_tracked_changes", "require_no_tracked_changes",
    },
    "pptx": {
        "contractVersion", "required_text", "forbidden_text", "min_slides", "max_slides",
        "required_slide_titles", "require_speaker_notes",
    },
    "xlsx": {
        "contractVersion", "required_sheets", "required_named_ranges", "no_numeric_hardcodes_in",
        "min_rows", "sentinels", "require_formula_cache", "tie_outs", "reconciliations",
        "formula_patterns", "required_provenance",
    },
}


def _validate_contract_keys(contract: dict[str, Any], artifact_format: str) -> None:
    if not isinstance(contract, dict):
        _die("CONTRACT_SCHEMA_FAILED: validation contract root must be an object", 3)
    unknown = sorted(set(contract) - VALIDATION_CONTRACT_KEYS[artifact_format])
    if unknown:
        _die(
            f"CONTRACT_SCHEMA_FAILED: unknown {artifact_format} validation field(s): {', '.join(unknown)}",
            3,
        )
    if "contractVersion" in contract and (
        type(contract["contractVersion"]) is not int or contract["contractVersion"] != 2
    ):
        _die("CONTRACT_SCHEMA_FAILED: contractVersion must be 2", 3)
    string_arrays = {
        "required_text", "forbidden_text", "required_sheets", "required_named_ranges",
        "no_numeric_hardcodes_in", "required_styles", "required_slide_titles",
    }
    object_arrays = {"tie_outs", "reconciliations", "formula_patterns"}
    integers = {"min_paragraphs", "min_tables", "min_comments", "min_slides", "max_slides"}
    booleans = {
        "require_formula_cache", "no_heading_level_skips", "require_alt_text",
        "require_table_header_rows", "require_fixed_table_layout", "require_tracked_changes",
        "require_no_tracked_changes", "require_speaker_notes",
    }
    for key in string_arrays & set(contract):
        if not isinstance(contract[key], list) or not all(
            isinstance(item, str) for item in contract[key]
        ):
            _die(f"CONTRACT_SCHEMA_FAILED: {key} must be an array of strings", 3)
    for key in object_arrays & set(contract):
        if not isinstance(contract[key], list) or not all(
            isinstance(item, dict) for item in contract[key]
        ):
            _die(f"CONTRACT_SCHEMA_FAILED: {key} must be an array of objects", 3)
    for key in integers & set(contract):
        if type(contract[key]) is not int or contract[key] < 0:
            _die(f"CONTRACT_SCHEMA_FAILED: {key} must be a non-negative integer", 3)
    for key in booleans & set(contract):
        if type(contract[key]) is not bool:
            _die(f"CONTRACT_SCHEMA_FAILED: {key} must be a boolean", 3)
    if "min_rows" in contract and (
        not isinstance(contract["min_rows"], dict)
        or not all(
            isinstance(name, str) and type(count) is int and count >= 0
            for name, count in contract["min_rows"].items()
        )
    ):
        _die("CONTRACT_SCHEMA_FAILED: min_rows must map sheets to non-negative integers", 3)
    for key in {"sentinels", "required_provenance"} & set(contract):
        if not isinstance(contract[key], dict):
            _die(f"CONTRACT_SCHEMA_FAILED: {key} must be an object", 3)


def _contract_evidence(path: Path, contract_path: str) -> dict[str, str]:
    contract = _validate_path(contract_path)
    return {
        "artifactSha256": hashlib.sha256(path.read_bytes()).hexdigest(),
        "contractSha256": hashlib.sha256(contract.read_bytes()).hexdigest(),
    }


def _xlsx_reference_value(workbook: Any, reference: str) -> Any:
    sheet_name, separator, coordinate = str(reference).rpartition("!")
    if not separator or sheet_name not in workbook.sheetnames:
        raise ValueError(f"invalid workbook reference: {reference}")
    sheet = workbook[sheet_name]
    if not hasattr(sheet, "iter_rows"):
        raise ValueError(f"workbook reference targets a chart sheet, not a worksheet: {reference}")
    return sheet[coordinate].value


def _xlsx_relative_formula_pattern(formula: str, coordinate: str) -> str:
    try:
        from openpyxl.utils.cell import coordinate_to_tuple, column_index_from_string  # type: ignore
    except ImportError:
        _missing("openpyxl")
    origin_row, origin_col = coordinate_to_tuple(coordinate)
    reference_re = re.compile(r"(?<![A-Za-z0-9_])(?P<col>\$?[A-Z]{1,3})(?P<row>\$?\d+)")

    def replace(match: re.Match[str]) -> str:
        raw_col = match.group("col")
        raw_row = match.group("row")
        col = column_index_from_string(raw_col.replace("$", ""))
        row = int(raw_row.replace("$", ""))
        row_token = str(row) if raw_row.startswith("$") else f"[{row - origin_row}]"
        col_token = str(col) if raw_col.startswith("$") else f"[{col - origin_col}]"
        return f"R{row_token}C{col_token}"

    return reference_re.sub(replace, formula.upper())


def _xlsx_custom_properties(path: Path) -> dict[str, str]:
    with zipfile.ZipFile(path) as archive:
        if "docProps/custom.xml" not in archive.namelist():
            return {}
        root = ET.fromstring(archive.read("docProps/custom.xml"))
    properties: dict[str, str] = {}
    for prop in root:
        name = prop.attrib.get("name", "")
        value = "".join(child.text or "" for child in prop)
        if name:
            properties[name] = value
    return properties


def _validate_xlsx_contract(path: Path, contract_path: str) -> dict[str, Any]:
    try:
        import openpyxl  # type: ignore
    except ImportError:
        _missing("openpyxl")

    contract = _read_json(contract_path)
    _validate_contract_keys(contract, "xlsx")
    _, _, inspect_formula_cache = _load_xlsx_renderer()
    wb = openpyxl.load_workbook(str(path), data_only=False, read_only=False)
    values_wb = openpyxl.load_workbook(str(path), data_only=True, read_only=True)
    errors: list[dict[str, Any]] = []
    checks: dict[str, Any] = {}
    try:
        required_sheets = [str(name) for name in contract.get("required_sheets", [])]
        missing_sheets = [name for name in required_sheets if name not in wb.sheetnames]
        checks["requiredSheets"] = {"required": required_sheets, "missing": missing_sheets}
        for name in missing_sheets:
            errors.append({"code": "sheet.missing", "sheet": name})

        required_names = [str(name) for name in contract.get("required_named_ranges", [])]
        existing_names = {str(item.name) for item in wb.defined_names.values()}
        missing_names = [name for name in required_names if name not in existing_names]
        checks["requiredNamedRanges"] = {"required": required_names, "missing": missing_names}
        for name in missing_names:
            errors.append({"code": "named_range.missing", "name": name})

        hardcode_findings: dict[str, list[str]] = {}
        for sheet_name in contract.get("no_numeric_hardcodes_in", []):
            if sheet_name not in wb.sheetnames:
                continue
            sheet = wb[sheet_name]
            if not hasattr(sheet, "iter_rows"):
                errors.append({"code": "sheet.not_worksheet", "sheet": str(sheet_name)})
                continue
            locations: list[str] = []
            for row in sheet.iter_rows():
                for cell in row:
                    if isinstance(cell.value, (int, float)) and not isinstance(cell.value, bool):
                        locations.append(cell.coordinate)
            if locations:
                hardcode_findings[str(sheet_name)] = locations[:200]
                errors.append({
                    "code": "formula.numeric_hardcode",
                    "sheet": str(sheet_name),
                    "locations": locations[:200],
                    "truncated": len(locations) > 200,
                })
        checks["numericHardcodes"] = hardcode_findings

        row_checks: dict[str, Any] = {}
        for sheet_name, minimum in contract.get("min_rows", {}).items():
            if sheet_name not in wb.sheetnames:
                continue
            sheet = wb[sheet_name]
            if not hasattr(sheet, "max_row"):
                errors.append({"code": "sheet.not_worksheet", "sheet": str(sheet_name)})
                continue
            actual = sheet.max_row
            row_checks[str(sheet_name)] = {"minimum": int(minimum), "actual": actual}
            if actual < int(minimum):
                errors.append({
                    "code": "rows.minimum",
                    "sheet": str(sheet_name),
                    "minimum": int(minimum),
                    "actual": actual,
                })
        checks["minimumRows"] = row_checks

        sentinel_checks: dict[str, Any] = {}
        for reference, expected in contract.get("sentinels", {}).items():
            sheet_name, separator, coordinate = str(reference).rpartition("!")
            if not separator or sheet_name not in values_wb.sheetnames:
                errors.append({"code": "sentinel.invalid_reference", "reference": reference})
                continue
            try:
                actual = _xlsx_reference_value(values_wb, str(reference))
            except ValueError:
                errors.append({"code": "sentinel.invalid_reference", "reference": reference})
                continue
            matches = actual == expected
            if isinstance(actual, (int, float)) and isinstance(expected, (int, float)):
                matches = abs(float(actual) - float(expected)) <= 1e-9
            sentinel_checks[str(reference)] = {
                "expected": expected,
                "actual": actual,
                "matches": matches,
            }
            if not matches:
                errors.append({
                    "code": "sentinel.mismatch",
                    "reference": reference,
                    "expected": expected,
                    "actual": actual,
                })
        checks["sentinels"] = sentinel_checks
        tie_out_checks: list[dict[str, Any]] = []
        for item in contract.get("tie_outs", []):
            if not isinstance(item, dict):
                errors.append({"code": "tie_out.invalid"})
                continue
            left_ref = str(item.get("left", ""))
            right_ref = str(item.get("right", ""))
            tolerance = float(item.get("tolerance", 0) or 0)
            try:
                left = _xlsx_reference_value(values_wb, left_ref)
                right = _xlsx_reference_value(values_wb, right_ref)
                difference = abs(float(left) - float(right))
                matches = difference <= tolerance
            except (TypeError, ValueError) as error:
                left = right = None
                difference = None
                matches = False
                errors.append({"code": "tie_out.invalid_reference", "detail": str(error)})
            check = {
                "left": left_ref, "right": right_ref, "leftValue": left,
                "rightValue": right, "difference": difference, "tolerance": tolerance,
                "matches": matches,
            }
            tie_out_checks.append(check)
            if not matches and difference is not None:
                errors.append({"code": "tie_out.mismatch", **check})
        checks["tieOuts"] = tie_out_checks

        reconciliation_checks: list[dict[str, Any]] = []
        for item in contract.get("reconciliations", []):
            if not isinstance(item, dict):
                errors.append({"code": "reconciliation.invalid"})
                continue
            range_ref = str(item.get("sumRange", ""))
            equals_ref = str(item.get("equals", ""))
            tolerance = float(item.get("tolerance", 0) or 0)
            try:
                sheet_name, separator, coordinates = range_ref.rpartition("!")
                if not separator or sheet_name not in values_wb.sheetnames:
                    raise ValueError(f"invalid range: {range_ref}")
                sheet = values_wb[sheet_name]
                if not hasattr(sheet, "iter_rows"):
                    raise ValueError(f"range targets chart sheet: {range_ref}")
                values = [
                    cell.value for row in sheet[coordinates] for cell in row
                    if isinstance(cell.value, (int, float)) and not isinstance(cell.value, bool)
                ]
                total = float(sum(values))
                expected = float(_xlsx_reference_value(values_wb, equals_ref))
                difference = abs(total - expected)
                matches = difference <= tolerance
            except (TypeError, ValueError) as error:
                total = expected = difference = None
                matches = False
                errors.append({"code": "reconciliation.invalid_reference", "detail": str(error)})
            check = {
                "sumRange": range_ref, "equals": equals_ref, "sum": total,
                "expected": expected, "difference": difference, "tolerance": tolerance,
                "matches": matches,
            }
            reconciliation_checks.append(check)
            if not matches and difference is not None:
                errors.append({"code": "reconciliation.mismatch", **check})
        checks["reconciliations"] = reconciliation_checks

        formula_pattern_checks: list[dict[str, Any]] = []
        for item in contract.get("formula_patterns", []):
            if not isinstance(item, dict):
                errors.append({"code": "formula_pattern.invalid"})
                continue
            sheet_name = str(item.get("sheet", ""))
            coordinate_range = str(item.get("range", ""))
            pattern = str(item.get("pattern", ".*"))
            minimum = int(item.get("minMatches", 1))
            formulas: list[tuple[str, str]] = []
            if sheet_name in wb.sheetnames and coordinate_range:
                sheet = wb[sheet_name]
                if not hasattr(sheet, "iter_rows"):
                    errors.append({"code": "sheet.not_worksheet", "sheet": sheet_name})
                    sheet = None
                rows = sheet[coordinate_range] if sheet is not None else []
                for row in rows:
                    for cell in row:
                        if isinstance(cell.value, str) and cell.value.startswith("="):
                            formulas.append((cell.coordinate, cell.value))
            matches = [(coordinate, formula) for coordinate, formula in formulas if re.search(pattern, formula)]
            relative_patterns = {
                _xlsx_relative_formula_pattern(formula, coordinate)
                for coordinate, formula in formulas
            }
            consistent = len(relative_patterns) <= 1
            check = {
                "sheet": sheet_name, "range": coordinate_range, "pattern": pattern,
                "formulaCells": len(formulas), "matches": len(matches),
                "minimum": minimum, "relativePatternCount": len(relative_patterns),
                "consistentRelativePattern": consistent,
            }
            formula_pattern_checks.append(check)
            if len(matches) < minimum:
                errors.append({"code": "formula_pattern.minimum", **check})
            if item.get("requireConsistentRelativePattern") and not consistent:
                errors.append({"code": "formula_pattern.inconsistent", **check})
        checks["formulaPatterns"] = formula_pattern_checks

        provenance = _xlsx_custom_properties(path)
        required_provenance = contract.get("required_provenance", {})
        provenance_mismatches = {
            str(name): {"expected": expected, "actual": provenance.get(str(name))}
            for name, expected in required_provenance.items()
            if provenance.get(str(name)) != str(expected)
        } if isinstance(required_provenance, dict) else {"<contract>": {"expected": "object", "actual": type(required_provenance).__name__}}
        checks["provenance"] = {
            "properties": provenance,
            "required": required_provenance,
            "mismatches": provenance_mismatches,
        }
        for name, detail in provenance_mismatches.items():
            errors.append({"code": "provenance.mismatch", "name": name, **detail})
        calculation = inspect_formula_cache(path)
        checks["formulaCache"] = calculation
        if contract.get("require_formula_cache") and calculation["coverage"] < 1.0:
            errors.append({
                "code": "formula.cache_missing",
                "formulaCells": calculation["formulaCells"],
                "cachedFormulaCells": calculation["cachedFormulaCells"],
            })
    finally:
        wb.close()
        values_wb.close()

    return {
        "status": "fail" if errors else "pass",
        "errors": errors,
        "checks": checks,
        "evidence": _contract_evidence(path, contract_path),
    }


def _contract_text_checks(text: str, contract: dict[str, Any]) -> tuple[list[dict[str, Any]], dict[str, Any]]:
    errors: list[dict[str, Any]] = []
    required = [str(item) for item in contract.get("required_text", [])]
    missing = [item for item in required if item not in text]
    forbidden = [str(item) for item in contract.get("forbidden_text", [])]
    present = [item for item in forbidden if item in text]
    for item in missing:
        errors.append({"code": "required_text.missing", "text": item})
    for item in present:
        errors.append({"code": "forbidden_text.present", "text": item})
    return errors, {
        "requiredText": {"required": required, "missing": missing},
        "forbiddenText": {"forbidden": forbidden, "present": present},
    }


def _validate_docx_contract(path: Path, contract_path: str) -> dict[str, Any]:
    contract = _read_json(contract_path)
    _validate_contract_keys(contract, "docx")
    skills_root = Path(__file__).resolve().parents[2]
    review_dir = skills_root / "docx-document-design" / "scripts"
    if str(review_dir) not in sys.path:
        sys.path.insert(0, str(review_dir))
    from docx_review_editor import UNSUPPORTED_REVISION_ELEMENTS  # type: ignore

    revision_names = {"ins", "del", *UNSUPPORTED_REVISION_ELEMENTS}
    word_ns = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
    with zipfile.ZipFile(path) as archive:
        names = set(archive.namelist())
        document_root = ET.fromstring(archive.read("word/document.xml"))
        styles: dict[str, str] = {}
        if "word/styles.xml" in names:
            styles_root = ET.fromstring(archive.read("word/styles.xml"))
            for style in styles_root.iter(f"{{{word_ns}}}style"):
                style_id = style.attrib.get(f"{{{word_ns}}}styleId", "")
                style_name = next(
                    (
                        child.attrib.get(f"{{{word_ns}}}val", "")
                        for child in style
                        if child.tag == f"{{{word_ns}}}name"
                    ),
                    style_id,
                )
                if style_id:
                    styles[style_id] = style_name
        story_parts = [
            name
            for name in archive.namelist()
            if name == "word/document.xml"
            or re.fullmatch(
                r"word/(?:header[0-9]+|footer[0-9]+|comments(?:Extended)?|footnotes|endnotes)\.xml",
                name,
                re.IGNORECASE,
            )
        ]
        stories = [
            "".join(
                element.text or ""
                for element in ET.fromstring(archive.read(name)).iter(f"{{{word_ns}}}t")
            )
            for name in story_parts
        ]
    paragraphs = list(document_root.iter(f"{{{word_ns}}}p"))
    errors, checks = _contract_text_checks("\n".join(stories), contract)
    measurements = {
        "paragraphs": len(paragraphs),
        "tables": sum(1 for _ in document_root.iter(f"{{{word_ns}}}tbl")),
    }
    checks["measurements"] = measurements
    for key, actual in measurements.items():
        minimum = contract.get(f"min_{key}")
        if minimum is not None and actual < int(minimum):
            errors.append({
                "code": f"{key}.minimum",
                "minimum": int(minimum),
                "actual": actual,
            })
    required_styles = [str(item) for item in contract.get("required_styles", [])]
    paragraph_style_ids = [
        next(
            (
                item.attrib.get(f"{{{word_ns}}}val", "")
                for item in paragraph.iter(f"{{{word_ns}}}pStyle")
            ),
            "",
        )
        for paragraph in paragraphs
    ]
    def display_style(style_id: str) -> str:
        if match := re.fullmatch(r"Heading([1-9])", style_id, re.IGNORECASE):
            return f"Heading {match.group(1)}"
        return styles.get(style_id, style_id)

    used_styles = {display_style(style_id) for style_id in paragraph_style_ids if style_id}
    used_style_keys = {style.casefold() for style in used_styles}
    missing_styles = [style for style in required_styles if style.casefold() not in used_style_keys]
    checks["requiredStyles"] = {"required": required_styles, "missing": missing_styles}
    for style in missing_styles:
        errors.append({"code": "style.missing", "style": style})

    heading_levels: list[int] = []
    for style_id in paragraph_style_ids:
        style_name = display_style(style_id)
        match = re.fullmatch(r"Heading ?([1-6])", style_name, re.IGNORECASE)
        if match:
            heading_levels.append(int(match.group(1)))
    heading_jumps = [
        {"from": previous, "to": current, "position": index}
        for index, (previous, current) in enumerate(zip(heading_levels, heading_levels[1:]), start=2)
        if current > previous + 1
    ]
    if heading_levels and heading_levels[0] != 1:
        heading_jumps.insert(0, {"from": 0, "to": heading_levels[0], "position": 1})
    checks["headingOrder"] = {"levels": heading_levels, "jumps": heading_jumps}
    if contract.get("no_heading_level_skips") and heading_jumps:
        errors.append({"code": "heading.level_skip", "jumps": heading_jumps})

    package_checks = {
        "pictures": 0,
        "picturesWithAltText": 0,
        "tables": 0,
        "tablesWithHeaderRows": 0,
        "tablesWithFixedLayout": 0,
        "languageValues": [],
        "comments": 0,
        "trackedInsertions": 0,
        "trackedDeletions": 0,
        "trackedChanges": {},
    }
    with zipfile.ZipFile(path) as archive:
        for name in archive.namelist():
            if not name.startswith("word/") or not name.endswith(".xml"):
                continue
            root = ET.fromstring(archive.read(name))
            for element in root.iter():
                local = element.tag.rsplit("}", 1)[-1]
                if local == "docPr":
                    package_checks["pictures"] += 1
                    if str(element.attrib.get("descr", "")).strip():
                        package_checks["picturesWithAltText"] += 1
                elif local == "tbl":
                    package_checks["tables"] += 1
                    first_row = next((child for child in element if child.tag.rsplit("}", 1)[-1] == "tr"), None)
                    if first_row is not None and any(
                        descendant.tag.rsplit("}", 1)[-1] == "tblHeader"
                        for descendant in first_row.iter()
                    ):
                        package_checks["tablesWithHeaderRows"] += 1
                    if any(
                        descendant.tag.rsplit("}", 1)[-1] == "tblLayout"
                        and descendant.attrib.get("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}type") == "fixed"
                        for descendant in element.iter()
                    ):
                        package_checks["tablesWithFixedLayout"] += 1
                elif local == "lang":
                    for attribute, value in element.attrib.items():
                        if attribute.rsplit("}", 1)[-1] in {"val", "eastAsia", "bidi"} and value:
                            package_checks["languageValues"].append(value)
                elif local == "comment" and name == "word/comments.xml":
                    package_checks["comments"] += 1
                elif local == "ins":
                    package_checks["trackedInsertions"] += 1
                elif local == "del":
                    package_checks["trackedDeletions"] += 1
                if local in revision_names:
                    tracked_changes = package_checks["trackedChanges"]
                    tracked_changes[local] = tracked_changes.get(local, 0) + 1
    package_checks["languageValues"] = sorted(set(package_checks["languageValues"]))
    checks["accessibilityAndLayout"] = package_checks
    if (
        contract.get("require_alt_text")
        and package_checks["picturesWithAltText"] != package_checks["pictures"]
    ):
        errors.append({
            "code": "image.alt_text_missing",
            "pictures": package_checks["pictures"],
            "withAltText": package_checks["picturesWithAltText"],
        })
    if (
        contract.get("require_table_header_rows")
        and package_checks["tablesWithHeaderRows"] != package_checks["tables"]
    ):
        errors.append({
            "code": "table.header_row_missing",
            "tables": package_checks["tables"],
            "withHeaderRows": package_checks["tablesWithHeaderRows"],
        })
    if (
        contract.get("require_fixed_table_layout")
        and package_checks["tablesWithFixedLayout"] != package_checks["tables"]
    ):
        errors.append({
            "code": "table.fixed_layout_missing",
            "tables": package_checks["tables"],
            "withFixedLayout": package_checks["tablesWithFixedLayout"],
        })
    required_language = contract.get("required_language")
    if required_language and str(required_language) not in package_checks["languageValues"]:
        errors.append({
            "code": "document.language_missing",
            "required": str(required_language),
            "actual": package_checks["languageValues"],
        })
    minimum_comments = int(contract.get("min_comments", 0) or 0)
    if package_checks["comments"] < minimum_comments:
        errors.append({
            "code": "comments.minimum",
            "minimum": minimum_comments,
            "actual": package_checks["comments"],
        })
    revision_count = sum(package_checks["trackedChanges"].values())
    if contract.get("require_tracked_changes") and revision_count == 0:
        errors.append({"code": "tracked_changes.missing"})
    if contract.get("require_no_tracked_changes") and revision_count:
        errors.append({
            "code": "tracked_changes.present",
            "counts": package_checks["trackedChanges"],
        })
    return {
        "status": "fail" if errors else "pass",
        "errors": errors,
        "checks": checks,
        "evidence": _contract_evidence(path, contract_path),
    }


def _validate_pptx_contract(path: Path, contract_path: str) -> dict[str, Any]:
    contract = _read_json(contract_path)
    _validate_contract_keys(contract, "pptx")
    skills_root = Path(__file__).resolve().parents[2]
    audit_dir = skills_root / "pptx-presentation-design" / "scripts"
    if str(audit_dir) not in sys.path:
        sys.path.insert(0, str(audit_dir))
    try:
        import pptx_audit  # type: ignore
    except ImportError:
        _missing("pptx_audit")
    package = pptx_audit.audit(path)
    fragments: list[str] = []
    titles: list[str] = []
    slides_with_notes = 0
    with zipfile.ZipFile(path) as archive:
        for slide in package.get("slide_details", []):
            slide_part = str(slide.get("part", ""))
            slide_root = pptx_audit.parse_xml(pptx_audit.read_text(archive, slide_part))
            slide_text = pptx_audit.slide_text(slide_root)
            if slide_text:
                fragments.append(slide_text)
            text_shapes = [
                shape
                for shape in slide.get("shape_details", [])
                if str(shape.get("text", "")).strip()
            ]
            title_shapes = [shape for shape in text_shapes if shape.get("isTitle")]
            if title_shapes:
                titles.append(str(title_shapes[0]["text"]).strip())
            elif text_shapes:
                # Imported decks often use an ordinary text box as the visual
                # title. The audit inventory is position-sorted, so the first
                # non-empty text shape is the same deterministic fallback used
                # for normal PPTX and macro/template packages.
                titles.append(str(text_shapes[0]["text"]).strip())
            notes_text: list[str] = []
            for relationship in pptx_audit.rel_targets(archive, slide_part):
                if (
                    relationship.get("target_mode") != "External"
                    and relationship.get("type", "").rsplit("/", 1)[-1] == "notesSlide"
                ):
                    notes = pptx_audit.parse_xml(
                        pptx_audit.read_text(archive, str(relationship.get("target", "")))
                    )
                    if notes is not None:
                        notes_text.extend(
                            (item.text or "").strip()
                            for item in notes.findall(".//a:t", pptx_audit.NS)
                            if (item.text or "").strip()
                        )
            if notes_text:
                slides_with_notes += 1
                fragments.extend(notes_text)
    errors, checks = _contract_text_checks("\n".join(fragments), contract)
    slide_count = int(package.get("slides", 0))
    checks["measurements"] = {
        "slides": slide_count,
        "slidesWithNotes": slides_with_notes,
        "titles": titles,
    }
    if contract.get("min_slides") is not None and slide_count < int(contract["min_slides"]):
        errors.append({
            "code": "slides.minimum",
            "minimum": int(contract["min_slides"]),
            "actual": slide_count,
        })
    if contract.get("max_slides") is not None and slide_count > int(contract["max_slides"]):
        errors.append({
            "code": "slides.maximum",
            "maximum": int(contract["max_slides"]),
            "actual": slide_count,
        })
    required_titles = [str(item) for item in contract.get("required_slide_titles", [])]
    missing_titles = [item for item in required_titles if item not in titles]
    checks["requiredSlideTitles"] = {"required": required_titles, "missing": missing_titles}
    for title in missing_titles:
        errors.append({"code": "slide_title.missing", "title": title})
    if contract.get("require_speaker_notes") and slides_with_notes != slide_count:
        errors.append({
            "code": "speaker_notes.missing",
            "slides": slide_count,
            "slidesWithNotes": slides_with_notes,
        })
    return {
        "status": "fail" if errors else "pass",
        "errors": errors,
        "checks": checks,
        "evidence": _contract_evidence(path, contract_path),
    }


def _validate_artifact_contract(path: Path, contract_path: str) -> dict[str, Any]:
    if _ext(path) in XLSX_EXTENSIONS:
        return _validate_xlsx_contract(path, contract_path)
    if _ext(path) in DOCX_EXTENSIONS:
        return _validate_docx_contract(path, contract_path)
    if _ext(path) in PPTX_EXTENSIONS:
        return _validate_pptx_contract(path, contract_path)
    _die("ERROR: validation contracts require DOCX, PPTX, or XLSX", 3)


def cmd_validate(args: argparse.Namespace) -> int:
    path = _validate_path(args.path)
    ext = _ext(path)
    result: dict[str, Any] = {"format": ext, "path": str(path), "status": "pass"}
    if ext in OOXML_EXTENSIONS:
        structural = validate_ooxml_package(path)
        result["structural"] = structural.to_dict()
        if structural.status == "fail":
            result["status"] = "fail"
            print(json.dumps(result, ensure_ascii=False, indent=2))
            return 1
    if ext in DOCX_EXTENSIONS:
        if ext == "docx":
            try:
                import docx  # type: ignore
            except ImportError:
                _missing("python-docx")
            doc = docx.Document(str(path))
            result["backend"] = {
                "id": "python-docx",
                "paragraphs": len(doc.paragraphs),
                "tables": len(doc.tables),
            }
        else:
            with zipfile.ZipFile(path) as archive:
                root = ET.fromstring(archive.read("word/document.xml"))
            result["backend"] = {
                "id": "wordprocessingml-package",
                "paragraphs": sum(1 for item in root.iter() if item.tag.rsplit("}", 1)[-1] == "p"),
                "tables": sum(1 for item in root.iter() if item.tag.rsplit("}", 1)[-1] == "tbl"),
                "macroTemplatePreservation": True,
            }
    elif ext in PPTX_EXTENSIONS:
        if ext == "pptx":
            try:
                from pptx import Presentation  # type: ignore
            except ImportError:
                _missing("python-pptx")
            prs = Presentation(str(path))
            result["backend"] = {"id": "python-pptx", "slides": len(prs.slides)}
        else:
            result["backend"] = {"id": "presentationml-package", "macroTemplatePreservation": True}
        skills_root = Path(__file__).resolve().parents[2]
        audit_dir = skills_root / "pptx-presentation-design" / "scripts"
        if str(audit_dir) not in sys.path:
            sys.path.insert(0, str(audit_dir))
        try:
            import pptx_audit  # type: ignore

            package_graph = pptx_audit.audit(path)
            result["packageGraph"] = package_graph
            if package_graph.get("validation_errors"):
                result["status"] = "fail"
        except Exception as exc:  # noqa: BLE001
            result["packageGraphWarning"] = f"{type(exc).__name__}: {exc}"
            if result["status"] == "pass":
                result["status"] = "warn"
    elif ext in XLSX_EXTENSIONS:
        try:
            import openpyxl  # type: ignore
        except ImportError:
            _missing("openpyxl")
        wb = openpyxl.load_workbook(
            str(path),
            data_only=False,
            read_only=True,
            keep_vba=ext in {"xlsm", "xltm"},
        )
        sheet_count = len(wb.worksheets)
        sheet_names = ",".join(wb.sheetnames)
        wb.close()
        total_errors, errors = _scan_xlsx_formula_errors(path)
        formula_count = _count_xlsx_formulas(path)
        result["backend"] = {
            "id": "openpyxl",
            "sheets": sheet_count,
            "sheetNames": sheet_names.split(",") if sheet_names else [],
            "formulas": formula_count,
            "formulaErrors": total_errors,
        }
        if errors:
            result["status"] = "fail"
            result["cachedFormulaErrors"] = errors
            print(json.dumps(result, ensure_ascii=False, indent=2))
            return 1
        try:
            _, audit_xlsx_formula_integrity, inspect_formula_cache = _load_xlsx_renderer()
            formula_qa = audit_xlsx_formula_integrity(path)
            result["formulaQa"] = formula_qa
            result["calculation"] = inspect_formula_cache(path)
            if formula_qa["status"] == "fail":
                result["status"] = "fail"
                print(json.dumps(result, ensure_ascii=False, indent=2))
                return 1
        except SystemExit:
            raise
        except Exception as exc:  # noqa: BLE001
            result["formulaQaWarning"] = f"{type(exc).__name__}: {exc}"
            if result["status"] == "pass":
                result["status"] = "warn"
    elif ext == "pdf":
        try:
            from pypdf import PdfReader  # type: ignore
        except ImportError:
            _missing("pypdf")
        reader = PdfReader(str(path))
        result["backend"] = {"id": "pypdf", "pages": len(reader.pages)}
    else:
        _die(f"ERROR: validate does not support .{ext}", 3)
    if getattr(args, "contract", None):
        contract = _validate_artifact_contract(path, args.contract)
        result["contract"] = contract
        if contract["status"] == "fail":
            result["status"] = "fail"
    if args.json:
        print(json.dumps(result, ensure_ascii=False, indent=2))
    else:
        backend = result.get("backend", {})
        print(f"VALID {ext.upper()} backend={backend.get('id', 'structural')} status={result['status']}")
    return 0 if result["status"] != "fail" else 1


def cmd_convert(args: argparse.Namespace) -> int:
    path = _validate_path(args.path)
    outdir = Path(args.outdir).resolve() if args.outdir else path.parent
    try:
        outdir.relative_to(Path.cwd().resolve())
    except ValueError:
        _die(f"ERROR: --outdir escapes workspace: {outdir}", 3)
    outdir.mkdir(parents=True, exist_ok=True)
    completed = _run_soffice_convert(path, args.to, outdir)
    if completed.stdout:
        print(completed.stdout.strip())
    if completed.stderr:
        print(completed.stderr.strip(), file=sys.stderr)
    if completed.returncode != 0:
        return completed.returncode
    print(f"converted {path.name} -> .{args.to} in {outdir}")
    return 0


# ---------------------------------------------------------------------------
# argparse
# ---------------------------------------------------------------------------

def build_parser() -> argparse.ArgumentParser:
    p = argparse.ArgumentParser(
        prog="edit_doc.py",
        description="Edit existing DOCX/PPTX/PDF/XLSX documents from run_shell.",
    )
    p.add_argument("--path", help="Absolute path to the target file")
    sub = p.add_subparsers(dest="cmd", required=True)

    p_check = sub.add_parser("check", help="Report available backends")
    p_check.add_argument("--json", action="store_true", help="Emit machine-readable readiness JSON")
    p_check.add_argument(
        "--format",
        choices=["docx", "xlsx", "pptx", "pdf", "all"],
        default=None,
        help="Make only this format's package(s) required; without it missing format packages are advisory.",
    )
    p_check.set_defaults(func=cmd_check)

    p_rep = sub.add_parser("replace", help="Replace text (docx/pptx/xlsx)")
    p_rep.add_argument("--find", required=True)
    p_rep.add_argument("--replace", default="")
    p_rep.add_argument("--dry-run", action="store_true")
    p_rep.add_argument("--expected-sha256", default=None, help="Fail if the artifact changed since inspection")
    p_rep.add_argument("--expected-count", type=int, default=None, help="Fail unless exactly this many matches exist")
    p_rep.add_argument("--scope", default=None, help="DOCX story scopes: body,table,textbox,header,footer,comments,footnotes,endnotes")
    p_rep.add_argument("--occurrence", type=int, default=None, help="Replace only this 1-based match within the selected DOCX scopes")
    p_rep.add_argument("--allow-style-merge", action="store_true", help="Allow a DOCX replacement to cross incompatible style/hyperlink runs")
    p_rep.set_defaults(func=cmd_replace)

    p_red = sub.add_parser("redact", help="Redact text (docx/pptx/xlsx)")
    p_red.add_argument("--find", required=True)
    p_red.add_argument("--replace", default=None)
    p_red.add_argument("--dry-run", action="store_true")
    p_red.add_argument("--expected-sha256", default=None, help="Fail if the artifact changed since inspection")
    p_red.add_argument("--expected-count", type=int, default=None, help="Fail unless exactly this many matches exist")
    p_red.add_argument("--scope", default=None, help="DOCX story scopes; this command is visible-text replacement, not secure redaction")
    p_red.add_argument("--occurrence", type=int, default=None)
    p_red.add_argument("--allow-style-merge", action="store_true")
    p_red.set_defaults(func=cmd_redact)

    p_secure_red = sub.add_parser(
        "secure_redact",
        help="Redact DOCX package text and prove the original is absent; fail on uninspectable media/embeddings",
    )
    p_secure_red.add_argument("--find", required=True)
    p_secure_red.add_argument("--replace", default="[REDACTED]")
    p_secure_red.add_argument("--expected-count", type=int, default=None)
    p_secure_red.add_argument("--expected-sha256", default=None)
    p_secure_red.add_argument("--privacy-scrub", action="store_true")
    p_secure_red.set_defaults(func=cmd_secure_redact)

    p_ext = sub.add_parser("extract", help="Extract plain text (docx/pdf/pptx/xlsx)")
    p_ext.add_argument("--pages", default=None, help="e.g. 1-3 or 1,3,5")
    p_ext.add_argument("--sheets", default=None, help="Comma-separated XLSX sheet names")
    p_ext.set_defaults(func=cmd_extract)

    p_ins = sub.add_parser("insert_slide", help="Insert a slide into a pptx")
    p_ins.add_argument("--after", type=int, default=0)
    p_ins.add_argument("--title", default="")
    p_ins.add_argument("--body", default="")
    p_ins.set_defaults(func=cmd_insert_slide)

    p_ver = sub.add_parser("version", help="Snapshot file to .nexa/doc-history")
    p_ver.set_defaults(func=cmd_version)

    p_cd = sub.add_parser("create_docx", help="Create a DOCX using python-docx")
    p_cd.add_argument("--spec", default=None, help="Absolute DOCX Spec v2 JSON path")
    p_cd.add_argument("--title", default="")
    p_cd.add_argument("--subtitle", default="")
    p_cd.add_argument("--body", default="")
    p_cd.add_argument("--input-md", default=None, help="Absolute path to markdown source")
    p_cd.add_argument("--template", default=None, help="Optional absolute .docx template path")
    p_cd.add_argument("--font", default="Calibri")
    p_cd.add_argument("--footer", default="")
    p_cd.add_argument("--author", default="Nexa")
    p_cd.set_defaults(func=cmd_create_docx)

    p_cx = sub.add_parser("create_xlsx", help="Create an XLSX workbook from a JSON spec")
    p_cx.add_argument("--spec", required=True, help="Absolute path to workbook JSON spec")
    p_cx.set_defaults(func=cmd_create_xlsx)

    p_cp = sub.add_parser("create_pptx", help="Create a PPTX presentation from a JSON spec")
    p_cp.add_argument("--spec", required=True, help="Absolute path to deck JSON spec, or '-' to read JSON from stdin")
    p_cp.add_argument("--template", default=None, help="Optional absolute .pptx template path")
    p_cp.set_defaults(func=cmd_create_pptx)

    p_edit_pptx = sub.add_parser("edit_pptx", help="Apply typed PPTX display-order, shape text, clone, or transition edits")
    p_edit_pptx.add_argument("--spec", required=True, help="Absolute JSON file containing an operations array")
    p_edit_pptx.set_defaults(func=cmd_edit_pptx)

    p_review_docx = sub.add_parser("review_docx", help="Apply typed DOCX comment or tracked-change operations")
    p_review_docx.add_argument("--spec", required=True, help="Absolute JSON file containing an operations array")
    p_review_docx.set_defaults(func=cmd_review_docx)

    p_comments_docx = sub.add_parser("comments_docx", help="Extract DOCX comments as JSON")
    p_comments_docx.set_defaults(func=cmd_comments_docx)

    p_chp = sub.add_parser("create_html_pptx", help="Create a PPTX from an HTML-first deck spec")
    p_chp.add_argument("--spec", required=True, help="Absolute path to HTML deck JSON spec, or '-' to read JSON from stdin")
    p_chp.add_argument("--outdir", required=True, help="Absolute output project directory for HTML, render artifacts, manifest, and QA")
    p_chp.add_argument("--mode", choices=["hybrid", "native", "raster"], default="hybrid")
    p_chp.add_argument("--screenshot", choices=["auto", "require", "skip"], default="auto")
    p_chp.set_defaults(func=cmd_create_html_pptx)

    p_unpack = sub.add_parser("unpack", help="Unpack DOCX/PPTX/XLSX into editable OOXML")
    p_unpack.add_argument("--outdir", required=True, help="Absolute output directory")
    p_unpack.add_argument("--overwrite", action="store_true")
    p_unpack.set_defaults(func=cmd_unpack)

    p_pack = sub.add_parser("pack", help="Pack an unpacked OOXML directory back into DOCX/PPTX/XLSX")
    p_pack.add_argument("--input-dir", required=True, help="Absolute unpacked OOXML directory")
    p_pack.set_defaults(func=cmd_pack)

    p_render = sub.add_parser("render", help="Render DOCX/PPTX/XLSX/PDF pages or slides to images for visual QA")
    p_render.add_argument("--outdir", required=True, help="Absolute output directory for page images")
    p_render.add_argument("--dpi", type=int, default=150)
    p_render.add_argument("--format", default="png", choices=["png", "jpeg"])
    p_render.add_argument("--sheets", default=None, help="XLSX only: all, active, or comma-separated visible sheet names")
    p_render.set_defaults(func=cmd_render)

    p_lint = sub.add_parser("lint_xlsx", help="Lint XLSX formulas without LibreOffice or Excel automation")
    p_lint.add_argument(
        "--contract",
        default=None,
        help="Optional absolute JSON workbook contract with sheets, names, hardcodes, rows, and sentinels",
    )
    p_lint.set_defaults(func=cmd_lint_xlsx)

    p_edit_xlsx = sub.add_parser("edit_xlsx", help="Apply typed direct-OOXML XLSX value, formula, range, clear, or style edits")
    p_edit_xlsx.add_argument("--spec", required=True, help="Absolute JSON file containing an operations array")
    p_edit_xlsx.set_defaults(func=cmd_edit_xlsx)

    p_recalc = sub.add_parser("recalc_xlsx", help="Recalculate and resave XLSX with LibreOffice")
    p_recalc.add_argument(
        "--allow-risky",
        action="store_true",
        help="Allow LibreOffice round-trip after reviewing preservation-sensitive features",
    )
    p_recalc.set_defaults(func=cmd_recalc_xlsx)

    p_val = sub.add_parser("validate", help="Validate OOXML structure, relationships, content types, and backend open")
    p_val.add_argument("--json", action="store_true", help="Emit machine-readable validation JSON")
    p_val.add_argument(
        "--contract",
        default=None,
        help="Optional absolute JSON contract for format-specific content and structure checks",
    )
    p_val.set_defaults(func=cmd_validate)

    p_conv = sub.add_parser("convert", help="Convert via LibreOffice headless")
    p_conv.add_argument("--to", required=True, help="Output extension/filter, e.g. pdf, docx, xlsx")
    p_conv.add_argument("--outdir", default=None, help="Optional absolute output directory")
    p_conv.set_defaults(func=cmd_convert)

    return p


def main(argv: list[str] | None = None) -> int:
    parser = build_parser()
    args = parser.parse_args(argv)
    try:
        return args.func(args)
    except SystemExit:
        raise
    except Exception as e:  # noqa: BLE001
        print(f"ERROR: {e}", file=sys.stderr)
        return 1


if __name__ == "__main__":
    sys.exit(main())
