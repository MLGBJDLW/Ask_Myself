from __future__ import annotations

import json
import hashlib
import hmac
import re
import sys
import tempfile
import unittest
import zipfile
from pathlib import Path
from unittest import mock

import office_artifact_engine
import office_artifact_service
from office_artifact_engine import OfficeArtifactEngine, OfficeArtifactError


class OfficeArtifactEngineTests(unittest.TestCase):
    def test_windows_process_liveness_probe_never_uses_kill(self) -> None:
        kernel = mock.MagicMock()
        kernel.OpenProcess.return_value = 11

        def write_still_active(process, output):
            output._obj.value = 259
            return 1

        kernel.GetExitCodeProcess.side_effect = write_still_active
        kernel.CloseHandle.return_value = 1
        with mock.patch.object(office_artifact_engine.os, "name", "nt"), mock.patch.object(
            office_artifact_engine.ctypes, "WinDLL", return_value=kernel
        ), mock.patch.object(office_artifact_engine.os, "kill", side_effect=AssertionError("must not kill")):
            self.assertTrue(office_artifact_engine._process_is_alive(4242))
        kernel.CloseHandle.assert_called_once_with(11)

    def setUp(self) -> None:
        self.temp = tempfile.TemporaryDirectory()
        self.root = Path(self.temp.name).resolve()
        self.engine = OfficeArtifactEngine(self.root)

    def tearDown(self) -> None:
        self.temp.cleanup()

    def _docx_request(self, destination: Path) -> dict:
        return {
            "requestVersion": 2,
            "format": "docx",
            "intent": "create",
            "destination": str(destination),
            "operations": [{
                "op": "create",
                "title": "Candidate report",
                "body": "Verified body",
            }],
            "guarantees": {
                "quality": "standard",
                "preservation": "strict",
                "render": "none",
            },
            "validation": {"contractVersion": 2, "required_text": ["Verified body"]},
        }

    def test_execute_publish_restore_lifecycle_keeps_destination_gated(self) -> None:
        destination = self.root / "delivery.docx"
        outcome = self.engine.execute(self._docx_request(destination))

        self.assertEqual("candidate", outcome["status"])
        self.assertFalse(destination.exists())
        candidate = Path(outcome["candidatePath"])
        self.assertTrue(candidate.exists())

        published = self.engine.decide(outcome["candidateId"], "publish")
        self.assertEqual("published", published["status"])
        self.assertTrue(destination.exists())
        self.assertTrue((self.root / "delivery.docx.manifest.json").exists())

        restored = self.engine.restore(published["receiptId"])
        self.assertEqual("restored", restored["status"])
        self.assertFalse(destination.exists())

    def test_restore_refuses_to_overwrite_newer_destination(self) -> None:
        destination = self.root / "delivery.docx"
        candidate = self.engine.execute(self._docx_request(destination))
        published = self.engine.decide(candidate["candidateId"], "publish")
        destination.write_bytes(destination.read_bytes() + b"newer")

        with self.assertRaisesRegex(OfficeArtifactError, "changed after publication"):
            self.engine.restore(published["receiptId"])

    def test_restore_refuses_changed_sidecar_and_tampered_receipt(self) -> None:
        destination = self.root / "delivery.docx"
        candidate = self.engine.execute(self._docx_request(destination))
        published = self.engine.decide(candidate["candidateId"], "publish")
        manifest = self.root / "delivery.docx.manifest.json"
        manifest.write_text('{"user":"newer manifest"}\n', encoding="utf-8")

        with self.assertRaisesRegex(OfficeArtifactError, "sidecar changed"):
            self.engine.restore(published["receiptId"])
        self.assertEqual('{"user":"newer manifest"}\n', manifest.read_text(encoding="utf-8"))
        self.assertTrue(destination.exists())

        # A separate publication proves that changing any receipt field is
        # rejected before a path or snapshot is trusted.
        second_destination = self.root / "second.docx"
        second_candidate = self.engine.execute(self._docx_request(second_destination))
        second = self.engine.decide(second_candidate["candidateId"], "publish")
        receipt_path = (
            self.root
            / ".nexa"
            / "office-artifacts"
            / "receipts"
            / f"{second['receiptId']}.json"
        )
        receipt = json.loads(receipt_path.read_text(encoding="utf-8"))
        receipt["destination"] = str(self.root / "forged.docx")
        receipt_path.write_text(json.dumps(receipt, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
        with self.assertRaisesRegex(OfficeArtifactError, "HMAC|integrity"):
            self.engine.restore(second["receiptId"])
        self.assertTrue(second_destination.exists())

    def test_publish_uses_destination_cas_and_exclusive_lock(self) -> None:
        destination = self.root / "cas.docx"
        candidate = self.engine.execute(self._docx_request(destination))
        destination.write_bytes(b"external writer")
        with self.assertRaisesRegex(OfficeArtifactError, "existence changed"):
            self.engine.decide(candidate["candidateId"], "publish")

        destination.unlink()
        lock = self.engine._acquire_destination_lock(destination, "f" * 32)
        try:
            with self.assertRaisesRegex(OfficeArtifactError, "owns the destination lock"):
                self.engine.decide(candidate["candidateId"], "publish")
        finally:
            lock.unlink(missing_ok=True)

    def test_shared_manifest_role_cas_prevents_cross_destination_overwrite(self) -> None:
        manifest = self.root / "shared-manifest.json"
        first_request = self._docx_request(self.root / "first.docx")
        second_request = self._docx_request(self.root / "second.docx")
        first_request["delivery"] = {"manifest": str(manifest)}
        second_request["delivery"] = {"manifest": str(manifest)}
        first = self.engine.execute(first_request)
        second = self.engine.execute(second_request)

        self.engine.decide(first["candidateId"], "publish")
        first_manifest = manifest.read_bytes()
        with self.assertRaisesRegex(OfficeArtifactError, "role existence changed"):
            self.engine.decide(second["candidateId"], "publish")
        self.assertEqual(first_manifest, manifest.read_bytes())
        self.assertFalse((self.root / "second.docx").exists())

    def test_candidate_state_hmac_blocks_destination_and_manifest_retargeting(self) -> None:
        destination = self.root / "signed.docx"
        forged_destination = self.root / "forged.docx"
        candidate = self.engine.execute(self._docx_request(destination))
        state_path = Path(candidate["candidatePath"]).parent / "state.json"
        state = json.loads(state_path.read_text(encoding="utf-8"))
        state["destination"] = str(forged_destination)
        state["requestedManifest"] = str(self.root / "source-code.py")
        state_path.write_text(json.dumps(state), encoding="utf-8")

        with self.assertRaisesRegex(OfficeArtifactError, "candidate state HMAC"):
            self.engine.decide(candidate["candidateId"], "publish")
        self.assertFalse(destination.exists())
        self.assertFalse(forged_destination.exists())

    def test_restore_reinstates_existing_destination_and_manifest(self) -> None:
        import docx

        destination = self.root / "existing.docx"
        manifest = self.root / "existing-manifest.json"
        original = docx.Document()
        original.add_paragraph("Original destination")
        original.save(destination)
        original_hash = destination.read_bytes()
        manifest.write_text('{"status":"old"}\n', encoding="utf-8")
        old_manifest = manifest.read_bytes()
        request = self._docx_request(destination)
        request["delivery"] = {"manifest": str(manifest)}

        candidate = self.engine.execute(request)
        published = self.engine.decide(candidate["candidateId"], "publish")
        self.assertNotEqual(original_hash, destination.read_bytes())
        receipt = json.loads(
            (self.root / ".nexa" / "office-artifacts" / "receipts" / f"{published['receiptId']}.json")
            .read_text(encoding="utf-8")
        )
        self.assertTrue(receipt["existedBefore"])
        self.assertIsNotNone(receipt["snapshot"])

        self.engine.restore(published["receiptId"])
        self.assertEqual(original_hash, destination.read_bytes())
        self.assertEqual(old_manifest, manifest.read_bytes())

    def test_restore_journal_recovers_partial_multi_role_restore(self) -> None:
        import docx

        destination = self.root / "restore-crash.docx"
        manifest = self.root / "restore-crash-manifest.json"
        original = docx.Document()
        original.add_paragraph("Original destination")
        original.save(destination)
        original_bytes = destination.read_bytes()
        manifest.write_text('{"status":"original"}\n', encoding="utf-8")
        original_manifest = manifest.read_bytes()
        request = self._docx_request(destination)
        request["delivery"] = {"manifest": str(manifest)}
        candidate = self.engine.execute(request)
        published = self.engine.decide(candidate["candidateId"], "publish")
        published_bytes = destination.read_bytes()
        real_rollback = office_artifact_engine.rollback_published_artifact

        def fail_destination(target, snapshot, workspace_root):
            if Path(target).resolve() == destination.resolve():
                raise OSError("injected destination restore failure")
            return real_rollback(target, snapshot, workspace_root)

        with mock.patch.object(
            office_artifact_engine,
            "rollback_published_artifact",
            side_effect=fail_destination,
        ):
            with self.assertRaisesRegex(OSError, "injected destination"):
                self.engine.restore(published["receiptId"])

        self.assertEqual(original_manifest, manifest.read_bytes())
        self.assertEqual(published_bytes, destination.read_bytes())
        journals = self.root / ".nexa" / "office-artifacts" / "journals"
        self.assertEqual(1, len(list(journals.glob("restore-*.json"))))

        with mock.patch.object(office_artifact_engine, "_process_is_alive", return_value=False):
            recovered_engine = OfficeArtifactEngine(self.root)
        self.assertEqual(original_bytes, destination.read_bytes())
        self.assertEqual(original_manifest, manifest.read_bytes())
        self.assertEqual([], list(journals.glob("restore-*.json")))
        receipt = json.loads(
            (self.root / ".nexa" / "office-artifacts" / "receipts" / f"{published['receiptId']}.json")
            .read_text(encoding="utf-8")
        )
        self.assertEqual("restored", receipt["status"])
        self.assertTrue(
            hmac.compare_digest(
                receipt["integrity"]["value"],
                recovered_engine._receipt_integrity(receipt)["value"],
            )
        )

    def test_restore_journal_reconciles_receipt_written_before_candidate_state(self) -> None:
        destination = self.root / "restore-receipt-window.docx"
        candidate = self.engine.execute(self._docx_request(destination))
        published = self.engine.decide(candidate["candidateId"], "publish")
        real_write_state = self.engine._write_candidate_state
        failed_once = False

        def fail_first_restored_state(path, state):
            nonlocal failed_once
            if state.get("status") == "restored" and not failed_once:
                failed_once = True
                raise OSError("injected state checkpoint failure")
            return real_write_state(path, state)

        with mock.patch.object(
            self.engine,
            "_write_candidate_state",
            side_effect=fail_first_restored_state,
        ):
            recovered = self.engine.restore(published["receiptId"])

        self.assertTrue(failed_once)
        self.assertEqual("restored", recovered["status"])
        self.assertEqual("committed", recovered["recovery"]["status"])
        self.assertIn("checkpoint failure", recovered["recovery"]["originalError"])
        self.assertFalse(destination.exists())
        journals = self.root / ".nexa" / "office-artifacts" / "journals"
        locks = self.root / ".nexa" / "office-artifacts" / "locks"
        self.assertEqual([], list(journals.glob("restore-*.json")))
        self.assertEqual([], list(locks.glob("*.lock")))
        _, state = self.engine._load_candidate(candidate["candidateId"])
        self.assertEqual("restored", state["status"])
        receipt_path = (
            self.root / ".nexa" / "office-artifacts" / "receipts"
            / f"{published['receiptId']}.json"
        )
        receipt = json.loads(receipt_path.read_text(encoding="utf-8"))
        self.assertEqual("restored", receipt["status"])
        self.assertEqual(office_artifact_engine._sha256(receipt_path), state["receiptSha256"])

    def test_restore_recovery_rejects_valid_receipt_replay_before_mutation(self) -> None:
        first = self.engine.execute(self._docx_request(self.root / "first-replay.docx"))
        first_published = self.engine.decide(first["candidateId"], "publish")
        second = self.engine.execute(self._docx_request(self.root / "second-replay.docx"))
        second_published = self.engine.decide(second["candidateId"], "publish")

        receipts = self.root / ".nexa" / "office-artifacts" / "receipts"
        first_receipt_path = receipts / f"{first_published['receiptId']}.json"
        second_receipt_path = receipts / f"{second_published['receiptId']}.json"
        first_receipt = json.loads(first_receipt_path.read_text(encoding="utf-8"))
        first_destination = Path(first_receipt["destination"])
        published_bytes = first_destination.read_bytes()

        # Simulate internally inconsistent, but individually authenticated,
        # candidate state. Recovery must bind all three records before it
        # performs the first destructive restore step.
        first_state_path, first_state = self.engine._load_candidate(first["candidateId"])
        first_state["receiptId"] = second_published["receiptId"]
        first_state["receiptSha256"] = office_artifact_engine._sha256(second_receipt_path)
        self.engine._write_candidate_state(first_state_path, first_state)

        journal = {
            "kind": "officeArtifactRestoreJournal",
            "version": 1,
            "status": "active",
            "pid": 2_147_483_647,
            "candidateId": first["candidateId"],
            "receiptId": first_published["receiptId"],
            "destination": str(first_destination),
            "lockRolePaths": [str(first_destination)],
            "roles": [{
                "path": str(first_destination),
                "publishedSha256": first_receipt["destinationSha256"],
                "snapshot": None,
                "snapshotSha256": None,
                "restoredSha256": None,
                "existedBefore": False,
                "restored": False,
            }],
        }
        journals = self.root / ".nexa" / "office-artifacts" / "journals"
        journals.mkdir(parents=True, exist_ok=True)
        journal_path = journals / f"restore-{first_published['receiptId']}.json"
        self.engine._write_journal(journal_path, journal)

        with mock.patch.object(office_artifact_engine, "_process_is_alive", return_value=False):
            OfficeArtifactEngine(self.root)

        self.assertEqual(published_bytes, first_destination.read_bytes())
        blocked = json.loads(journal_path.read_text(encoding="utf-8"))
        self.assertEqual("recovery_blocked", blocked["status"])
        self.assertIn("binding failed", blocked["recoveryBlockers"][0])

    def test_manifest_fault_after_artifact_publish_rolls_back_every_role(self) -> None:
        import docx

        destination = self.root / "fault.docx"
        manifest = self.root / "fault-manifest.json"
        original = docx.Document()
        original.add_paragraph("Before fault")
        original.save(destination)
        before = destination.read_bytes()
        manifest.write_text('{"status":"before"}\n', encoding="utf-8")
        manifest_before = manifest.read_bytes()
        request = self._docx_request(destination)
        request["delivery"] = {"manifest": str(manifest)}
        candidate = self.engine.execute(request)
        real_publish = self.engine._journal_publish_role

        def fail_requested_manifest(journal_path, journal, staged, target, *, validate):
            if Path(target).resolve() == manifest.resolve():
                raise OSError("injected manifest fault")
            return real_publish(journal_path, journal, staged, target, validate=validate)

        with mock.patch.object(
            self.engine,
            "_journal_publish_role",
            side_effect=fail_requested_manifest,
        ):
            with self.assertRaisesRegex(OSError, "injected manifest fault"):
                self.engine.decide(candidate["candidateId"], "publish")

        self.assertEqual(before, destination.read_bytes())
        self.assertEqual(manifest_before, manifest.read_bytes())
        state = json.loads(
            (Path(candidate["candidatePath"]).parent / "state.json").read_text(encoding="utf-8")
        )
        self.assertEqual("candidate", state["status"])
        self.assertEqual([], list((self.root / ".nexa" / "office-artifacts" / "locks").glob("*.lock")))

    def test_publish_rollback_fault_retains_all_locks_until_journal_recovery(self) -> None:
        import docx

        destination = self.root / "rollback-fault.docx"
        manifest = self.root / "rollback-fault-manifest.json"
        original = docx.Document()
        original.add_paragraph("Original destination")
        original.save(destination)
        destination_before = destination.read_bytes()
        manifest.write_text('{"status":"original"}\n', encoding="utf-8")
        manifest_before = manifest.read_bytes()
        request = self._docx_request(destination)
        request["delivery"] = {"manifest": str(manifest)}
        candidate = self.engine.execute(request)
        real_publish = self.engine._journal_publish_role
        real_rollback = office_artifact_engine.rollback_published_artifact

        def fail_manifest(journal_path, journal, staged, target, *, validate):
            if Path(target).resolve() == manifest.resolve():
                raise OSError("injected manifest publication failure")
            return real_publish(journal_path, journal, staged, target, validate=validate)

        def fail_destination_rollback(target, snapshot, workspace_root):
            if Path(target).resolve() == destination.resolve():
                raise OSError("injected destination rollback failure")
            return real_rollback(target, snapshot, workspace_root)

        with mock.patch.object(
            self.engine,
            "_journal_publish_role",
            side_effect=fail_manifest,
        ), mock.patch.object(
            office_artifact_engine,
            "rollback_published_artifact",
            side_effect=fail_destination_rollback,
        ):
            with self.assertRaisesRegex(OSError, "manifest publication"):
                self.engine.decide(candidate["candidateId"], "publish")

        journals = self.root / ".nexa" / "office-artifacts" / "journals"
        locks = self.root / ".nexa" / "office-artifacts" / "locks"
        journal_paths = list(journals.glob("*.json"))
        self.assertEqual(1, len(journal_paths))
        blocked = json.loads(journal_paths[0].read_text(encoding="utf-8"))
        self.assertEqual("recovery_blocked", blocked["status"])
        self.assertEqual(0, blocked["pid"])
        self.assertEqual(2, len(list(locks.glob("*.lock"))))
        self.assertNotEqual(destination_before, destination.read_bytes())
        self.assertEqual(manifest_before, manifest.read_bytes())

        # A fresh engine in the same process must retry pid=0 blocked recovery,
        # restore every role, and only then release the protected locks.
        OfficeArtifactEngine(self.root)
        self.assertEqual(destination_before, destination.read_bytes())
        self.assertEqual(manifest_before, manifest.read_bytes())
        self.assertEqual([], list(journals.glob("*.json")))
        self.assertEqual([], list(locks.glob("*.lock")))
        _, state = self.engine._load_candidate(candidate["candidateId"])
        self.assertEqual("recovered_rolled_back", state["status"])

    def test_late_publish_checkpoint_fault_returns_committed_recovery_outcome(self) -> None:
        destination = self.root / "late-publish.docx"
        candidate = self.engine.execute(self._docx_request(destination))
        real_write_journal = self.engine._write_journal
        failed_once = False

        def fail_committed_checkpoint(path, journal):
            nonlocal failed_once
            if (
                journal.get("kind") == "officeArtifactPublishJournal"
                and journal.get("status") == "committed"
                and not failed_once
            ):
                failed_once = True
                raise OSError("injected committed journal checkpoint failure")
            return real_write_journal(path, journal)

        with mock.patch.object(
            self.engine,
            "_write_journal",
            side_effect=fail_committed_checkpoint,
        ):
            published = self.engine.decide(candidate["candidateId"], "publish")

        self.assertTrue(failed_once)
        self.assertEqual("published", published["status"])
        self.assertEqual("committed", published["recovery"]["status"])
        self.assertIn("checkpoint failure", published["recovery"]["originalError"])
        self.assertTrue(destination.is_file())
        journals = self.root / ".nexa" / "office-artifacts" / "journals"
        locks = self.root / ".nexa" / "office-artifacts" / "locks"
        self.assertEqual([], list(journals.glob("*.json")))
        self.assertEqual([], list(locks.glob("*.lock")))

    def test_startup_recovers_crash_after_artifact_publication(self) -> None:
        destination = self.root / "crash.docx"
        destination.write_bytes(b"before")
        snapshot = office_artifact_engine.snapshot_file(destination, self.root)
        self.assertIsNotNone(snapshot)
        destination.write_bytes(b"published")
        candidate_id = "b" * 32
        journals = self.root / ".nexa" / "office-artifacts" / "journals"
        journals.mkdir(parents=True, exist_ok=True)
        journal = {
            "kind": "officeArtifactPublishJournal",
            "version": 1,
            "candidateId": candidate_id,
            "status": "active",
            "pid": 2_147_483_647,
            "destination": str(destination),
            "roles": [{
                "path": str(destination),
                "existedBefore": True,
                "preexistingSha256": office_artifact_engine.hashlib.sha256(b"before").hexdigest(),
                "snapshot": str(snapshot),
                "snapshotSha256": office_artifact_engine._sha256(snapshot),
                "intendedSha256": office_artifact_engine.hashlib.sha256(b"published").hexdigest(),
                "published": True,
            }],
        }
        journal["integrity"] = self.engine._journal_integrity(journal)
        journal_path = journals / f"{candidate_id}.json"
        journal_path.write_text(json.dumps(journal), encoding="utf-8")

        OfficeArtifactEngine(self.root)
        self.assertEqual(b"before", destination.read_bytes())
        self.assertFalse(journal_path.exists())

    def test_startup_quarantines_forged_journal_without_touching_user_files(self) -> None:
        destination = self.root / "important.docx"
        destination.write_bytes(b"keep me")
        candidate_id = "c" * 32
        journals = self.root / ".nexa" / "office-artifacts" / "journals"
        journals.mkdir(parents=True, exist_ok=True)
        forged = {
            "kind": "officeArtifactPublishJournal",
            "version": 1,
            "candidateId": candidate_id,
            "status": "active",
            "pid": 2_147_483_647,
            "destination": str(destination),
            "roles": [{
                "path": str(destination),
                "existedBefore": False,
                "snapshot": None,
                "intendedSha256": office_artifact_engine._sha256(destination),
            }],
            "integrity": {"algorithm": "HMAC-SHA256", "value": "0" * 64},
        }
        journal_path = journals / f"{candidate_id}.json"
        journal_path.write_text(json.dumps(forged), encoding="utf-8")

        OfficeArtifactEngine(self.root)
        self.assertEqual(b"keep me", destination.read_bytes())
        self.assertFalse(journal_path.exists())
        self.assertEqual(1, len(list(journals.glob(f"{candidate_id}.json.invalid-*.quarantine"))))

    def test_startup_removes_signed_dead_orphan_lock_without_journal(self) -> None:
        destination = self.root / "orphan.docx"
        lock_path = self.engine._acquire_destination_lock(destination, "d" * 32)
        lock = json.loads(lock_path.read_text(encoding="utf-8"))
        lock["pid"] = 2_147_483_647
        lock["integrity"] = self.engine._lock_integrity(lock)
        lock_path.write_text(json.dumps(lock), encoding="utf-8")

        with mock.patch.object(office_artifact_engine, "_process_is_alive", return_value=False):
            OfficeArtifactEngine(self.root)
        self.assertFalse(lock_path.exists())

    def test_assessment_reports_unsatisfied_publish_render_backend(self) -> None:
        request = self._docx_request(self.root / "publish.docx")
        request["guarantees"]["quality"] = "publish"
        request["guarantees"].pop("render")
        assessment = self.engine.assess(request)

        backend_status = {item["id"]: item for item in self.engine.capabilities()["backends"]}
        if all(
            backend_status[adapter]["status"] == "ready"
            for adapter in ("libreoffice", "openxml-sdk")
        ):
            self.assertTrue(assessment["ready"])
        else:
            self.assertFalse(assessment["ready"])
            self.assertIn("backend.unavailable", {item["code"] for item in assessment["blockers"]})
        render_step = next(
            step for step in assessment["adapterPlan"]["steps"]
            if step["step"] == "render"
        )
        self.assertEqual("libreoffice", render_step["adapter"])

        request = self._docx_request(self.root / "publish-bypass.docx")
        request["guarantees"].update({"quality": "publish", "render": "none"})
        with self.assertRaisesRegex(OfficeArtifactError, "requires final candidate render"):
            self.engine.assess(request)

    def test_path_roles_and_candidate_ids_are_validated(self) -> None:
        destination = self.root / "conflict.docx"
        request = self._docx_request(destination)
        request["delivery"] = {"manifest": str(destination)}
        with self.assertRaisesRegex(OfficeArtifactError, "manifest must be distinct"):
            self.engine.execute(request)
        request = self._docx_request(destination)
        request["delivery"] = {
            "manifest": str(
                self.root
                / ".nexa"
                / "office-artifacts"
                / "receipts"
                / f"{'a' * 32}.json"
            )
        }
        with self.assertRaisesRegex(OfficeArtifactError, "reserved .nexa state"):
            self.engine.assess(request)
        xlsx_destination = self.root / "conflict.xlsx"
        spec = self.root / "conflict-spec.json"
        spec.write_text(
            json.dumps({"sheets": [{"name": "Sheet1", "rows": [["value"]]}]}),
            encoding="utf-8",
        )
        request = {
            "requestVersion": 2,
            "format": "xlsx",
            "intent": "create",
            "destination": str(xlsx_destination),
            "operations": [{"op": "create", "spec": str(spec)}],
        }
        request["delivery"] = {"manifest": str(xlsx_destination.with_suffix(".xlsx.qa.json"))}
        with self.assertRaisesRegex(OfficeArtifactError, "XLSX QA sidecar"):
            self.engine.assess(request)
        request["delivery"] = {"manifest": str(spec)}
        with self.assertRaisesRegex(OfficeArtifactError, "cannot overwrite request input"):
            self.engine.assess(request)
        with self.assertRaisesRegex(OfficeArtifactError, "invalid candidate id"):
            self.engine.decide("../escape", "discard")

    def test_request_operation_and_contract_schemas_reject_unknown_fields(self) -> None:
        destination = self.root / "strict.docx"
        request = self._docx_request(destination)
        request["typo"] = True
        with self.assertRaisesRegex(OfficeArtifactError, "unknown field.*request"):
            self.engine.assess(request)

        request = self._docx_request(destination)
        request["operations"][0]["titel"] = "typo"
        with self.assertRaisesRegex(OfficeArtifactError, r"operations\[0\]"):
            self.engine.assess(request)

        request = self._docx_request(destination)
        request["validation"] = {"contractVersion": 2, "required_tex": ["missing t"]}
        with self.assertRaisesRegex(OfficeArtifactError, "unknown field.*validation"):
            self.engine.assess(request)

        request = self._docx_request(destination)
        request["validation"] = 7
        with self.assertRaisesRegex(OfficeArtifactError, "validation must be"):
            self.engine.assess(request)

        request = self._docx_request(destination)
        request["validation"] = {"contractVersion": 2.9, "required_text": "BA"}
        with self.assertRaisesRegex(OfficeArtifactError, "contractVersion must be 2"):
            self.engine.assess(request)

        request = self._docx_request(destination)
        request["validation"] = {"contractVersion": 2, "required_text": "BA"}
        with self.assertRaisesRegex(OfficeArtifactError, "array of strings"):
            self.engine.assess(request)

        contract_path = self.root / "missing-version-contract.json"
        contract_path.write_text(json.dumps({"required_text": ["Verified body"]}), encoding="utf-8")
        request = self._docx_request(destination)
        request["validation"] = str(contract_path)
        with self.assertRaisesRegex(OfficeArtifactError, "contractVersion is required"):
            self.engine.assess(request)

        request = self._docx_request(destination)
        request["operations"] = [{
            "op": "replace",
            "find": "a",
            "replace": "b",
            "allowStyleMerge": "false",
        }]
        with self.assertRaisesRegex(OfficeArtifactError, "allowStyleMerge must be a boolean"):
            self.engine.assess(request)

        from pptx import Presentation

        pptx_source = self.root / "strict-source.pptx"
        Presentation().save(pptx_source)
        request = {
            "requestVersion": 2,
            "format": "pptx",
            "intent": "modify",
            "source": str(pptx_source),
            "preconditions": {"sourceSha256": hashlib.sha256(pptx_source.read_bytes()).hexdigest()},
            "destination": str(self.root / "strict-output.pptx"),
            "operations": [{"op": "replace", "find": "a", "replace": "b", "scope": "body"}],
        }
        with self.assertRaisesRegex(OfficeArtifactError, "unknown field.*scope"):
            self.engine.assess(request)

        request = self._docx_request(destination)
        request["requestVersion"] = 2.9
        with self.assertRaisesRegex(OfficeArtifactError, "requestVersion must be 2"):
            self.engine.assess(request)

    def test_source_sha_precondition_is_format_wide(self) -> None:
        source = self.root / "precondition.docx"
        source.write_bytes(b"changed")
        request = {
            "requestVersion": 2,
            "format": "docx",
            "intent": "verify",
            "source": str(source),
            "destination": str(self.root / "verified.docx"),
            "operations": [],
        }
        with self.assertRaisesRegex(OfficeArtifactError, "sourceSha256 from inspect"):
            self.engine.assess(request)
        request["preconditions"] = {"sourceSha256": "0" * 64}
        with self.assertRaisesRegex(OfficeArtifactError, "source SHA-256"):
            self.engine.assess(request)

        request["preconditions"] = {
            "sourceSha256": hashlib.sha256(source.read_bytes()).hexdigest()
        }
        request["delivery"] = {"mode": "publish"}
        with self.assertRaisesRegex(OfficeArtifactError, "verify is read-only"):
            self.engine.assess(request)

        request.pop("delivery")
        source_doc = self.root / "read-only-source.docx"
        import docx

        document = docx.Document()
        document.add_paragraph("Read-only evidence")
        document.save(source_doc)
        request.update({
            "source": str(source_doc),
            "destination": str(self.root / "must-never-publish.docx"),
            "preconditions": {
                "sourceSha256": hashlib.sha256(source_doc.read_bytes()).hexdigest()
            },
            "validation": {
                "contractVersion": 2,
                "required_text": ["Read-only evidence"],
            },
        })
        verified = self.engine.execute(request)
        with self.assertRaisesRegex(OfficeArtifactError, "evidence-only"):
            self.engine.decide(verified["candidateId"], "publish")
        self.assertFalse(Path(request["destination"]).exists())

    def test_docx_replace_verifier_rejects_unrelated_header_drift(self) -> None:
        import docx

        source = self.root / "strict-scope-source.docx"
        document = docx.Document()
        document.add_paragraph("Body target")
        document.sections[0].header.paragraphs[0].text = "Header must stay"
        document.save(source)
        request = {
            "requestVersion": 2,
            "format": "docx",
            "intent": "modify",
            "source": str(source),
            "preconditions": {"sourceSha256": hashlib.sha256(source.read_bytes()).hexdigest()},
            "destination": str(self.root / "strict-scope-output.docx"),
            "operations": [{
                "op": "replace",
                "find": "Body target",
                "replace": "Body updated",
                "scope": "body",
                "expectedMatches": 1,
            }],
            "preconditions": {"sourceSha256": hashlib.sha256(source.read_bytes()).hexdigest()},
            "guarantees": {"quality": "standard", "preservation": "strict", "render": "none"},
            "validation": {"contractVersion": 2, "required_text": ["Body updated"]},
        }
        candidate = self.engine.execute(request)
        candidate_document = docx.Document(candidate["candidatePath"])
        self.assertEqual("Body updated", candidate_document.paragraphs[0].text)
        self.assertEqual(
            "Header must stay",
            candidate_document.sections[0].header.paragraphs[0].text,
        )

        real_run_editor = office_artifact_service._run_editor

        def inject_header_drift(path, command, arguments, actions, workspace_root, *, timeout=180):
            result = real_run_editor(
                path,
                command,
                arguments,
                actions,
                workspace_root,
                timeout=timeout,
            )
            if command == "replace":
                staged = path.with_name(f"{path.name}.scope-fixture")
                with zipfile.ZipFile(path) as archive, zipfile.ZipFile(staged, "w") as output:
                    for info in archive.infolist():
                        data = archive.read(info.filename)
                        if info.filename.startswith("word/header") and info.filename.endswith(".xml"):
                            data = data.replace(b"Header must stay", b"Header was drifted")
                        output.writestr(info, data)
                staged.replace(path)
            return result

        request["destination"] = str(self.root / "strict-scope-injected.docx")
        with mock.patch.object(
            office_artifact_service,
            "_run_editor",
            side_effect=inject_header_drift,
        ):
            with self.assertRaisesRegex(OfficeArtifactError, "outside the requested scope"):
                self.engine.execute(request)
        self.assertFalse(Path(request["destination"]).exists())

    def test_docx_inspect_reads_content_controls_and_dotm_without_python_docx_loader(self) -> None:
        import docx

        source = self.root / "content-control-source.docx"
        document = docx.Document()
        document.add_paragraph("Controlled insight")
        document.save(source)
        request = {
            "requestVersion": 2,
            "format": "docx",
            "intent": "modify",
            "source": str(source),
            "destination": str(self.root / "content-control.docx"),
            "operations": [{
                "op": "wrap_content_control",
                "find": "Controlled insight",
                "tag": "decision-insight",
            }],
            "preconditions": {"sourceSha256": hashlib.sha256(source.read_bytes()).hexdigest()},
            "guarantees": {"quality": "standard", "preservation": "strict", "render": "none"},
            "validation": {"contractVersion": 2, "required_text": ["Controlled insight"]},
        }
        candidate = self.engine.execute(request)
        inspection = self.engine.inspect(candidate["candidatePath"], "docx")
        self.assertIn("Controlled insight", inspection["profile"]["textPreview"])
        self.assertEqual("direct-openxml", inspection["profile"]["profileEngine"])
        self.assertEqual("decision-insight", inspection["profile"]["contentControls"][0]["tag"])

        dotm = self.root / "content-control.dotm"
        with zipfile.ZipFile(candidate["candidatePath"]) as archive, zipfile.ZipFile(dotm, "w") as output:
            for info in archive.infolist():
                data = archive.read(info.filename)
                if info.filename == "[Content_Types].xml":
                    data = data.replace(
                        b"application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml",
                        b"application/vnd.ms-word.template.macroEnabledTemplate.main+xml",
                    )
                output.writestr(info, data)
        dotm_inspection = self.engine.inspect(str(dotm), "docx")
        self.assertIn("Controlled insight", dotm_inspection["profile"]["textPreview"])

    def test_xlsx_chartsheet_can_be_inspected_and_verified_without_row_iteration(self) -> None:
        import openpyxl
        from openpyxl.chart import BarChart, Reference

        source = self.root / "chartsheet-source.xlsx"
        workbook = openpyxl.Workbook()
        worksheet = workbook.active
        worksheet.title = "Data"
        worksheet.append(["Quarter", "Value"])
        worksheet.append(["Q1", 10])
        worksheet.append(["Q2", 20])
        chart = BarChart()
        chart.add_data(
            Reference(worksheet, min_col=2, min_row=1, max_row=3),
            titles_from_data=True,
        )
        chart.set_categories(Reference(worksheet, min_col=1, min_row=2, max_row=3))
        workbook.create_chartsheet("Decision chart").add_chart(chart)
        workbook.save(source)
        workbook.close()
        request = {
            "requestVersion": 2,
            "format": "xlsx",
            "intent": "verify",
            "source": str(source),
            "destination": str(self.root / "chartsheet-verified.xlsx"),
            "operations": [],
            "preconditions": {"sourceSha256": hashlib.sha256(source.read_bytes()).hexdigest()},
            "guarantees": {
                "quality": "standard",
                "preservation": "strict",
                "calculation": "static",
                "render": "none",
            },
            "validation": {
                "contractVersion": 2,
                "required_sheets": ["Data", "Decision chart"],
            },
        }
        inspection = self.engine.inspect(str(source), "xlsx")
        self.assertEqual("pass", inspection["structural"]["status"])
        self.assertEqual(
            [("Data", "worksheet"), ("Decision chart", "chartsheet")],
            [
                (sheet["name"], sheet["type"])
                for sheet in inspection["profile"]["sheet_details"]
            ],
        )
        candidate = self.engine.execute(request)
        self.assertEqual("pass", candidate["validation"]["backend"]["contract"]["status"])

    def test_macro_enabled_format_families_preserve_vba_bytes_end_to_end(self) -> None:
        from xml.etree import ElementTree as ET
        import docx
        import openpyxl
        from pptx import Presentation

        fixtures = []
        docx_base = self.root / "macro-base.docx"
        document = docx.Document()
        document.add_paragraph("Macro document")
        document.save(docx_base)
        fixtures.append(("docx", docx_base, ".docm", "word", "document.xml"))
        xlsx_base = self.root / "macro-base.xlsx"
        workbook = openpyxl.Workbook()
        workbook.active["A1"] = "Macro workbook"
        workbook.save(xlsx_base)
        workbook.close()
        fixtures.append(("xlsx", xlsx_base, ".xlsm", "xl", "workbook.xml"))
        pptx_base = self.root / "macro-base.pptx"
        presentation = Presentation()
        presentation.slides.add_slide(presentation.slide_layouts[6])
        presentation.save(pptx_base)
        fixtures.append(("pptx", pptx_base, ".pptm", "ppt", "presentation.xml"))

        macro_content_types = {
            "docx": "application/vnd.ms-word.document.macroEnabled.main+xml",
            "xlsx": "application/vnd.ms-excel.sheet.macroEnabled.main+xml",
            "pptx": "application/vnd.ms-powerpoint.presentation.macroEnabled.main+xml",
        }
        payload = b"NEXA-VBA-PROJECT-BYTES"
        for artifact_format, base, extension, folder, main_name in fixtures:
            source = base.with_suffix(extension)
            rels_name = f"{folder}/_rels/{main_name}.rels"
            vba_part = f"{folder}/vbaProject.bin"
            with zipfile.ZipFile(base) as archive:
                content_types = ET.fromstring(archive.read("[Content_Types].xml"))
                for override in content_types:
                    if override.attrib.get("PartName") == f"/{folder}/{main_name}":
                        override.set("ContentType", macro_content_types[artifact_format])
                ET.SubElement(content_types, "{http://schemas.openxmlformats.org/package/2006/content-types}Override", {
                    "PartName": f"/{vba_part}",
                    "ContentType": "application/vnd.ms-office.vbaProject",
                })
                relationships = ET.fromstring(archive.read(rels_name))
                ET.SubElement(relationships, "{http://schemas.openxmlformats.org/package/2006/relationships}Relationship", {
                    "Id": "rIdNexaVba",
                    "Type": "http://schemas.microsoft.com/office/2006/relationships/vbaProject",
                    "Target": "vbaProject.bin",
                })
                with zipfile.ZipFile(source, "w") as output:
                    for info in archive.infolist():
                        data = archive.read(info.filename)
                        if info.filename == "[Content_Types].xml":
                            data = ET.tostring(content_types, encoding="utf-8", xml_declaration=True)
                        elif info.filename == rels_name:
                            data = ET.tostring(relationships, encoding="utf-8", xml_declaration=True)
                        output.writestr(info, data)
                    output.writestr(vba_part, payload)
            outcome = self.engine.execute({
                "requestVersion": 2,
                "format": artifact_format,
                "intent": "verify",
                "source": str(source),
                "preconditions": {"sourceSha256": hashlib.sha256(source.read_bytes()).hexdigest()},
                "destination": str(self.root / f"verified{extension}"),
                "operations": [],
                "guarantees": {"quality": "standard", "preservation": "strict", "render": "none"},
            })
            candidate = Path(outcome["candidatePath"])
            self.assertEqual(extension, candidate.suffix)
            with zipfile.ZipFile(candidate) as archive:
                self.assertEqual(payload, archive.read(vba_part))
            self.assertTrue(outcome["preservationEvidence"]["verified"])

    def test_template_format_families_preserve_exact_package_parts(self) -> None:
        from xml.etree import ElementTree as ET
        import docx
        import openpyxl
        from pptx import Presentation

        base_docx = self.root / "template.docx"
        doc = docx.Document()
        doc.add_paragraph("Template Word")
        doc.save(base_docx)
        base_xlsx = self.root / "template.xlsx"
        workbook = openpyxl.Workbook()
        workbook.active["A1"] = "Template Excel"
        workbook.save(base_xlsx)
        workbook.close()
        base_pptx = self.root / "template.pptx"
        presentation = Presentation()
        presentation.slides.add_slide(presentation.slide_layouts[6])
        presentation.save(base_pptx)
        fixtures = [
            ("docx", base_docx, ".dotx", "/word/document.xml", "application/vnd.openxmlformats-officedocument.wordprocessingml.template.main+xml"),
            ("xlsx", base_xlsx, ".xltx", "/xl/workbook.xml", "application/vnd.openxmlformats-officedocument.spreadsheetml.template.main+xml"),
            ("pptx", base_pptx, ".potx", "/ppt/presentation.xml", "application/vnd.openxmlformats-officedocument.presentationml.template.main+xml"),
        ]
        for artifact_format, base, extension, part_name, content_type in fixtures:
            source = base.with_suffix(extension)
            with zipfile.ZipFile(base) as archive, zipfile.ZipFile(source, "w") as output:
                content_types = ET.fromstring(archive.read("[Content_Types].xml"))
                for override in content_types:
                    if override.attrib.get("PartName") == part_name:
                        override.set("ContentType", content_type)
                for info in archive.infolist():
                    data = (
                        ET.tostring(content_types, encoding="utf-8", xml_declaration=True)
                        if info.filename == "[Content_Types].xml"
                        else archive.read(info.filename)
                    )
                    output.writestr(info, data)
            with zipfile.ZipFile(source) as archive:
                source_hashes = {
                    name: hashlib.sha256(archive.read(name)).hexdigest()
                    for name in archive.namelist()
                }
            outcome = self.engine.execute({
                "requestVersion": 2,
                "format": artifact_format,
                "intent": "verify",
                "source": str(source),
                "preconditions": {"sourceSha256": hashlib.sha256(source.read_bytes()).hexdigest()},
                "destination": str(self.root / f"verified{extension}"),
                "operations": [],
                "guarantees": {"quality": "standard", "preservation": "strict", "render": "none"},
            })
            candidate = Path(outcome["candidatePath"])
            self.assertEqual(extension, candidate.suffix)
            with zipfile.ZipFile(candidate) as archive:
                candidate_hashes = {
                    name: hashlib.sha256(archive.read(name)).hexdigest()
                    for name in archive.namelist()
                }
            self.assertEqual(source_hashes, candidate_hashes)

    def test_capabilities_validate_external_adapter_declarations_without_loading_code(self) -> None:
        adapter_dir = self.root / ".nexa" / "office-adapters"
        adapter_dir.mkdir(parents=True)
        (adapter_dir / "valid.json").write_text(json.dumps({
            "adapterVersion": 1,
            "id": "example-live",
            "deployment": "live-officejs",
            "formats": ["docx"],
            "operations": ["set_text"],
            "guarantees": {"preservation": ["native"], "calculation": [], "render": []},
            "limitations": ["declaration only"],
            "requires": ["officejs-host"],
        }), encoding="utf-8")
        (adapter_dir / "invalid.json").write_text('{"adapterVersion":2,"id":"Bad ID"}', encoding="utf-8")

        declarations = self.engine.capabilities()["externalAdapterDeclarations"]
        by_name = {Path(item["manifestPath"]).name: item for item in declarations}
        self.assertEqual("declared-not-loaded", by_name["valid.json"]["status"])
        self.assertEqual("invalid", by_name["invalid.json"]["status"])

    def test_capabilities_bind_hash_locked_python_supply_chain(self) -> None:
        lock = self.engine.capabilities()["pythonDependencyLock"]
        self.assertEqual("ready", lock["status"])
        self.assertTrue(lock["requireHashes"])
        self.assertRegex(lock["lockSha256"], r"^[0-9a-f]{64}$")
        self.assertRegex(lock["sbomSha256"], r"^[0-9a-f]{64}$")
        content = Path(__file__).with_name("requirements.lock").read_text(encoding="utf-8")
        self.assertIn("lxml==6.1.1", content)
        self.assertIn("pillow==12.3.0", content)
        self.assertIn('pywin32==312 ; sys_platform == "win32"', content)
        self.assertGreaterEqual(content.count("--hash=sha256:"), 200)

    def test_v2_translates_to_internal_plan_not_legacy_job_protocol(self) -> None:
        request = office_artifact_engine.ArtifactRequest.from_dict(
            self._docx_request(self.root / "plan.docx"),
            self.root,
        )
        payload = self.engine._execution_plan_payload(
            request,
            self.root / "candidate.docx",
            self.root / "candidate-manifest.json",
        )
        self.assertEqual(1, payload["planVersion"])
        self.assertNotIn("jobVersion", payload)
        plan = office_artifact_engine.OfficeExecutionPlan.from_internal_dict(payload, self.root)
        self.assertEqual("docx", plan.format)

    def test_xlsx_render_evidence_requires_surface_manifest_bound_to_candidate_sha(self) -> None:
        try:
            from PIL import Image
        except ImportError:
            self.skipTest("Pillow is not installed")
        candidate_dir = self.root / ".nexa" / "office-artifacts" / "candidates" / ("a" * 32)
        render_dir = candidate_dir / "artifact-rendered"
        render_dir.mkdir(parents=True)
        candidate = candidate_dir / "artifact.xlsx"
        candidate.write_bytes(b"candidate-bytes")
        pages = [render_dir / "sheet-001-page-1.png", render_dir / "sheet-002-page-1.png"]
        for index, page in enumerate(pages, start=1):
            Image.new("RGB", (640, 360), (40 * index, 80, 140)).save(page)
        manifest = {
            "kind": "officeRenderManifest",
            "format": "xlsx",
            "artifactSha256": office_artifact_engine._sha256(candidate),
            "expectedSurfaces": 2,
            "renderedSurfaces": 2,
            "complete": True,
        }
        (render_dir / "render-manifest.json").write_text(json.dumps(manifest), encoding="utf-8")
        execution = {"renderedPreviews": [str(page) for page in pages]}

        evidence = self.engine._render_evidence(candidate, execution, "xlsx", "all")
        self.assertTrue(evidence["complete"])
        self.assertEqual(2, evidence["expectedSurfaces"])
        self.assertEqual(2, evidence["renderedSurfaces"])

        candidate.write_bytes(b"tampered")
        stale = self.engine._render_evidence(candidate, execution, "xlsx", "all")
        self.assertFalse(stale["complete"])

    def test_discard_removes_only_owned_candidate_directory(self) -> None:
        destination = self.root / "discarded.docx"
        keep = self.root / "keep.txt"
        keep.write_text("safe", encoding="utf-8")
        candidate = self.engine.execute(self._docx_request(destination))
        candidate_dir = Path(candidate["candidatePath"]).parent

        discarded = self.engine.decide(candidate["candidateId"], "discard")
        self.assertEqual("discarded", discarded["status"])
        self.assertFalse(candidate_dir.exists())
        self.assertEqual("safe", keep.read_text(encoding="utf-8"))

    def test_published_manifest_is_machine_readable_outcome(self) -> None:
        destination = self.root / "manifested.docx"
        request = self._docx_request(destination)
        request["delivery"] = {"mode": "publish"}
        outcome = self.engine.execute(request)

        manifest = json.loads((self.root / "manifested.docx.manifest.json").read_text(encoding="utf-8"))
        self.assertEqual("published", outcome["status"])
        self.assertEqual(outcome["receiptId"], manifest["receiptId"])

    def test_typed_xlsx_edits_are_literal_formula_safe_and_part_precise(self) -> None:
        try:
            import openpyxl  # type: ignore
        except ImportError:
            self.skipTest("openpyxl is not installed")
        import zipfile

        source = self.root / "source.xlsx"
        destination = self.root / "result.xlsx"
        workbook = openpyxl.Workbook()
        workbook.active.title = "Inputs"
        workbook.active["A1"] = "old"
        workbook.create_sheet("Untouched")["A1"] = "stable"
        workbook.save(source)
        workbook.close()
        with zipfile.ZipFile(source) as archive:
            untouched_before = archive.read("xl/worksheets/sheet2.xml")

        outcome = self.engine.execute({
            "requestVersion": 2,
            "format": "xlsx",
            "intent": "modify",
            "source": str(source),
            "preconditions": {"sourceSha256": hashlib.sha256(source.read_bytes()).hexdigest()},
            "destination": str(destination),
            "operations": [
                {"op": "set_value", "sheet": "inputs", "cell": "A1", "value": "=WEBSERVICE(\"https://example.invalid\")"},
                {"op": "set_formula", "sheet": "Inputs", "cell": "B1", "formula": "=1+1"},
                {"op": "set_range", "sheet": "Inputs", "range": "A2:B2", "values": [[3, 4]]},
                {"op": "set_style", "sheet": "Inputs", "range": "A1:B2", "styleId": 0},
            ],
            "guarantees": {"quality": "standard", "preservation": "strict", "render": "none"},
        })

        candidate = Path(outcome["candidatePath"])
        self.assertFalse(destination.exists())
        self.assertEqual("static", outcome["calculationEvidence"]["profile"])
        self.assertFalse(outcome["calculationEvidence"]["excelNative"])
        workbook = openpyxl.load_workbook(candidate, data_only=False)
        try:
            sheet = workbook["Inputs"]
            self.assertEqual("s", sheet["A1"].data_type)
            self.assertEqual("=WEBSERVICE(\"https://example.invalid\")", sheet["A1"].value)
            self.assertEqual("f", sheet["B1"].data_type)
            self.assertEqual("=1+1", sheet["B1"].value)
            self.assertEqual([3, 4], [sheet["A2"].value, sheet["B2"].value])
        finally:
            workbook.close()
        with zipfile.ZipFile(candidate) as archive:
            self.assertEqual(untouched_before, archive.read("xl/worksheets/sheet2.xml"))

    def test_xlsx_object_operations_run_through_strict_candidate_lifecycle(self) -> None:
        import openpyxl

        source = self.root / "objects.xlsx"
        workbook = openpyxl.Workbook()
        sheet = workbook.active
        sheet.title = "Summary"
        sheet.append(["Metric", "Value"])
        sheet.append(["Revenue", 100])
        sheet.append(["Cost", 40])
        workbook.create_sheet("Inputs")["A1"] = "=Summary!B2"
        workbook.save(source)
        workbook.close()
        outcome = self.engine.execute({
            "requestVersion": 2,
            "format": "xlsx",
            "intent": "modify",
            "source": str(source),
            "preconditions": {"sourceSha256": hashlib.sha256(source.read_bytes()).hexdigest()},
            "destination": str(self.root / "objects-result.xlsx"),
            "operations": [
                {"op": "rename_sheet", "sheet": "Summary", "newName": "Data"},
                {"op": "set_defined_name", "name": "RevenueCell", "formula": "Data!$B$2"},
                {"op": "set_data_validation", "sheet": "Data", "range": "B2:B3", "validationType": "whole", "formula1": "0", "formula2": "1000"},
                {"op": "create_table", "sheet": "Data", "range": "A1:B3", "name": "DataTable"},
                {"op": "set_number_format", "sheet": "Data", "range": "B2:B3", "formatCode": "0.00"},
            ],
            "guarantees": {"quality": "standard", "preservation": "strict", "render": "none"},
        })
        self.assertTrue(outcome["preservationEvidence"]["verified"])
        workbook = openpyxl.load_workbook(Path(outcome["candidatePath"]), data_only=False)
        self.assertEqual("=Data!B2", workbook["Inputs"]["A1"].value)
        self.assertIn("DataTable", workbook["Data"].tables)
        self.assertEqual("0.00", workbook["Data"]["B2"].number_format)
        workbook.close()

    def test_xlsx_chart_data_is_atomic_across_source_formula_and_cache(self) -> None:
        import openpyxl
        from openpyxl.chart import BarChart, Reference

        source = self.root / "chart-source.xlsx"
        workbook = openpyxl.Workbook()
        sheet = workbook.active
        sheet.title = "Summary"
        sheet.append(["Region", "Amount"])
        sheet.append(["North", 100])
        sheet.append(["South", 80])
        chart = BarChart()
        chart.add_data(
            Reference(sheet, min_col=2, min_row=1, max_row=3),
            titles_from_data=True,
        )
        chart.set_categories(Reference(sheet, min_col=1, min_row=2, max_row=3))
        sheet.add_chart(chart, "D2")
        workbook.save(source)
        workbook.close()
        destination = self.root / "chart-result.xlsx"

        outcome = self.engine.execute({
            "requestVersion": 2,
            "format": "xlsx",
            "intent": "modify",
            "source": str(source),
            "preconditions": {"sourceSha256": hashlib.sha256(source.read_bytes()).hexdigest()},
            "destination": str(destination),
            "operations": [{
                "op": "set_chart_data",
                "chartPart": "xl/charts/chart1.xml",
                "seriesIndex": 1,
                "seriesName": "Updated amount",
                "categories": ["East", "West"],
                "values": [125, 95],
            }],
            "guarantees": {"quality": "standard", "preservation": "strict", "render": "none"},
        })

        self.assertFalse(destination.exists())
        self.assertTrue(outcome["preservationEvidence"]["verified"])
        candidate = Path(outcome["candidatePath"])
        workbook = openpyxl.load_workbook(candidate, data_only=False)
        self.assertEqual(["East", "West"], [workbook["Summary"]["A2"].value, workbook["Summary"]["A3"].value])
        self.assertEqual([125, 95], [workbook["Summary"]["B2"].value, workbook["Summary"]["B3"].value])
        self.assertEqual("Updated amount", workbook["Summary"]["B1"].value)
        workbook.close()
        inspection = self.engine.inspect(str(candidate), "xlsx")
        self.assertEqual([], inspection["profile"]["chart_validation_errors"])

    def test_assessment_requires_excel_native_for_dynamic_array_formulas(self) -> None:
        try:
            import openpyxl  # type: ignore
        except ImportError:
            self.skipTest("openpyxl is not installed")
        source = self.root / "dynamic.xlsx"
        workbook = openpyxl.Workbook()
        workbook.active["A1"] = 1
        workbook.active["A2"] = "=_xlfn.FILTER(A1:A1,A1:A1>0)"
        workbook.save(source)
        workbook.close()
        request = {
            "requestVersion": 2,
            "format": "xlsx",
            "intent": "modify",
            "source": str(source),
            "preconditions": {"sourceSha256": hashlib.sha256(source.read_bytes()).hexdigest()},
            "destination": str(self.root / "dynamic-result.xlsx"),
            "operations": [{"op": "set_value", "sheet": "Sheet", "cell": "B1", "value": 2}],
            "guarantees": {"calculation": "compatible", "quality": "standard", "render": "none"},
        }

        assessment = self.engine.assess(request)

        self.assertFalse(assessment["ready"])
        self.assertIn(
            "calculation.excel_native_required",
            {blocker["code"] for blocker in assessment["blockers"]},
        )
        self.assertIn(
            "function:FILTER",
            assessment["sourceProfile"]["formulaProfile"]["nativeFeatures"],
        )

    def test_docx_spec_v2_runs_through_candidate_validation(self) -> None:
        spec_path = self.root / "report-spec.json"
        destination = self.root / "professional.docx"
        spec_path.write_text(json.dumps({
            "schemaVersion": 2,
            "preset": "memo",
            "title": "Decision memo",
            "language": "en-US",
            "footer": {"text": "Internal", "pageNumber": True},
            "blocks": [
                {"type": "heading", "level": 1, "text": "Recommendation"},
                {"type": "paragraph", "text": "Approve the controlled rollout."},
                {
                    "type": "table",
                    "headers": ["Owner", "Date"],
                    "rows": [["Operations", "2026-08-20"]],
                    "columnWidths": [3.0, 2.0],
                    "repeatHeader": True,
                },
            ],
        }), encoding="utf-8")
        outcome = self.engine.execute({
            "requestVersion": 2,
            "format": "docx",
            "intent": "create",
            "destination": str(destination),
            "operations": [{"op": "create", "spec": str(spec_path)}],
            "guarantees": {"quality": "standard", "render": "none"},
            "validation": {
                "contractVersion": 2,
                "required_text": ["Approve the controlled rollout."],
                "min_tables": 1,
                "required_styles": ["Heading 1"],
                "no_heading_level_skips": True,
                "require_table_header_rows": True,
                "require_fixed_table_layout": True,
                "required_language": "en-US",
            },
        })
        self.assertEqual("candidate", outcome["status"])
        self.assertFalse(destination.exists())
        self.assertTrue(Path(outcome["candidatePath"]).exists())

    def test_pptx_exact_clone_copies_chart_workbook_and_targets_shape_by_id(self) -> None:
        try:
            from pptx import Presentation
            from pptx.chart.data import ChartData
            from pptx.enum.chart import XL_CHART_TYPE
            from pptx.util import Inches
        except ImportError:
            self.skipTest("python-pptx is not installed")
        import zipfile

        source = self.root / "source-deck.pptx"
        destination = self.root / "result-deck.pptx"
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        title = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(5), Inches(0.7))
        title.name = "Decision title"
        title.text = "Original decision"
        chart_data = ChartData()
        chart_data.categories = ["A", "B"]
        chart_data.add_series("Revenue", (10, 20))
        slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Inches(1), Inches(1.5), Inches(6), Inches(3.5),
            chart_data,
        )
        presentation.save(source)
        shape_id = title.shape_id

        outcome = self.engine.execute({
            "requestVersion": 2,
            "format": "pptx",
            "intent": "modify",
            "source": str(source),
            "preconditions": {"sourceSha256": hashlib.sha256(source.read_bytes()).hexdigest()},
            "destination": str(destination),
            "operations": [
                {"op": "clone_slide", "slideIndex": 1},
                {"op": "set_text", "slideIndex": 2, "shapeId": shape_id, "text": "Cloned decision"},
                {"op": "set_transition", "slideIndex": 2, "transition": "fade", "speed": "fast"},
            ],
            "guarantees": {"quality": "standard", "preservation": "strict", "render": "none"},
            "validation": {"contractVersion": 2, "min_slides": 2, "max_slides": 2, "required_text": ["Cloned decision"]},
        })

        candidate = Path(outcome["candidatePath"])
        self.assertFalse(destination.exists())
        cloned = Presentation(candidate)
        try:
            self.assertEqual(2, len(cloned.slides))
            source_title = next(shape for shape in cloned.slides[0].shapes if shape.shape_id == shape_id)
            clone_title = next(shape for shape in cloned.slides[1].shapes if shape.shape_id == shape_id)
            self.assertEqual("Original decision", source_title.text)
            self.assertEqual("Cloned decision", clone_title.text)
        finally:
            del cloned
        with zipfile.ZipFile(candidate) as archive:
            names = set(archive.namelist())
            chart_parts = sorted(name for name in names if re.fullmatch(r"ppt/charts/chart\d+\.xml", name))
            workbook_parts = sorted(name for name in names if name.startswith("ppt/embeddings/") and name.endswith(".xlsx"))
            self.assertEqual(2, len(chart_parts))
            self.assertEqual(2, len(workbook_parts))
            self.assertEqual(archive.read(chart_parts[0]), archive.read(chart_parts[1]))
            self.assertEqual(archive.read(workbook_parts[0]), archive.read(workbook_parts[1]))
            self.assertIn(b"transition", archive.read("ppt/slides/slide2.xml"))

    def test_pptx_insert_slide_preserves_existing_package_parts_under_strict_mode(self) -> None:
        try:
            from pptx import Presentation
        except ImportError:
            self.skipTest("python-pptx is not installed")

        source = self.root / "insert-source.pptx"
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = "Existing decision"
        slide.placeholders[1].text = "Existing evidence"
        presentation.save(source)
        with zipfile.ZipFile(source) as archive:
            source_hashes = {
                name: hashlib.sha256(archive.read(name)).hexdigest()
                for name in archive.namelist()
                if name.startswith(("ppt/slides/", "ppt/slideLayouts/", "ppt/slideMasters/", "ppt/theme/"))
            }

        outcome = self.engine.execute({
            "requestVersion": 2,
            "format": "pptx",
            "intent": "modify",
            "source": str(source),
            "preconditions": {"sourceSha256": hashlib.sha256(source.read_bytes()).hexdigest()},
            "destination": str(self.root / "insert-result.pptx"),
            "operations": [{
                "op": "insert_slide",
                "after": 1,
                "title": "Inserted decision",
                "body": "Inserted evidence",
            }],
            "guarantees": {"quality": "standard", "preservation": "strict", "render": "none"},
            "validation": {
                "contractVersion": 2,
                "min_slides": 2,
                "max_slides": 2,
                "required_text": ["Inserted decision", "Inserted evidence"],
            },
        })
        candidate = Path(outcome["candidatePath"])
        with zipfile.ZipFile(candidate) as archive:
            candidate_hashes = {
                name: hashlib.sha256(archive.read(name)).hexdigest()
                for name in source_hashes
            }
            self.assertIn("ppt/slides/slide2.xml", archive.namelist())
        self.assertEqual(source_hashes, candidate_hashes)

    def test_pptx_chart_data_is_atomic_across_embedded_workbook_and_cache(self) -> None:
        from pptx import Presentation
        from pptx.chart.data import ChartData
        from pptx.enum.chart import XL_CHART_TYPE
        from pptx.util import Inches

        source = self.root / "chart-deck.pptx"
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        data = ChartData()
        data.categories = ["North", "South"]
        data.add_series("Amount", (100, 80))
        chart_shape = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Inches(1), Inches(1), Inches(6), Inches(4), data,
        )
        shape_id = chart_shape.shape_id
        presentation.save(source)
        with zipfile.ZipFile(source) as archive:
            scripts = Path(__file__).resolve().parents[2] / "pptx-presentation-design" / "scripts"
            if str(scripts) not in sys.path:
                sys.path.insert(0, str(scripts))
            from pptx_structured_editor import presentation_order

            slide_id = presentation_order(archive)[0]["slideId"]
        destination = self.root / "chart-deck-result.pptx"

        outcome = self.engine.execute({
            "requestVersion": 2,
            "format": "pptx",
            "intent": "modify",
            "source": str(source),
            "preconditions": {"sourceSha256": hashlib.sha256(source.read_bytes()).hexdigest()},
            "destination": str(destination),
            "operations": [{
                "op": "set_chart_data",
                "slideId": slide_id,
                "shapeId": shape_id,
                "chartPart": "ppt/charts/chart1.xml",
                "seriesIndex": 1,
                "seriesName": "Updated amount",
                "categories": ["East", "West"],
                "values": [125, 95],
            }],
            "guarantees": {"quality": "standard", "preservation": "strict", "render": "none"},
        })

        self.assertFalse(destination.exists())
        self.assertTrue(outcome["preservationEvidence"]["verified"])
        inspection = self.engine.inspect(outcome["candidatePath"], "pptx")
        self.assertEqual([], inspection["profile"]["chart_validation_errors"])

    def test_pptx_notes_comments_and_alt_text_use_strict_typed_operations(self) -> None:
        from pptx import Presentation

        source = self.root / "review-deck.pptx"
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = "Decision"
        slide.placeholders[1].text = "Evidence"
        slide.notes_slide.notes_text_frame.text = "Original notes"
        shape_id = slide.shapes.title.shape_id
        presentation.save(source)
        with zipfile.ZipFile(source) as archive:
            scripts = Path(__file__).resolve().parents[2] / "pptx-presentation-design" / "scripts"
            if str(scripts) not in sys.path:
                sys.path.insert(0, str(scripts))
            from pptx_structured_editor import presentation_order

            slide_id = presentation_order(archive)[0]["slideId"]
        outcome = self.engine.execute({
            "requestVersion": 2,
            "format": "pptx",
            "intent": "modify",
            "source": str(source),
            "preconditions": {"sourceSha256": hashlib.sha256(source.read_bytes()).hexdigest()},
            "destination": str(self.root / "review-deck-result.pptx"),
            "operations": [
                {"op": "set_alt_text", "slideId": slide_id, "shapeId": shape_id, "altText": "Decision title"},
                {"op": "set_speaker_notes", "slideId": slide_id, "text": "Updated evidence notes"},
                {"op": "add_comment", "slideId": slide_id, "comment": "Confirm owner", "author": "Reviewer"},
            ],
            "guarantees": {"quality": "standard", "preservation": "strict", "render": "none"},
            "validation": {"contractVersion": 2, "require_speaker_notes": True},
        })
        self.assertTrue(outcome["preservationEvidence"]["verified"])
        with zipfile.ZipFile(Path(outcome["candidatePath"])) as archive:
            self.assertIn(b"Updated evidence notes", archive.read("ppt/notesSlides/notesSlide1.xml"))
            self.assertIn(b"Decision title", archive.read("ppt/slides/slide1.xml"))
            self.assertTrue(any(name.startswith("ppt/comments/comment") for name in archive.namelist()))
        inspection = self.engine.inspect(outcome["candidatePath"], "pptx")
        self.assertEqual("Confirm owner", inspection["profile"]["comments"][0]["text"])
        title_shape = next(
            shape
            for shape in inspection["profile"]["slide_details"][0]["shape_details"]
            if str(shape["shapeId"]) == str(shape_id)
        )
        self.assertEqual("Decision title", title_shape["altText"])

    def test_pptxgenjs_author_is_capability_solved_and_candidate_gated(self) -> None:
        import shutil
        from pptx import Presentation

        if not shutil.which("node"):
            self.skipTest("Node.js is not installed")
        status = {
            item["id"]: item for item in self.engine.capabilities()["backends"]
        }["pptxgenjs"]
        if status["status"] != "ready":
            self.skipTest(status.get("detail", "PptxGenJS is unavailable"))
        spec = self.root / "pptxgenjs-spec.json"
        spec.write_text(json.dumps({
            "schemaVersion": 1,
            "slides": [{
                "notes": ["Native author evidence"],
                "elements": [
                    {"type": "text", "text": "PptxGenJS decision", "x": 0.7, "y": 0.6, "w": 6.0, "h": 0.6, "options": {"fontSize": 26, "bold": True}},
                    {"type": "chart", "chartType": "column", "altText": "Decision chart", "data": [{"name": "Value", "labels": ["A", "B"], "values": [1, 2]}], "x": 0.7, "y": 1.5, "w": 7.0, "h": 4.0, "options": {}},
                ],
            }],
        }), encoding="utf-8")
        request = {
            "requestVersion": 2,
            "format": "pptx",
            "intent": "create",
            "destination": str(self.root / "pptxgenjs-result.pptx"),
            "operations": [{"op": "create", "spec": str(spec), "authorEngine": "pptxgenjs"}],
            "guarantees": {"quality": "standard", "preservation": "replace", "render": "none"},
        }
        assessment = self.engine.assess(request)
        self.assertTrue(assessment["ready"])
        self.assertIn("pptxgenjs", assessment["adapterPlan"]["requiredAdapters"])
        outcome = self.engine.execute(request)
        self.assertFalse((self.root / "pptxgenjs-result.pptx").exists())
        presentation = Presentation(Path(outcome["candidatePath"]))
        self.assertIn(
            "PptxGenJS decision",
            " ".join(shape.text for shape in presentation.slides[0].shapes if hasattr(shape, "text")),
        )

    def test_docx_review_operations_are_candidate_gated_and_contract_checked(self) -> None:
        import docx

        source = self.root / "review-source.docx"
        destination = self.root / "review-result.docx"
        document = docx.Document()
        document.add_paragraph("Approve old wording")
        document.save(source)

        outcome = self.engine.execute({
            "requestVersion": 2,
            "format": "docx",
            "intent": "modify",
            "source": str(source),
            "preconditions": {"sourceSha256": hashlib.sha256(source.read_bytes()).hexdigest()},
            "destination": str(destination),
            "operations": [
                {"op": "add_comment", "find": "Approve", "comment": "Owner confirmation required."},
                {"op": "reply_comment", "commentId": "0", "comment": "Owner confirmed."},
                {"op": "resolve_comment", "commentId": "0", "resolved": True},
                {"op": "tracked_replace", "find": "old", "replace": "new", "author": "Nexa"},
            ],
            "guarantees": {"quality": "standard", "preservation": "strict", "render": "none"},
            "validation": {"contractVersion": 2, "min_comments": 2, "require_tracked_changes": True},
        })

        self.assertEqual("candidate", outcome["status"])
        self.assertFalse(destination.exists())
        with zipfile.ZipFile(Path(outcome["candidatePath"])) as archive:
            xml = archive.read("word/document.xml")
            self.assertIn(b"commentRangeStart", xml)
            self.assertIn(b":ins", xml)
            self.assertIn(b":del", xml)
        inspection = self.engine.inspect(outcome["candidatePath"], "docx")
        self.assertEqual(2, len(inspection["profile"]["comments"]))
        self.assertEqual("0", inspection["profile"]["comments"][1]["parentId"])
        self.assertTrue(all(item["resolved"] for item in inspection["profile"]["comments"]))
        self.assertEqual({"del": 1, "ins": 1}, inspection["profile"]["trackedChanges"]["counts"])
        state_text = (Path(outcome["candidatePath"]).parent / "state.json").read_text(
            encoding="utf-8"
        )
        self.assertNotIn("Owner confirmation required.", state_text)
        self.assertNotIn('"find": "Approve"', state_text)
        self.assertNotIn('"replace": "new"', state_text)
        self.assertIn("requestSha256", state_text)

    def test_docx_fields_bookmarks_content_controls_and_protection_are_candidate_gated(self) -> None:
        import docx

        source = self.root / "structured-source.docx"
        document = docx.Document()
        document.add_paragraph("Decision anchor")
        document.add_paragraph("Reference placeholder")
        document.add_paragraph("Controlled value")
        document.save(source)
        outcome = self.engine.execute({
            "requestVersion": 2,
            "format": "docx",
            "intent": "modify",
            "source": str(source),
            "preconditions": {"sourceSha256": hashlib.sha256(source.read_bytes()).hexdigest()},
            "destination": str(self.root / "structured-result.docx"),
            "operations": [
                {"op": "add_bookmark", "find": "Decision", "bookmarkName": "DecisionAnchor"},
                {"op": "insert_field", "find": "placeholder", "instruction": "REF DecisionAnchor", "displayText": "Decision"},
                {"op": "wrap_content_control", "find": "Controlled", "tag": "decision", "lock": "content"},
                {
                    "op": "bind_template",
                    "bindings": {
                        "DecisionAnchor": "Bound decision",
                        "decision": "Bound controlled",
                    },
                },
                {"op": "set_protection", "mode": "trackedChanges"},
            ],
            "guarantees": {"quality": "standard", "preservation": "strict", "render": "none"},
        })
        self.assertTrue(outcome["preservationEvidence"]["verified"])
        with zipfile.ZipFile(Path(outcome["candidatePath"])) as archive:
            document_xml = archive.read("word/document.xml")
            settings_xml = archive.read("word/settings.xml")
            self.assertIn(b"bookmarkStart", document_xml)
            self.assertIn(b"fldSimple", document_xml)
            self.assertIn(b"sdtContent", document_xml)
            self.assertIn(b"documentProtection", settings_xml)
        inspection = self.engine.inspect(outcome["candidatePath"], "docx")["profile"]
        self.assertEqual("DecisionAnchor", inspection["bookmarks"][0]["name"])
        self.assertEqual("REF DecisionAnchor", inspection["fields"][0]["instruction"])
        self.assertEqual("decision", inspection["contentControls"][0]["tag"])
        self.assertEqual("Bound controlled", inspection["contentControls"][0]["text"])
        self.assertEqual("trackedChanges", inspection["protection"]["edit"])


if __name__ == "__main__":
    unittest.main()
