#!/usr/bin/env python3
"""Tests for the upgraded PPTX pipeline helpers."""

from __future__ import annotations

import json
import tempfile
import unittest
import zipfile
from xml.etree import ElementTree as ET
from pathlib import Path

import pptx_asset_pack
import pptx_deck_planner
import pptx_delivery_pack
import pptx_audit
import pptx_regression_suite
import pptx_rewrite_plan
import pptx_semantic_rewriter
import pptx_style_profile
import pptx_template_bind
import pptx_visual_qa


PRESENTATION_XML = """<?xml version="1.0" encoding="UTF-8"?>
<p:presentation xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
  <p:sldSz cx="12192000" cy="6858000"/>
</p:presentation>
"""


def _shape_xml(x: int, y: int, cx: int, cy: int, text: str, fill: str = "FFFFFF", color: str = "111111") -> str:
    return f"""
    <p:sp>
      <p:spPr>
        <a:xfrm>
          <a:off x="{x}" y="{y}"/>
          <a:ext cx="{cx}" cy="{cy}"/>
        </a:xfrm>
        <a:solidFill><a:srgbClr val="{fill}"/></a:solidFill>
      </p:spPr>
      <p:txBody>
        <a:bodyPr/>
        <a:lstStyle/>
        <a:p><a:r><a:rPr><a:solidFill><a:srgbClr val="{color}"/></a:solidFill></a:rPr><a:t>{text}</a:t></a:r></a:p>
      </p:txBody>
    </p:sp>
    """


def _write_minimal_pptx(path: Path, *, overlapping: bool = False) -> None:
    second_x = 960000 if overlapping else 4600000
    slide_xml = f"""<?xml version="1.0" encoding="UTF-8"?>
<p:sld xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
       xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:cSld><p:spTree>
    {_shape_xml(914400, 914400, 2743200, 914400, "First message")}
    {_shape_xml(second_x, 960000, 2743200, 914400, "Second message")}
  </p:spTree></p:cSld>
</p:sld>
"""
    with zipfile.ZipFile(path, "w") as zf:
        zf.writestr("ppt/presentation.xml", PRESENTATION_XML)
        zf.writestr("ppt/slides/slide1.xml", slide_xml)
        zf.writestr("ppt/media/image1.png", b"fake")
        zf.writestr(
            "ppt/slides/_rels/slide1.xml.rels",
            """<?xml version="1.0" encoding="UTF-8"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink" Target="https://example.com/source" TargetMode="External"/>
</Relationships>
""",
        )


def _write_theme_pptx(path: Path) -> None:
    theme_xml = """<?xml version="1.0" encoding="UTF-8"?>
<a:theme xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <a:themeElements>
    <a:clrScheme name="Custom">
      <a:dk1><a:srgbClr val="101010"/></a:dk1>
      <a:lt1><a:srgbClr val="FFFFFF"/></a:lt1>
      <a:accent1><a:srgbClr val="AA0000"/></a:accent1>
      <a:accent2><a:srgbClr val="00AA00"/></a:accent2>
    </a:clrScheme>
    <a:fontScheme name="Fonts">
      <a:majorFont><a:latin typeface="Georgia"/></a:majorFont>
      <a:minorFont><a:latin typeface="Arial"/></a:minorFont>
    </a:fontScheme>
  </a:themeElements>
</a:theme>
"""
    with zipfile.ZipFile(path, "w") as zf:
        zf.writestr("ppt/presentation.xml", PRESENTATION_XML)
        zf.writestr("ppt/theme/theme1.xml", theme_xml)
        zf.writestr("ppt/slides/slide1.xml", "<p:sld xmlns:p=\"http://schemas.openxmlformats.org/presentationml/2006/main\"/>")


class PptxUpgradePipelineTests(unittest.TestCase):
    def test_audit_uses_presentation_display_order_not_part_number(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            deck = Path(tmp) / "reordered.pptx"
            presentation = """<?xml version="1.0" encoding="UTF-8"?>
<p:presentation xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
 xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
 <p:sldIdLst>
  <p:sldId id="258" r:id="rId3"/><p:sldId id="256" r:id="rId1"/><p:sldId id="257" r:id="rId2"/>
 </p:sldIdLst><p:sldSz cx="12192000" cy="6858000"/>
</p:presentation>"""
            relationships = """<?xml version="1.0" encoding="UTF-8"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
 <Relationship Id="rId1" Type="x/slide" Target="slides/slide1.xml"/>
 <Relationship Id="rId2" Type="x/slide" Target="slides/slide2.xml"/>
 <Relationship Id="rId3" Type="x/slide" Target="slides/slide3.xml"/>
</Relationships>"""
            with zipfile.ZipFile(deck, "w") as archive:
                archive.writestr("ppt/presentation.xml", presentation)
                archive.writestr("ppt/_rels/presentation.xml.rels", relationships)
                for number, label in enumerate(("A", "B", "C"), start=1):
                    archive.writestr(
                        f"ppt/slides/slide{number}.xml",
                        f'<p:sld xmlns:p="{pptx_audit.NS["p"]}" xmlns:a="{pptx_audit.NS["a"]}"><p:cSld><p:spTree><a:t>{label}</a:t></p:spTree></p:cSld></p:sld>',
                    )

            report = pptx_audit.audit(deck)
            self.assertEqual(["C", "A", "B"], [slide["text"] for slide in report["slide_details"]])
            self.assertEqual(["258", "256", "257"], [slide["slide_id"] for slide in report["slide_details"]])
            self.assertEqual([], report["validation_errors"])

    def test_audit_rejects_chart_cache_dimension_mismatch(self) -> None:
        try:
            from pptx import Presentation
            from pptx.chart.data import ChartData
            from pptx.enum.chart import XL_CHART_TYPE
            from pptx.util import Inches
        except ImportError:
            self.skipTest("python-pptx is not installed")
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source = root / "chart.pptx"
            broken = root / "chart-broken.pptx"
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[5])
            data = ChartData()
            data.categories = ["A", "B"]
            data.add_series("Value", (1, 2))
            slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED,
                Inches(1), Inches(1), Inches(5), Inches(3), data,
            )
            presentation.save(source)
            with zipfile.ZipFile(source) as input_archive, zipfile.ZipFile(broken, "w") as output_archive:
                for info in input_archive.infolist():
                    payload = input_archive.read(info.filename)
                    if info.filename == "ppt/charts/chart1.xml":
                        chart = ET.fromstring(payload)
                        points = chart.findall(f".//{{{pptx_audit.NS['c']}}}val//{{{pptx_audit.NS['c']}}}pt")
                        parent = next(
                            element for element in chart.iter()
                            if points[0] in list(element)
                        )
                        parent.remove(points[0])
                        payload = ET.tostring(chart, encoding="utf-8", xml_declaration=True)
                    output_archive.writestr(info, payload)

            report = pptx_audit.audit(broken)
            self.assertTrue(any("dimension mismatch" in error for error in report["chart_validation_errors"]))

    def test_visual_qa_detects_overlap_and_repairs_dense_spec(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            deck = Path(tmp) / "deck.pptx"
            _write_minimal_pptx(deck, overlapping=True)

            report = pptx_visual_qa.analyze_pptx(deck)
            repaired = pptx_visual_qa.repair_spec(
                {"slides": [{"layout": "body", "title": "Dense", "bullets": list("ABCDEFG")}]}
            )

            self.assertEqual("fail", report["status"])
            self.assertTrue(any(issue["code"] == "shape_overlap" for issue in report["issues"]))
            self.assertEqual(2, len(repaired["slides"]))

    def test_visual_qa_scores_spec_design_and_flags_flat_renders(self) -> None:
        try:
            from PIL import Image
        except ImportError:
            self.skipTest("Pillow is not installed")

        spec = {
            "metadata": {"design_brief": {"industry": "finance"}},
            "slides": [
                {"layout": "title", "title": "Market Risk", "background_style": "spotlight", "design_role": "anchor"},
                {"layout": "body", "title": "Drivers", "bullets": ["One", "Two"], "icon": "shield", "design_role": "dense"},
                {"layout": "chart", "title": "Exposure", "categories": ["A"], "series": [{"name": "Value", "values": [1]}], "design_role": "breathing"},
            ],
        }
        review = pptx_visual_qa.evaluate_spec_design(spec)

        self.assertEqual("pass", review["status"])
        self.assertGreaterEqual(review["score"], 80)
        self.assertEqual(3, review["metrics"]["slides"])

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            deck = root / "deck.pptx"
            render_dir = root / "renders"
            render_dir.mkdir()
            _write_minimal_pptx(deck, overlapping=False)
            Image.new("RGB", (240, 135), (255, 255, 255)).save(render_dir / "slide-01.png")

            report = pptx_visual_qa.analyze_pptx(deck, render_dir)

            self.assertTrue(any(issue["code"] == "flat_render" for issue in report["issues"]))

    def test_style_profile_extracts_renderer_theme(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            deck = Path(tmp) / "theme.pptx"
            _write_theme_pptx(deck)

            profile = pptx_style_profile.profile_style(deck)

            self.assertEqual("AA0000", profile["renderer_theme"]["primary_color"])
            self.assertEqual("Georgia", profile["renderer_theme"]["title_font"])

    def test_deck_planner_preserves_links_and_selects_rich_layouts(self) -> None:
        text = """Growth Plan
Revenue increased 12%
Margin improved 18%
Roadmap Q1 launch
Source https://example.com/report
"""
        spec = pptx_deck_planner.plan_deck(text, audience="leadership", target_slides=5)

        self.assertEqual("pptx_deck_planner", spec["metadata"]["source"])
        self.assertIn("https://example.com/report", spec["metadata"]["source_links"])
        self.assertIn("visual_strategy", spec["metadata"])
        self.assertIn("design_brief", spec["metadata"])
        self.assertEqual("leadership", spec["metadata"]["design_brief"]["audience"])
        self.assertEqual(8, len(spec["metadata"]["design_brief"]["decision_points"]))
        self.assertIn(
            spec["theme"],
            {
                "consulting-clean",
                "executive-midnight",
                "product-energy",
                "editorial-ink",
                "nexa-dark",
                "nexa-light",
                "healthcare-trust",
                "finance-precision",
                "education-bright",
                "industrial-contrast",
            },
        )
        self.assertTrue(all(slide.get("design_role") in {"anchor", "breathing", "dense"} for slide in spec["slides"]))
        self.assertTrue(all(slide.get("background_style") for slide in spec["slides"]))
        self.assertTrue(any(slide["layout"] in {"chart", "timeline"} for slide in spec["slides"]))

    def test_deck_planner_infers_industry_visual_language(self) -> None:
        text = """Patient Access Plan
Healthcare teams need safer intake workflows
Patient wait time decreased 15%
Clinical quality should improve without more manual work
"""
        spec = pptx_deck_planner.plan_deck(text, audience="hospital executives", target_slides=5)
        brief = spec["metadata"]["design_brief"]

        self.assertEqual("healthcare", brief["industry"])
        self.assertEqual("healthcare-trust", spec["theme"])
        self.assertIn("clinical_grid", brief["visual_language"]["background_presets"])
        self.assertIn("trust", brief["visual_language"]["tone"])
        self.assertTrue(any(slide.get("icon") for slide in spec["slides"]))

    def test_deck_planner_preserves_years_and_does_not_mix_untyped_numeric_units_into_chart(self) -> None:
        text = """2026 年经营计划
营收增长 18%
预算为 ¥200 万
客户数量 35 家
路线图：第一季度完成试点
"""
        spec = pptx_deck_planner.plan_deck(text, audience="管理层", target_slides=5)

        self.assertEqual("2026 年经营计划", spec["slides"][0]["title"])
        self.assertEqual("disabled_without_typed_data", spec["metadata"]["chart_inference"])
        self.assertFalse(any(slide["layout"] == "chart" for slide in spec["slides"]))
        self.assertTrue(any(slide["layout"] == "timeline" for slide in spec["slides"]))

    def test_asset_pack_inventories_media_links_and_missing_spec_assets(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            deck = root / "deck.pptx"
            _write_minimal_pptx(deck)
            spec = root / "spec.json"
            spec.write_text(json.dumps({"slides": [{"image_path": "missing.png", "links": ["https://example.com"]}]}), encoding="utf-8")

            pptx_assets = pptx_asset_pack.inventory_pptx_assets(deck)
            spec_assets = pptx_asset_pack.validate_spec_assets(spec, root)

            self.assertEqual(1, pptx_assets["media_count"])
            self.assertEqual(1, pptx_assets["external_link_count"])
            self.assertEqual("fail", spec_assets["status"])

    def test_asset_pack_resolves_renderer_image_catalog_aliases(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            image = root / "hero.png"
            image.write_bytes(b"fake image")
            spec = root / "spec.json"
            spec.write_text(
                json.dumps(
                    {
                        "images": {"hero": "hero.png"},
                        "slides": [
                            {"layout": "title", "background_image_id": "hero"},
                            {"layout": "body", "image": "@hero"},
                        ],
                    }
                ),
                encoding="utf-8",
            )

            spec_assets = pptx_asset_pack.validate_spec_assets(spec, root)

            self.assertEqual("pass", spec_assets["status"])
            self.assertEqual(1, len(spec_assets["local_assets"]))

    def test_asset_pack_reports_image_semantics_for_catalog(self) -> None:
        try:
            from PIL import Image
        except ImportError:
            self.skipTest("Pillow is not installed")

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            hero = root / "hero.png"
            portrait = root / "portrait.png"
            Image.new("RGB", (1600, 900), (24, 48, 92)).save(hero)
            Image.new("RGB", (600, 1200), (220, 120, 40)).save(portrait)
            spec = root / "spec.json"
            spec.write_text(
                json.dumps(
                    {
                        "images": {
                            "hero": {"path": "hero.png", "role": "background"},
                            "speaker": {"path": "portrait.png"},
                        },
                        "slides": [
                            {"layout": "title", "background_image_id": "hero"},
                            {"layout": "body", "image_id": "speaker"},
                        ],
                    }
                ),
                encoding="utf-8",
            )

            spec_assets = pptx_asset_pack.validate_spec_assets(spec, root)
            by_alias = {item["alias"]: item for item in spec_assets["image_catalog"]}

            self.assertEqual("landscape", by_alias["hero"]["orientation"])
            self.assertEqual("background", by_alias["hero"]["recommended_usage"])
            self.assertEqual("portrait", by_alias["speaker"]["orientation"])
            self.assertEqual("inline_portrait", by_alias["speaker"]["recommended_usage"])

    def test_rewrite_plan_and_regression_samples_are_renderer_ready(self) -> None:
        report = {
            "warnings": ["slide 2 has no visual anchor"],
            "slide_details": [
                {"index": 1, "text": "Title", "has_visual_anchor": True, "text_chars": 20},
                {"index": 2, "text": "A long operating model summary", "has_visual_anchor": False, "text_chars": 900},
            ],
        }

        plan = pptx_rewrite_plan.build_rewrite_plan(report, target_slides=4)
        samples = pptx_regression_suite.sample_specs()

        self.assertEqual("rewrite-recommended", plan["status"])
        self.assertIn("recommended_spec", plan)
        self.assertIn("data_dashboard", samples)
        self.assertEqual("pptx_semantic_rewriter", plan["recommended_spec"]["metadata"]["source"])

    def test_template_bind_adds_layout_indices_and_style_tokens(self) -> None:
        spec = {"slides": [{"layout": "title", "title": "Quarterly Review"}, {"layout": "chart", "title": "Growth", "categories": ["Q1"], "series": [{"name": "Value", "values": [1]}]}]}
        template_profile = {
            "path": "template.pptx",
            "layouts": 2,
            "recommendations": {
                "title": {"layout_index": 0, "layout_name": "Title Slide", "score": 12},
                "chart": {"layout_index": 1, "layout_name": "Chart", "score": 10},
                "body": {"layout_index": 1, "layout_name": "Body", "score": 8},
            },
        }
        style_profile = {"renderer_theme": {"primary_color": "AA0000", "body_font": "Arial"}}

        bound = pptx_template_bind.bind_spec_to_template(spec, template_profile, style_profile=style_profile)

        self.assertEqual(0, bound["slides"][0]["template_layout_index"])
        self.assertEqual(1, bound["slides"][1]["template_layout_index"])
        self.assertEqual("AA0000", bound["theme"]["primary_color"])

    def test_semantic_rewriter_builds_decision_story(self) -> None:
        report = {
            "slides": 4,
            "slide_details": [
                {"index": 1, "text": "Market context and current customer demand are shifting."},
                {"index": 2, "text": "Revenue increased 12% and retention improved 18%."},
                {"index": 3, "text": "Option A reduces risk versus Option B but creates a tradeoff."},
                {"index": 4, "text": "Recommendation: prioritize Q2 launch and confirm next milestone."},
            ],
        }

        spec = pptx_semantic_rewriter.semantic_rewrite_from_report(report, target_slides=6)

        layouts = {slide["layout"] for slide in spec["slides"]}
        self.assertEqual("pptx_semantic_rewriter", spec["metadata"]["source"])
        self.assertIn("chart", layouts)
        self.assertIn("comparison", layouts)
        self.assertIn("timeline", layouts)

    def test_delivery_pack_writes_manifest(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            deck = root / "deck.pptx"
            out_dir = root / "delivery"
            _write_minimal_pptx(deck, overlapping=False)

            manifest = pptx_delivery_pack.create_delivery_pack(deck, out_dir)

            self.assertTrue((out_dir / "manifest.json").exists())
            self.assertEqual(str(deck), manifest["source"])


if __name__ == "__main__":
    unittest.main()
