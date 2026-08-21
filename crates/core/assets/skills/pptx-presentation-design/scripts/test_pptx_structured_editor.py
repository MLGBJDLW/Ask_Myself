from __future__ import annotations

import tempfile
import unittest
import zipfile
from pathlib import Path
from xml.etree import ElementTree as ET

import pptx_structured_editor
import pptx_audit


class PptxStructuredEditorTests(unittest.TestCase):
    def _source(self, root: Path) -> tuple[Path, str, int]:
        from pptx import Presentation

        path = root / "source.pptx"
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = "Decision"
        slide.placeholders[1].text = "Evidence"
        slide.notes_slide.notes_text_frame.text = "Original notes"
        title_shape_id = slide.shapes.title.shape_id
        presentation.save(path)
        with zipfile.ZipFile(path) as archive:
            stable_id = pptx_structured_editor.presentation_order(archive)[0]["slideId"]
        return path, stable_id, title_shape_id

    def test_notes_comments_and_alt_text_are_real_powerpoint_parts(self) -> None:
        from pptx import Presentation

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source, slide_id, shape_id = self._source(root)
            output = root / "edited.pptx"
            result = pptx_structured_editor.patch_pptx(source, output, [
                {"op": "set_alt_text", "slideId": slide_id, "shapeId": shape_id, "altText": "Decision title", "title": "Decision"},
                {"op": "set_speaker_notes", "slideId": slide_id, "text": "Updated presenter evidence"},
                {"op": "add_comment", "slideId": slide_id, "comment": "Confirm the decision owner", "author": "Reviewer", "initials": "RV"},
            ])

            self.assertIn("ppt/commentAuthors.xml", result["changedParts"])
            self.assertTrue(any(part.startswith("ppt/comments/comment") for part in result["changedParts"]))
            presentation = Presentation(output)
            self.assertEqual("Updated presenter evidence", presentation.slides[0].notes_slide.notes_text_frame.text)
            self.assertEqual(
                "Decision title",
                presentation.slides[0].shapes.title._element.nvSpPr.cNvPr.attrib["descr"],
            )
            with zipfile.ZipFile(output) as archive:
                authors = archive.read("ppt/commentAuthors.xml")
                comment_part = next(name for name in archive.namelist() if name.startswith("ppt/comments/comment"))
                comments = archive.read(comment_part)
                self.assertIn(b"Reviewer", authors)
                self.assertIn(b"Confirm the decision owner", comments)

    def test_speaker_notes_fails_closed_without_notes_relationship(self) -> None:
        from pptx import Presentation

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source = root / "without-notes.pptx"
            presentation = Presentation()
            presentation.slides.add_slide(presentation.slide_layouts[6])
            presentation.save(source)
            with zipfile.ZipFile(source) as archive:
                slide_id = pptx_structured_editor.presentation_order(archive)[0]["slideId"]
            # python-pptx may create a notes part lazily only when requested; a
            # package without it must not silently fabricate an invalid master graph.
            with zipfile.ZipFile(source) as archive:
                has_notes = any(name.startswith("ppt/notesSlides/notesSlide") for name in archive.namelist())
            if has_notes:
                self.skipTest("fixture writer eagerly created a notes relationship")
            with self.assertRaisesRegex(pptx_structured_editor.PptxEditError, "existing notes"):
                pptx_structured_editor.patch_pptx(source, root / "failed.pptx", [{
                    "op": "set_speaker_notes", "slideId": slide_id, "text": "Notes",
                }])

    def test_insert_slide_is_lossless_with_notes_comments_and_macro_templates(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source, slide_id, _ = self._source(root)
            reviewed = root / "reviewed.pptx"
            pptx_structured_editor.patch_pptx(source, reviewed, [{
                "op": "add_comment",
                "slideId": slide_id,
                "comment": "Keep this review thread",
                "author": "Reviewer",
            }])
            with zipfile.ZipFile(reviewed) as archive:
                original_slide_rels = archive.read("ppt/slides/_rels/slide1.xml.rels")
                original_comment = next(
                    archive.read(name)
                    for name in archive.namelist()
                    if name.startswith("ppt/comments/comment")
                )
            inserted = root / "inserted.pptx"
            result = pptx_structured_editor.patch_pptx(reviewed, inserted, [{
                "op": "insert_slide", "after": 1, "title": "Inserted", "body": "Evidence",
            }])
            self.assertIn("ppt/slides/slide2.xml", result["changedParts"])
            with zipfile.ZipFile(inserted) as archive:
                self.assertEqual(original_slide_rels, archive.read("ppt/slides/_rels/slide1.xml.rels"))
                self.assertEqual(
                    original_comment,
                    next(
                        archive.read(name)
                        for name in archive.namelist()
                        if name.startswith("ppt/comments/comment")
                    ),
                )
                self.assertIn(b"Inserted", archive.read("ppt/slides/slide2.xml"))
                self.assertIn(b"Evidence", archive.read("ppt/slides/slide2.xml"))

            potm = root / "reviewed.potm"
            with zipfile.ZipFile(reviewed) as archive, zipfile.ZipFile(potm, "w") as output:
                for info in archive.infolist():
                    data = archive.read(info.filename)
                    if info.filename == "[Content_Types].xml":
                        data = data.replace(
                            b"application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml",
                            b"application/vnd.ms-powerpoint.template.macroEnabled.main+xml",
                        )
                    output.writestr(info, data)
            potm_inserted = root / "inserted.potm"
            pptx_structured_editor.patch_pptx(potm, potm_inserted, [{
                "op": "insert_slide", "after": 0, "title": "Macro-safe insert",
            }])
            with zipfile.ZipFile(potm_inserted) as archive:
                self.assertIn(b"Macro-safe insert", archive.read("ppt/slides/slide2.xml"))

    def test_clone_preserves_animation_and_copy_on_write_smartart_ole_closure(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source, slide_id, _ = self._source(root)
            injected = root / "advanced-source.pptx"
            rel_ns = "http://schemas.openxmlformats.org/package/2006/relationships"
            p_ns = "http://schemas.openxmlformats.org/presentationml/2006/main"
            ct_ns = "http://schemas.openxmlformats.org/package/2006/content-types"
            additions = {
                "ppt/diagrams/data1.xml": b'<?xml version="1.0"?><dgm:dataModel xmlns:dgm="http://schemas.openxmlformats.org/drawingml/2006/diagram"/>',
                "ppt/diagrams/layout1.xml": b'<?xml version="1.0"?><dgm:layoutDef xmlns:dgm="http://schemas.openxmlformats.org/drawingml/2006/diagram"/>',
                "ppt/diagrams/colors1.xml": b'<?xml version="1.0"?><dgm:colorsDef xmlns:dgm="http://schemas.openxmlformats.org/drawingml/2006/diagram"/>',
                "ppt/diagrams/quickStyle1.xml": b'<?xml version="1.0"?><dgm:styleDef xmlns:dgm="http://schemas.openxmlformats.org/drawingml/2006/diagram"/>',
                "ppt/embeddings/oleObject1.bin": b"OLE-COPY-ON-WRITE-EVIDENCE",
            }
            content_types = {
                "ppt/diagrams/data1.xml": "application/vnd.openxmlformats-officedocument.drawingml.diagramData+xml",
                "ppt/diagrams/layout1.xml": "application/vnd.openxmlformats-officedocument.drawingml.diagramLayout+xml",
                "ppt/diagrams/colors1.xml": "application/vnd.openxmlformats-officedocument.drawingml.diagramColors+xml",
                "ppt/diagrams/quickStyle1.xml": "application/vnd.openxmlformats-officedocument.drawingml.diagramStyle+xml",
                "ppt/embeddings/oleObject1.bin": "application/vnd.openxmlformats-officedocument.oleObject",
            }
            with zipfile.ZipFile(source) as archive, zipfile.ZipFile(injected, "w") as output:
                for info in archive.infolist():
                    data = archive.read(info.filename)
                    if info.filename == "ppt/slides/slide1.xml":
                        slide = ET.fromstring(data)
                        timing = ET.SubElement(slide, f"{{{p_ns}}}timing")
                        timing.set("advTm", "2500")
                        data = ET.tostring(slide, encoding="utf-8", xml_declaration=True)
                    elif info.filename == "ppt/slides/_rels/slide1.xml.rels":
                        relationships = ET.fromstring(data)
                        ET.SubElement(relationships, f"{{{rel_ns}}}Relationship", {
                            "Id": "rId900", "Type": f"{pptx_structured_editor.R_NS}/diagramData",
                            "Target": "../diagrams/data1.xml",
                        })
                        ET.SubElement(relationships, f"{{{rel_ns}}}Relationship", {
                            "Id": "rId901", "Type": f"{pptx_structured_editor.R_NS}/oleObject",
                            "Target": "../embeddings/oleObject1.bin",
                        })
                        data = ET.tostring(relationships, encoding="utf-8", xml_declaration=True)
                    elif info.filename == "[Content_Types].xml":
                        types = ET.fromstring(data)
                        for part, content_type in content_types.items():
                            ET.SubElement(types, f"{{{ct_ns}}}Override", {
                                "PartName": f"/{part}", "ContentType": content_type,
                            })
                        data = ET.tostring(types, encoding="utf-8", xml_declaration=True)
                    output.writestr(info, data)
                diagram_relationships = ET.Element(f"{{{rel_ns}}}Relationships")
                for index, (leaf, target) in enumerate((
                    ("diagramLayout", "layout1.xml"),
                    ("diagramColors", "colors1.xml"),
                    ("diagramQuickStyle", "quickStyle1.xml"),
                ), start=1):
                    ET.SubElement(diagram_relationships, f"{{{rel_ns}}}Relationship", {
                        "Id": f"rId{index}",
                        "Type": f"{pptx_structured_editor.R_NS}/{leaf}",
                        "Target": target,
                    })
                output.writestr(
                    "ppt/diagrams/_rels/data1.xml.rels",
                    ET.tostring(diagram_relationships, encoding="utf-8", xml_declaration=True),
                )
                for name, data in additions.items():
                    output.writestr(name, data)

            output_path = root / "advanced-clone.pptx"
            result = pptx_structured_editor.patch_pptx(
                injected,
                output_path,
                [{"op": "clone_slide", "slideId": slide_id}],
            )
            mapping = result["operations"][0]["detail"]["clonedParts"]
            for original in additions:
                self.assertIn(original, mapping)
                self.assertNotEqual(original, mapping[original])
            with zipfile.ZipFile(output_path) as archive:
                self.assertIn(b"timing", archive.read("ppt/slides/slide2.xml"))
                self.assertEqual(
                    archive.read("ppt/embeddings/oleObject1.bin"),
                    archive.read(mapping["ppt/embeddings/oleObject1.bin"]),
                )
                cloned_data = mapping["ppt/diagrams/data1.xml"]
                cloned_rels = ET.fromstring(archive.read(pptx_structured_editor._rels_path(cloned_data)))
                targets = {item.attrib["Target"] for item in cloned_rels}
                self.assertNotIn("layout1.xml", targets)
                self.assertTrue(any("layout" in target for target in targets))
                self.assertEqual(
                    additions["ppt/diagrams/data1.xml"],
                    archive.read("ppt/diagrams/data1.xml"),
                )

    def test_new_slides_are_available_to_later_operations_in_same_request(self) -> None:
        from pptx import Presentation

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source, _, _ = self._source(root)
            output = root / "sequential-slides.pptx"
            result = pptx_structured_editor.patch_pptx(source, output, [
                {
                    "op": "insert_slide",
                    "after": 1,
                    "title": "Inserted title",
                    "body": "Inserted body",
                },
                {
                    "op": "set_text",
                    "slideIndex": 2,
                    "shapeName": "Nexa title",
                    "text": "Edited inserted title",
                },
                {
                    "op": "add_comment",
                    "slideIndex": 2,
                    "comment": "Comment on staged slide",
                },
                {"op": "clone_slide", "slideIndex": 2},
                {
                    "op": "set_text",
                    "slideIndex": 3,
                    "shapeName": "Nexa title",
                    "text": "Edited clone title",
                },
                {
                    "op": "set_transition",
                    "slideIndex": 3,
                    "transition": "fade",
                },
            ])

            self.assertEqual(3, len(Presentation(output).slides))
            with zipfile.ZipFile(output) as archive:
                order = pptx_structured_editor.presentation_order(archive)
                second = archive.read(order[1]["part"])
                third = archive.read(order[2]["part"])
                has_comment_part = any(
                    name.startswith("ppt/comments/comment")
                    for name in archive.namelist()
                )
            self.assertIn(b"Edited inserted title", second)
            self.assertIn(b"Edited clone title", third)
            self.assertIn(b"transition", third)
            self.assertTrue(has_comment_part)
            self.assertEqual(6, len(result["operations"]))

    def test_chart_data_updates_embedded_workbook_and_caches_atomically(self) -> None:
        import io
        import openpyxl
        from pptx import Presentation
        from pptx.chart.data import ChartData
        from pptx.enum.chart import XL_CHART_TYPE
        from pptx.util import Inches

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source = root / "chart-source.pptx"
            output = root / "chart-output.pptx"
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[5])
            data = ChartData()
            data.categories = ["North", "South"]
            data.add_series("Amount", (100, 80))
            chart_shape = slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED,
                Inches(1), Inches(1), Inches(6), Inches(4), data,
            )
            shape_id = chart_shape.shape_id
            presentation.save(source)
            with zipfile.ZipFile(source) as archive:
                slide_id = pptx_structured_editor.presentation_order(archive)[0]["slideId"]

            result = pptx_structured_editor.patch_pptx(source, output, [{
                "op": "set_chart_data",
                "slideId": slide_id,
                "shapeId": shape_id,
                "chartPart": "ppt/charts/chart1.xml",
                "seriesIndex": 1,
                "seriesName": "Updated amount",
                "categoryRange": "Sheet1!$A$2:$A$4",
                "valueRange": "Sheet1!$B$2:$B$4",
                "categories": ["East", "West", "Central"],
                "values": [125, 95.5, 60],
            }])

            changed = set(result["changedParts"])
            workbook_part = result["operations"][0]["detail"]["workbookPart"]
            self.assertIn("ppt/charts/chart1.xml", changed)
            self.assertIn(workbook_part, changed)
            audit = pptx_audit.audit(output)
            self.assertEqual([], audit["chart_validation_errors"])
            self.assertIn("dataVisual", audit["slide_details"][0]["frame_map"])
            chart_frame = next(
                item
                for item in audit["slide_details"][0]["shape_details"]
                if str(item["shapeId"]) == str(shape_id)
            )
            self.assertEqual("dataVisual", chart_frame["semanticRole"])
            with zipfile.ZipFile(output) as archive:
                embedded = archive.read(workbook_part)
                chart = archive.read("ppt/charts/chart1.xml")
            workbook = openpyxl.load_workbook(io.BytesIO(embedded), data_only=False)
            sheet = workbook["Sheet1"]
            self.assertEqual(["East", "West", "Central"], [sheet[f"A{row}"].value for row in range(2, 5)])
            self.assertEqual([125, 95.5, 60], [sheet[f"B{row}"].value for row in range(2, 5)])
            self.assertEqual("Updated amount", sheet["B1"].value)
            workbook.close()
            self.assertIn(b"Sheet1!$A$2:$A$4", chart)
            self.assertIn(b"Sheet1!$B$2:$B$4", chart)

            chained = root / "chart-clone-output.pptx"
            chained_result = pptx_structured_editor.patch_pptx(source, chained, [
                {"op": "clone_slide", "slideId": slide_id},
                {
                    "op": "set_chart_data",
                    "slideIndex": 2,
                    "shapeId": shape_id,
                    "seriesIndex": 1,
                    "categoryRange": "Sheet1!$A$2:$A$3",
                    "valueRange": "Sheet1!$B$2:$B$3",
                    "categories": ["Clone A", "Clone B"],
                    "values": [11, 22],
                },
            ])
            chained_detail = chained_result["operations"][1]["detail"]
            self.assertNotEqual("ppt/charts/chart1.xml", chained_detail["chartPart"])
            with zipfile.ZipFile(chained) as archive:
                cloned_workbook = archive.read(chained_detail["workbookPart"])
            workbook = openpyxl.load_workbook(io.BytesIO(cloned_workbook), data_only=False)
            self.assertEqual(["Clone A", "Clone B"], [workbook["Sheet1"][cell].value for cell in ("A2", "A3")])
            self.assertEqual([11, 22], [workbook["Sheet1"][cell].value for cell in ("B2", "B3")])
            workbook.close()
            self.assertEqual([], pptx_audit.audit(chained)["chart_validation_errors"])


if __name__ == "__main__":
    unittest.main()
