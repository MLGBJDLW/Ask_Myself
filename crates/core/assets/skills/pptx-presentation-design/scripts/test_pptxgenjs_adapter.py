from __future__ import annotations

import json
import os
import shutil
import subprocess
import tempfile
import unittest
import zipfile
from pathlib import Path


class PptxGenJsAdapterTests(unittest.TestCase):
    def setUp(self) -> None:
        self.node = shutil.which("node")
        if not self.node:
            self.skipTest("Node.js is not installed")
        probe = subprocess.run(
            [self.node, "-e", "require.resolve('pptxgenjs')"],
            cwd=Path(__file__).resolve().parents[6],
            capture_output=True,
            text=True,
            check=False,
        )
        if probe.returncode:
            self.skipTest("version-locked pptxgenjs module is not installed")
        self.adapter = Path(__file__).with_name("pptxgenjs_adapter.mjs")

    def test_authors_editable_master_chart_table_notes_and_alt_text(self) -> None:
        from pptx import Presentation

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            spec = root / "deck.json"
            output = root / "deck.pptx"
            spec.write_text(json.dumps({
                "schemaVersion": 1,
                "title": "Native author adapter",
                "masters": [{
                    "title": "BRAND",
                    "background": {"color": "F5F7FA"},
                    "objects": [{"rect": {"x": 0, "y": 0, "w": 13.333, "h": 0.18, "fill": {"color": "2457D6"}, "line": {"color": "2457D6"}}}],
                    "slideNumber": {"x": 12.4, "y": 7.1, "w": 0.5, "h": 0.2},
                }],
                "slides": [{
                    "masterName": "BRAND",
                    "notes": ["Presenter evidence"],
                    "elements": [
                        {"type": "text", "text": "Decision", "x": 0.7, "y": 0.5, "w": 5.5, "h": 0.6, "options": {"fontSize": 26, "bold": True}},
                        {"type": "table", "rows": [["Metric", "Value"], ["Revenue", 100]], "x": 0.7, "y": 1.4, "w": 4.2, "h": 1.2, "options": {"border": {"type": "solid", "color": "C7D2E5", "pt": 1}}},
                        {"type": "chart", "chartType": "column", "altText": "Revenue chart", "data": [{"name": "Revenue", "labels": ["Q1", "Q2"], "values": [10, 20]}], "x": 5.3, "y": 1.4, "w": 6.5, "h": 3.8, "options": {"showLegend": False, "showTitle": True, "title": "Revenue"}},
                    ],
                }],
            }), encoding="utf-8")
            run = subprocess.run(
                [self.node, str(self.adapter), "--spec", str(spec), "--out", str(output), "--workspace", str(root)],
                cwd=Path(__file__).resolve().parents[6],
                capture_output=True,
                text=True,
                check=False,
            )
            self.assertEqual(0, run.returncode, run.stderr)
            result = json.loads(run.stdout)
            self.assertEqual("4.0.1", result["engineVersion"])
            presentation = Presentation(output)
            self.assertEqual(1, len(presentation.slides))
            self.assertIn("Decision", " ".join(shape.text for shape in presentation.slides[0].shapes if hasattr(shape, "text")))
            with zipfile.ZipFile(output) as archive:
                names = set(archive.namelist())
                self.assertIn("ppt/charts/chart1.xml", names)
                self.assertTrue(any(name.startswith("ppt/embeddings/") and name.endswith(".xlsx") for name in names))
                self.assertTrue(any(name.startswith("ppt/notesSlides/notesSlide") for name in names))
                self.assertTrue(any(name.startswith("ppt/slideMasters/slideMaster") for name in names))

    def test_blocks_urls_unc_and_vulnerable_image_parser_extensions(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            spec = root / "unsafe.json"
            output = root / "unsafe.pptx"
            spec.write_text(json.dumps({
                "schemaVersion": 1,
                "slides": [{"elements": [{
                    "type": "image", "path": "https://example.invalid/a.icns", "altText": "unsafe",
                    "x": 0, "y": 0, "w": 1, "h": 1,
                }]}],
            }), encoding="utf-8")
            run = subprocess.run(
                [self.node, str(self.adapter), "--spec", str(spec), "--out", str(output), "--workspace", str(root)],
                cwd=Path(__file__).resolve().parents[6],
                capture_output=True,
                text=True,
                check=False,
            )
            self.assertNotEqual(0, run.returncode)
            self.assertIn("URLs", run.stderr)
            self.assertFalse(output.exists())

    def test_blocks_workspace_symlink_assets_that_resolve_outside(self) -> None:
        with tempfile.TemporaryDirectory() as workspace_tmp, tempfile.TemporaryDirectory() as outside_tmp:
            root = Path(workspace_tmp)
            outside = Path(outside_tmp) / "outside.png"
            outside.write_bytes(b"not-read-because-link-is-blocked")
            linked = root / "linked.png"
            try:
                os.symlink(outside, linked)
            except (OSError, NotImplementedError) as error:
                self.skipTest(f"file symlinks unavailable: {error}")
            spec = root / "symlink.json"
            output = root / "symlink.pptx"
            spec.write_text(json.dumps({
                "schemaVersion": 1,
                "slides": [{"elements": [{
                    "type": "image", "path": linked.name, "altText": "blocked",
                    "x": 0, "y": 0, "w": 1, "h": 1,
                }]}],
            }), encoding="utf-8")
            run = subprocess.run(
                [self.node, str(self.adapter), "--spec", str(spec), "--out", str(output), "--workspace", str(root)],
                cwd=Path(__file__).resolve().parents[6],
                capture_output=True,
                text=True,
                check=False,
            )
            self.assertNotEqual(0, run.returncode)
            self.assertIn("symbolic link", run.stderr)
            self.assertFalse(output.exists())


if __name__ == "__main__":
    unittest.main()
