#!/usr/bin/env python3
"""PPTX renderer for the pptx-presentation-design skill.

Creates editable PowerPoint decks from a compact JSON spec. This module owns
presentation-specific rendering; doc-script-editor only keeps a compatibility
wrapper for the historical `create_pptx` command.
"""

from __future__ import annotations

import argparse
import json
import re
import sys
import tempfile
from pathlib import Path
from typing import Any
from urllib.parse import urlparse
from urllib.request import Request, urlopen


PPTX_THEME_PRESETS = {
    "nexa-light": {
        "primary_color": "2563EB",
        "accent_color": "0F766E",
        "background_color": "F8FAFC",
        "surface_color": "FFFFFF",
        "muted_surface_color": "E2E8F0",
        "text_color": "111827",
        "muted_text_color": "64748B",
        "title_color": "0F172A",
        "inverse_text_color": "FFFFFF",
        "title_font": "Aptos Display",
        "body_font": "Aptos",
    },
    "nexa-dark": {
        "primary_color": "38BDF8",
        "accent_color": "A3E635",
        "background_color": "0B1120",
        "surface_color": "111827",
        "muted_surface_color": "1F2937",
        "text_color": "E5E7EB",
        "muted_text_color": "94A3B8",
        "title_color": "F8FAFC",
        "inverse_text_color": "0B1120",
        "title_font": "Aptos Display",
        "body_font": "Aptos",
    },
    "consulting-clean": {
        "primary_color": "005587",
        "accent_color": "C41230",
        "background_color": "F7F8FA",
        "surface_color": "FFFFFF",
        "muted_surface_color": "E5E7EB",
        "text_color": "1F2937",
        "muted_text_color": "6B7280",
        "title_color": "111827",
        "inverse_text_color": "FFFFFF",
        "title_font": "Aptos Display",
        "body_font": "Aptos",
        "background_style": "diagonal",
    },
    "executive-midnight": {
        "primary_color": "93C5FD",
        "accent_color": "F59E0B",
        "background_color": "0F172A",
        "surface_color": "111827",
        "muted_surface_color": "1E293B",
        "text_color": "E5E7EB",
        "muted_text_color": "CBD5E1",
        "title_color": "F8FAFC",
        "inverse_text_color": "0F172A",
        "title_font": "Aptos Display",
        "body_font": "Aptos",
        "background_style": "gradient_mesh",
    },
    "editorial-ink": {
        "primary_color": "6D2E46",
        "accent_color": "0F766E",
        "background_color": "FBF7F2",
        "surface_color": "FFFFFF",
        "muted_surface_color": "E8DED4",
        "text_color": "2D2A26",
        "muted_text_color": "6B625C",
        "title_color": "231F20",
        "inverse_text_color": "FFFFFF",
        "title_font": "Georgia",
        "body_font": "Aptos",
        "background_style": "soft_geometry",
    },
    "product-energy": {
        "primary_color": "2563EB",
        "accent_color": "F97316",
        "background_color": "F8FAFC",
        "surface_color": "FFFFFF",
        "muted_surface_color": "DBEAFE",
        "text_color": "111827",
        "muted_text_color": "475569",
        "title_color": "0F172A",
        "inverse_text_color": "FFFFFF",
        "title_font": "Aptos Display",
        "body_font": "Aptos",
        "background_style": "gradient_mesh",
    },
    "healthcare-trust": {
        "primary_color": "0F766E",
        "accent_color": "2563EB",
        "background_color": "F6FBFA",
        "surface_color": "FFFFFF",
        "muted_surface_color": "DDF3EF",
        "text_color": "12332F",
        "muted_text_color": "4B6965",
        "title_color": "0B2724",
        "inverse_text_color": "FFFFFF",
        "title_font": "Aptos Display",
        "body_font": "Aptos",
        "background_style": "clinical_grid",
    },
    "finance-precision": {
        "primary_color": "1E3A8A",
        "accent_color": "16A34A",
        "background_color": "F8FAFC",
        "surface_color": "FFFFFF",
        "muted_surface_color": "E2E8F0",
        "text_color": "111827",
        "muted_text_color": "475569",
        "title_color": "0F172A",
        "inverse_text_color": "FFFFFF",
        "title_font": "Aptos Display",
        "body_font": "Aptos",
        "background_style": "data_grid",
    },
    "education-bright": {
        "primary_color": "7C3AED",
        "accent_color": "F59E0B",
        "background_color": "FFFBEB",
        "surface_color": "FFFFFF",
        "muted_surface_color": "FDE68A",
        "text_color": "312E24",
        "muted_text_color": "6B6250",
        "title_color": "1F2937",
        "inverse_text_color": "FFFFFF",
        "title_font": "Trebuchet MS",
        "body_font": "Aptos",
        "background_style": "paper_texture",
    },
    "industrial-contrast": {
        "primary_color": "374151",
        "accent_color": "F59E0B",
        "background_color": "F3F4F6",
        "surface_color": "FFFFFF",
        "muted_surface_color": "D1D5DB",
        "text_color": "111827",
        "muted_text_color": "4B5563",
        "title_color": "030712",
        "inverse_text_color": "FFFFFF",
        "title_font": "Aptos Display",
        "body_font": "Aptos",
        "background_style": "blueprint_grid",
    },
}

PPTX_SUPPORTED_LAYOUTS = {
    "title",
    "agenda",
    "body",
    "two_column",
    "stat",
    "quote",
    "section",
    "image_full",
    "table",
    "timeline",
    "process",
    "comparison",
    "matrix",
    "chart",
}

PPTX_SUPPORTED_CHART_TYPES = {
    "area",
    "bar",
    "bar_stacked",
    "column",
    "column_stacked",
    "doughnut",
    "line",
    "pie",
    "stacked_bar",
    "stacked_column",
}

PPTX_BACKGROUND_STYLES = {
    "none",
    "solid",
    "flat",
    "diagonal",
    "section",
    "editorial",
    "mesh",
    "gradient_mesh",
    "soft_geometry",
    "ambient",
    "blueprint_grid",
    "paper_texture",
    "clinical_grid",
    "data_grid",
    "spotlight",
}

PPTX_TRANSITION_TYPES = {
    "fade",
    "push",
    "wipe",
    "split",
    "cover",
    "pull",
    "cut",
}

EMU_PER_INCH = 914400
IMAGE_SUFFIXES = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tif", ".tiff"}


def _die(msg: str, code: int = 1) -> None:
    print(msg, file=sys.stderr)
    sys.exit(code)


def _missing(pkg: str) -> None:
    print(f"MISSING_DEP: {pkg}", file=sys.stderr)
    print(f"Install with: python -m pip install {pkg}", file=sys.stderr)
    sys.exit(2)


def _workspace_root(workspace_root: str | Path | None = None) -> Path:
    return Path(workspace_root).resolve() if workspace_root else Path.cwd().resolve()


def _validate_path(raw: str, *, must_exist: bool = True, workspace_root: str | Path | None = None) -> Path:
    if not raw:
        _die("ERROR: path is required", 3)
    p = Path(raw)
    if not p.is_absolute():
        _die(f"ERROR: path must be absolute: {raw}", 3)
    try:
        resolved = p.resolve()
        resolved.relative_to(_workspace_root(workspace_root))
    except ValueError:
        _die(f"ERROR: path escapes workspace: {raw}", 3)
    except OSError as exc:
        _die(f"ERROR: cannot resolve path: {exc}", 3)
    if must_exist and not resolved.exists():
        _die(f"ERROR: file not found: {resolved}", 3)
    return resolved


def _ext(p: Path) -> str:
    return p.suffix.lower().lstrip(".")


def _validate_output_path(raw: str, *, workspace_root: str | Path | None = None) -> Path:
    p = _validate_path(raw, must_exist=False, workspace_root=workspace_root)
    if _ext(p) != "pptx":
        _die("ERROR: output path must end with .pptx", 3)
    p.parent.mkdir(parents=True, exist_ok=True)
    return p


def _read_json(path: str, *, workspace_root: str | Path | None = None) -> dict[str, Any]:
    if path == "-":
        data = json.load(sys.stdin)
    else:
        spec_path = _validate_path(path, workspace_root=workspace_root)
        with spec_path.open("r", encoding="utf-8") as f:
            data = json.load(f)
    if not isinstance(data, dict):
        _die("ERROR: JSON spec root must be an object", 3)
    return data


def _hex_to_rgb(value: str):
    try:
        from pptx.dml.color import RGBColor  # type: ignore
    except ImportError:
        _missing("python-pptx")

    raw = str(value or "").strip().lstrip("#")
    if len(raw) == 3:
        raw = "".join(ch * 2 for ch in raw)
    if len(raw) != 6 or any(ch not in "0123456789abcdefABCDEF" for ch in raw):
        _die(f"ERROR: invalid PPTX theme color: {value}", 3)
    return RGBColor(int(raw[0:2], 16), int(raw[2:4], 16), int(raw[4:6], 16))


def _normalize_theme(theme_spec: Any) -> dict[str, str]:
    if theme_spec is None:
        return dict(PPTX_THEME_PRESETS["nexa-light"])
    if isinstance(theme_spec, str):
        key = theme_spec.strip().lower()
        key = {"light": "nexa-light", "dark": "nexa-dark"}.get(key, key)
        if key not in PPTX_THEME_PRESETS:
            presets = ", ".join(sorted(PPTX_THEME_PRESETS))
            _die(f"ERROR: unsupported PPTX theme. Use one of: {presets}; or pass a custom theme object.", 3)
        return dict(PPTX_THEME_PRESETS[key])
    if not isinstance(theme_spec, dict):
        _die("ERROR: PPTX theme must be a string or object", 3)

    theme = dict(PPTX_THEME_PRESETS["nexa-light"])
    for key, value in theme_spec.items():
        if value is not None:
            theme[str(key)] = str(value).strip().lstrip("#")
    return theme


def _rgb(theme: dict[str, str], key: str):
    return _hex_to_rgb(theme[key])


def _inches(value: float):
    try:
        from pptx.util import Inches  # type: ignore
    except ImportError:
        _missing("python-pptx")
    return Inches(float(value))


def _pt(value: float):
    try:
        from pptx.util import Pt  # type: ignore
    except ImportError:
        _missing("python-pptx")
    return Pt(float(value))


def _slide_size_inches(prs: Any) -> tuple[float, float]:
    return (float(prs.slide_width) / EMU_PER_INCH, float(prs.slide_height) / EMU_PER_INCH)


def _blank_layout(prs: Any):
    for layout in prs.slide_layouts:
        if "blank" in getattr(layout, "name", "").lower():
            return layout
    return prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]


def _remove_existing_slides(prs: Any) -> None:
    """Clear sample slides from a template while preserving masters/layouts."""
    try:
        slide_id_list = prs.slides._sldIdLst
        slide_ids = list(slide_id_list)
        for slide_id in slide_ids:
            prs.part.drop_rel(slide_id.rId)
            slide_id_list.remove(slide_id)
    except Exception as exc:
        _die(f"ERROR: could not clear template slides: {exc}", 3)


def _bound_template_layout(prs: Any, slide_spec: dict[str, Any]):
    raw = slide_spec.get("template_layout_index")
    if raw in (None, ""):
        raw = slide_spec.get("layout_index")
    if raw in (None, ""):
        return None
    try:
        index = int(raw)
    except (TypeError, ValueError):
        _die(f"ERROR: template_layout_index must be an integer: {raw}", 3)
    if index < 0 or index >= len(prs.slide_layouts):
        _die(f"ERROR: template_layout_index out of range: {index}", 3)
    return prs.slide_layouts[index]


def _set_slide_size(prs: Any, spec: dict[str, Any], has_template: bool) -> None:
    if has_template:
        return
    size = spec.get("slide_size") or spec.get("size") or "wide"
    if isinstance(size, str):
        key = size.strip().lower()
        if key in {"wide", "16:9", "ppt169", "widescreen"}:
            prs.slide_width = _inches(13.333333)
            prs.slide_height = _inches(7.5)
            return
        if key in {"standard", "4:3", "ppt43"}:
            prs.slide_width = _inches(10)
            prs.slide_height = _inches(7.5)
            return
        _die("ERROR: slide_size must be 'wide', 'standard', or an object", 3)
    if isinstance(size, dict):
        width = size.get("width") or size.get("w")
        height = size.get("height") or size.get("h")
        if not width or not height:
            _die("ERROR: slide_size object requires width and height in inches", 3)
        prs.slide_width = _inches(float(width))
        prs.slide_height = _inches(float(height))


def _set_background(slide: Any, theme: dict[str, str], color_key: str = "background_color") -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(theme, color_key)


def _rgb_from_theme_or_hex(theme: dict[str, str], value: Any, fallback_key: str):
    if value in (None, ""):
        return _rgb(theme, fallback_key)
    raw = str(value).strip()
    if raw in theme:
        return _rgb(theme, raw)
    return _hex_to_rgb(raw)


def _style_shape_rgb(shape: Any, fill_rgb: Any, line_rgb: Any | None = None, transparency: int = 0) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if transparency:
        shape.fill.transparency = max(0, min(100, int(transparency)))
    if line_rgb is not None:
        shape.line.color.rgb = line_rgb
    else:
        shape.line.fill.background()


def _style_shape(shape: Any, theme: dict[str, str], fill_key: str, line_key: str | None = None, transparency: int = 0) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(theme, fill_key)
    if transparency:
        shape.fill.transparency = transparency
    if line_key:
        shape.line.color.rgb = _rgb(theme, line_key)
    else:
        shape.line.fill.background()


def _apply_font(paragraph: Any, theme: dict[str, str], size: float, color_key: str, *, bold: bool = False, font_key: str = "body_font") -> None:
    paragraph.font.name = theme.get(font_key) or theme["body_font"]
    paragraph.font.size = _pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = _rgb(theme, color_key)


def _clear_text_frame(tf: Any) -> None:
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = _inches(0.05)
    tf.margin_right = _inches(0.05)
    tf.margin_top = _inches(0.03)
    tf.margin_bottom = _inches(0.03)


def _add_text(
    slide: Any,
    text: Any,
    left: float,
    top: float,
    width: float,
    height: float,
    theme: dict[str, str],
    *,
    size: float = 18,
    color_key: str = "text_color",
    bold: bool = False,
    font_key: str = "body_font",
    align: str | None = None,
    valign: str | None = None,
) -> Any:
    try:
        from pptx.enum.text import MSO_ANCHOR, PP_ALIGN  # type: ignore
    except ImportError:
        _missing("python-pptx")

    box = slide.shapes.add_textbox(_inches(left), _inches(top), _inches(width), _inches(height))
    tf = box.text_frame
    _clear_text_frame(tf)
    if valign == "middle":
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif valign == "bottom":
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM

    lines = str(text or "").splitlines() or [""]
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        if align == "center":
            p.alignment = PP_ALIGN.CENTER
        elif align == "right":
            p.alignment = PP_ALIGN.RIGHT
        _apply_font(p, theme, size, color_key, bold=bold, font_key=font_key)
    return box


def _add_title(slide: Any, title: Any, theme: dict[str, str], slide_w: float, top: float = 0.45) -> None:
    if title:
        _add_text(
            slide,
            title,
            0.65,
            top,
            slide_w - 1.3,
            0.65,
            theme,
            size=30,
            color_key="title_color",
            bold=True,
            font_key="title_font",
        )


def _add_bullets(slide: Any, items: list[Any], left: float, top: float, width: float, height: float, theme: dict[str, str], *, size: float = 17) -> None:
    box = slide.shapes.add_textbox(_inches(left), _inches(top), _inches(width), _inches(height))
    tf = box.text_frame
    _clear_text_frame(tf)
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = f"\u2022 {item}"
        p.space_after = _pt(8)
        _apply_font(p, theme, size, "text_color")


def _add_footer(slide: Any, slide_spec: dict[str, Any], theme: dict[str, str], slide_w: float, slide_h: float, index: int, total: int) -> None:
    footer = slide_spec.get("footer")
    if footer is False:
        return
    label = str(footer) if footer else f"{index}/{total}"
    _add_text(slide, label, slide_w - 1.25, slide_h - 0.45, 0.65, 0.2, theme, size=8, color_key="muted_text_color", align="right")


def _table_rows(value: Any) -> list[list[Any]]:
    if isinstance(value, dict):
        rows = value.get("rows") or value.get("data") or []
        headers = value.get("headers") or value.get("columns")
        if headers:
            rows = [headers] + list(rows)
        value = rows
    if not isinstance(value, list) or not value:
        return []
    return [row if isinstance(row, list) else [row] for row in value]


def _table_options(value: Any) -> dict[str, Any]:
    return value if isinstance(value, dict) else {}


def _format_table_value(value: Any, fmt: str | None = None) -> str:
    if value is None:
        return ""
    if fmt and isinstance(value, (int, float)):
        key = fmt.lower()
        if key in {"percent", "percentage"}:
            return f"{float(value):.0%}" if abs(float(value)) <= 1 else f"{float(value):.0f}%"
        if key in {"currency", "usd"}:
            return f"${float(value):,.0f}"
        if key in {"comma", "number"}:
            return f"{float(value):,.0f}"
        if key.startswith("."):
            return format(float(value), key)
    return str(value)


def _add_table(slide: Any, table: Any, left: float, top: float, width: float, height: float, theme: dict[str, str]) -> None:
    rows = _table_rows(table)
    if not rows:
        return
    options = _table_options(table)
    number_format = options.get("number_format")
    cols = max(len(row) for row in rows)
    shape = slide.shapes.add_table(len(rows), cols, _inches(left), _inches(top), _inches(width), _inches(height))
    column_widths = options.get("column_widths")
    if isinstance(column_widths, list) and len(column_widths) == cols:
        total = sum(float(item) for item in column_widths) or 1
        for ci, item in enumerate(column_widths):
            shape.table.columns[ci].width = _inches(width * float(item) / total)
    for ri, row in enumerate(rows):
        for ci in range(cols):
            cell = shape.table.cell(ri, ci)
            cell.text = _format_table_value(row[ci], number_format) if ci < len(row) else ""
            cell.margin_left = _inches(0.05)
            cell.margin_right = _inches(0.05)
            if ri == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = _rgb(theme, "primary_color")
            elif options.get("banded_rows", True) and ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = _rgb(theme, "muted_surface_color")
            for paragraph in cell.text_frame.paragraphs:
                _apply_font(
                    paragraph,
                    theme,
                    11 if len(rows) > 6 else 12,
                    "inverse_text_color" if ri == 0 else "text_color",
                    bold=ri == 0,
                )
    if options.get("caption"):
        _add_text(slide, options["caption"], left, top + height + 0.08, width, 0.25, theme, size=9, color_key="muted_text_color")


def _item_dicts(value: Any) -> list[dict[str, Any]]:
    if not isinstance(value, list):
        return []
    items: list[dict[str, Any]] = []
    for item in value:
        if isinstance(item, dict):
            items.append(item)
        else:
            items.append({"title": item})
    return items


def _item_title(item: dict[str, Any], fallback: str = "") -> str:
    for key in ("title", "label", "heading", "name", "date"):
        value = item.get(key)
        if value not in (None, ""):
            return str(value)
    return fallback


def _item_detail(item: dict[str, Any]) -> str:
    for key in ("detail", "description", "body", "caption", "text"):
        value = item.get(key)
        if value not in (None, ""):
            return str(value)
    return ""


def _catalog_ref(value: Any) -> str | None:
    if isinstance(value, str) and value.strip():
        return value.strip()
    if isinstance(value, dict):
        for key in ("path", "url", "src", "image", "image_path", "image_url", "file"):
            item = value.get(key)
            if isinstance(item, str) and item.strip():
                return item.strip()
    return None


def _normalize_image_catalog(images: Any) -> dict[str, str]:
    catalog: dict[str, str] = {}
    if isinstance(images, dict):
        for key, value in images.items():
            ref = _catalog_ref(value)
            if ref:
                catalog[str(key).strip().lstrip("@")] = ref
    elif isinstance(images, list):
        for value in images:
            if not isinstance(value, dict):
                continue
            alias = value.get("id") or value.get("name") or value.get("key") or value.get("alias")
            ref = _catalog_ref(value)
            if alias and ref:
                catalog[str(alias).strip().lstrip("@")] = ref
    return catalog


def _icon_catalog_ref(value: Any) -> str | None:
    if isinstance(value, str) and value.strip():
        return value.strip()
    if isinstance(value, dict):
        for key in ("name", "icon", "shape", "kind", "value", "path", "url", "src", "file"):
            item = value.get(key)
            if isinstance(item, str) and item.strip():
                return item.strip()
    return None


def _normalize_icon_catalog(icons: Any) -> dict[str, str]:
    catalog: dict[str, str] = {}
    if isinstance(icons, dict):
        for key, value in icons.items():
            ref = _icon_catalog_ref(value)
            if ref:
                catalog[str(key).strip().lstrip("@")] = ref
    elif isinstance(icons, list):
        for value in icons:
            if not isinstance(value, dict):
                continue
            alias = value.get("id") or value.get("name") or value.get("key") or value.get("alias")
            ref = _icon_catalog_ref(value)
            if alias and ref:
                catalog[str(alias).strip().lstrip("@")] = ref
    return catalog


def _catalog_lookup(value: Any, catalog: dict[str, str]) -> Any:
    if not isinstance(value, str):
        return value
    raw = value.strip()
    key = raw[1:] if raw.startswith("@") else raw
    return catalog.get(key, value)


def _apply_image_catalog_to_slide(slide_spec: dict[str, Any], catalog: dict[str, str]) -> dict[str, Any]:
    if not catalog:
        return dict(slide_spec)

    def visit(value: Any) -> Any:
        if isinstance(value, dict):
            resolved = {
                str(key): visit(item) if isinstance(item, (dict, list)) else item
                for key, item in value.items()
            }
            if "image_id" in resolved and not any(resolved.get(key) for key in ("image", "image_path", "image_url")):
                resolved["image"] = _catalog_lookup(resolved["image_id"], catalog)
            if "background_image_id" in resolved and not any(
                resolved.get(key) for key in ("background_image", "background_image_path", "background_image_url")
            ):
                resolved["background_image"] = _catalog_lookup(resolved["background_image_id"], catalog)
            for key in (
                "image",
                "image_path",
                "image_url",
                "background_image",
                "background_image_path",
                "background_image_url",
            ):
                if key in resolved:
                    resolved[key] = _catalog_lookup(resolved[key], catalog)
            return resolved
        if isinstance(value, list):
            return [visit(item) for item in value]
        return value

    return visit(slide_spec)


def _apply_icon_catalog_to_slide(slide_spec: dict[str, Any], catalog: dict[str, str]) -> dict[str, Any]:
    if not catalog:
        return dict(slide_spec)

    def visit(value: Any) -> Any:
        if isinstance(value, dict):
            resolved = {
                str(key): visit(item) if isinstance(item, (dict, list)) else item
                for key, item in value.items()
            }
            if "icon_id" in resolved and not resolved.get("icon"):
                alias = str(resolved["icon_id"]).strip().lstrip("@")
                resolved["icon"] = _catalog_lookup(alias, catalog)
                resolved.setdefault("icon_label", alias)
            elif isinstance(resolved.get("icon"), str):
                raw = str(resolved["icon"]).strip()
                if raw.startswith("@"):
                    alias = raw[1:]
                    resolved["icon"] = _catalog_lookup(raw, catalog)
                    resolved.setdefault("icon_label", alias)
            return resolved
        if isinstance(value, list):
            return [visit(item) for item in value]
        return value

    return visit(slide_spec)


def _image_ref(slide_spec: dict[str, Any]) -> str | None:
    for key in ("image_path", "image_url", "image"):
        value = slide_spec.get(key)
        if value:
            return str(value)
    return None


def _image_fit(slide_spec: dict[str, Any], default: str = "cover") -> str:
    value = slide_spec.get("image_fit") or slide_spec.get("fit") or slide_spec.get("crop")
    key = str(value or default).strip().lower()
    return key if key in {"cover", "contain", "stretch"} else default


def _resolve_image_reference(ref: str, temp_paths: list[Path], workspace_root: Path) -> Path:
    parsed = urlparse(ref)
    if parsed.scheme in {"http", "https"}:
        suffix = Path(parsed.path).suffix.lower()
        if suffix not in IMAGE_SUFFIXES:
            suffix = ".png"
        tmp = tempfile.NamedTemporaryFile(delete=False, suffix=suffix)
        temp_path = Path(tmp.name)
        tmp.close()
        request = Request(ref, headers={"User-Agent": "Nexa PPTX Renderer/1.0"})
        with urlopen(request, timeout=20) as response:
            temp_path.write_bytes(response.read())
        temp_paths.append(temp_path)
        return temp_path

    path = Path(ref)
    if not path.is_absolute():
        path = workspace_root / path
    return _validate_path(str(path), workspace_root=workspace_root)


def _add_image_or_placeholder(
    slide: Any,
    ref: str | None,
    left: float,
    top: float,
    width: float,
    height: float,
    theme: dict[str, str],
    temp_paths: list[Path],
    workspace_root: Path,
    *,
    caption: Any = None,
    fill_width_only: bool = False,
    fit: str = "cover",
) -> None:
    try:
        if ref:
            image_path = _resolve_image_reference(ref, temp_paths, workspace_root)
            if fill_width_only:
                slide.shapes.add_picture(str(image_path), _inches(left), _inches(top), width=_inches(width))
            else:
                _add_picture_with_fit(slide, image_path, left, top, width, height, fit)
            if caption:
                _add_text(slide, caption, left, top + height + 0.05, width, 0.25, theme, size=9, color_key="muted_text_color")
            return
    except Exception as exc:
        print(f"WARN: could not add image {ref}: {exc}", file=sys.stderr)

    try:
        from pptx.enum.shapes import MSO_SHAPE  # type: ignore
    except ImportError:
        _missing("python-pptx")
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, _inches(left), _inches(top), _inches(width), _inches(height))
    _style_shape(shape, theme, "muted_surface_color", "muted_text_color")
    _add_text(slide, "Image unavailable", left + 0.15, top + height / 2 - 0.15, width - 0.3, 0.3, theme, size=12, color_key="muted_text_color", align="center")


def _safe_token(value: Any, fallback: str = "item") -> str:
    raw = str(value or fallback).strip().lower()
    token = re.sub(r"[^a-z0-9]+", "-", raw).strip("-")
    return token or fallback


def _looks_like_image_ref(value: str, workspace_root: Path) -> bool:
    parsed = urlparse(value)
    if parsed.scheme in {"http", "https"}:
        return Path(parsed.path).suffix.lower() in IMAGE_SUFFIXES
    path = Path(value)
    if not path.is_absolute():
        path = workspace_root / path
    return path.suffix.lower() in IMAGE_SUFFIXES


def _add_builtin_icon(slide: Any, icon: str, label: str, left: float, top: float, size: float, theme: dict[str, str]) -> None:
    try:
        from pptx.enum.shapes import MSO_SHAPE  # type: ignore
    except ImportError:
        _missing("python-pptx")

    key = _safe_token(icon, "symbol")
    name = f"icon-{_safe_token(label, key)}"
    base = slide.shapes.add_shape(MSO_SHAPE.OVAL, _inches(left), _inches(top), _inches(size), _inches(size))
    base.name = name
    _style_shape(base, theme, "primary_color")

    if key in {"shield", "risk", "security", "trust"}:
        mark = slide.shapes.add_shape(getattr(MSO_SHAPE, "PENTAGON", MSO_SHAPE.DIAMOND), _inches(left + size * 0.28), _inches(top + size * 0.22), _inches(size * 0.44), _inches(size * 0.52))
        mark.name = f"{name}-mark"
        _style_shape(mark, theme, "inverse_text_color", transparency=0)
    elif key in {"trend", "growth", "up", "signal"}:
        mark = slide.shapes.add_shape(getattr(MSO_SHAPE, "UP_ARROW", MSO_SHAPE.RIGHT_ARROW), _inches(left + size * 0.28), _inches(top + size * 0.24), _inches(size * 0.46), _inches(size * 0.5))
        mark.name = f"{name}-mark"
        _style_shape(mark, theme, "inverse_text_color")
    elif key in {"network", "workflow", "process"}:
        for idx, (dx, dy) in enumerate([(0.26, 0.28), (0.56, 0.28), (0.41, 0.58)], start=1):
            node = slide.shapes.add_shape(MSO_SHAPE.OVAL, _inches(left + size * dx), _inches(top + size * dy), _inches(size * 0.16), _inches(size * 0.16))
            node.name = f"{name}-node-{idx}"
            _style_shape(node, theme, "inverse_text_color")
    elif key in {"spark", "idea", "insight"}:
        mark = slide.shapes.add_shape(getattr(MSO_SHAPE, "SUN", MSO_SHAPE.DIAMOND), _inches(left + size * 0.27), _inches(top + size * 0.27), _inches(size * 0.46), _inches(size * 0.46))
        mark.name = f"{name}-mark"
        _style_shape(mark, theme, "inverse_text_color")
    elif key in {"check", "done", "success"}:
        mark = slide.shapes.add_shape(getattr(MSO_SHAPE, "CHEVRON", MSO_SHAPE.RIGHT_TRIANGLE), _inches(left + size * 0.24), _inches(top + size * 0.33), _inches(size * 0.5), _inches(size * 0.32))
        mark.name = f"{name}-mark"
        _style_shape(mark, theme, "inverse_text_color")
    else:
        mark = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, _inches(left + size * 0.31), _inches(top + size * 0.31), _inches(size * 0.38), _inches(size * 0.38))
        mark.name = f"{name}-mark"
        _style_shape(mark, theme, "inverse_text_color")


def _add_icon_or_image(
    slide: Any,
    icon: Any,
    label: Any,
    left: float,
    top: float,
    size: float,
    theme: dict[str, str],
    temp_paths: list[Path],
    workspace_root: Path,
) -> None:
    if icon in (None, ""):
        return
    value = str(icon).strip()
    name = str(label or value).strip()
    if _looks_like_image_ref(value, workspace_root):
        _add_image_or_placeholder(slide, value, left, top, size, size, theme, temp_paths, workspace_root, fit="contain")
        try:
            slide.shapes[-1].name = f"icon-{_safe_token(name)}"
        except Exception:
            pass
        return
    _add_builtin_icon(slide, value, name, left, top, size, theme)


def _add_slide_icon(slide: Any, slide_spec: dict[str, Any], theme: dict[str, str], slide_w: float, temp_paths: list[Path], workspace_root: Path) -> None:
    icon = slide_spec.get("icon")
    if not icon:
        return
    _add_icon_or_image(
        slide,
        icon,
        slide_spec.get("icon_label") or slide_spec.get("icon_id") or icon,
        slide_w - 1.15,
        0.43,
        0.48,
        theme,
        temp_paths,
        workspace_root,
    )


def _background_dict(slide_spec: dict[str, Any]) -> dict[str, Any]:
    value = slide_spec.get("background")
    return value if isinstance(value, dict) else {}


def _background_value(slide_spec: dict[str, Any], *keys: str) -> Any:
    background = _background_dict(slide_spec)
    for key in keys:
        if key in background and background[key] not in (None, ""):
            return background[key]
        if key in slide_spec and slide_spec[key] not in (None, ""):
            return slide_spec[key]
    return None


def _background_image_ref(slide_spec: dict[str, Any]) -> str | None:
    background = _background_dict(slide_spec)
    for key in ("image", "image_path", "image_url", "background_image", "background_image_path", "background_image_url"):
        value = background.get(key)
        if value:
            return str(value)
    for key in ("background_image", "background_image_path", "background_image_url"):
        value = slide_spec.get(key)
        if value:
            return str(value)
    return None


def _background_int(slide_spec: dict[str, Any], default: int, *keys: str) -> int:
    value = _background_value(slide_spec, *keys)
    if value in (None, ""):
        return default
    try:
        return max(0, min(100, int(value)))
    except (TypeError, ValueError):
        return default


def _background_bool(slide_spec: dict[str, Any], default: bool, *keys: str) -> bool:
    value = _background_value(slide_spec, *keys)
    if value in (None, ""):
        return default
    if isinstance(value, bool):
        return value
    return str(value).strip().lower() in {"1", "true", "yes", "on"}


def _image_pixel_size(path: Path) -> tuple[int, int] | None:
    try:
        from PIL import Image  # type: ignore
    except ImportError:
        return None
    try:
        with Image.open(path) as image:
            return image.size
    except Exception:
        return None


def _add_picture_with_fit(slide: Any, image_path: Path, left: float, top: float, width: float, height: float, fit: str) -> Any:
    size = _image_pixel_size(image_path)
    fit_mode = str(fit or "cover").lower()
    if fit_mode == "stretch" or not size:
        return slide.shapes.add_picture(str(image_path), _inches(left), _inches(top), width=_inches(width), height=_inches(height))

    image_w, image_h = size
    image_ratio = image_w / image_h
    frame_ratio = width / height

    if fit_mode == "contain":
        if image_ratio > frame_ratio:
            pic_w = width
            pic_h = width / image_ratio
        else:
            pic_h = height
            pic_w = height * image_ratio
        return slide.shapes.add_picture(
            str(image_path),
            _inches(left + (width - pic_w) / 2),
            _inches(top + (height - pic_h) / 2),
            width=_inches(pic_w),
            height=_inches(pic_h),
        )

    picture = slide.shapes.add_picture(str(image_path), _inches(left), _inches(top), width=_inches(width), height=_inches(height))
    try:
        if image_ratio > frame_ratio:
            visible = image_h * frame_ratio
            crop = max(0.0, min(0.45, (image_w - visible) / image_w / 2))
            picture.crop_left = crop
            picture.crop_right = crop
        elif image_ratio < frame_ratio:
            visible = image_w / frame_ratio
            crop = max(0.0, min(0.45, (image_h - visible) / image_h / 2))
            picture.crop_top = crop
            picture.crop_bottom = crop
    except Exception:
        pass
    return picture


def _add_background_picture(
    slide: Any,
    ref: str,
    slide_w: float,
    slide_h: float,
    temp_paths: list[Path],
    workspace_root: Path,
    *,
    fit: str = "cover",
) -> bool:
    try:
        image_path = _resolve_image_reference(ref, temp_paths, workspace_root)
        if image_path.suffix.lower() == ".svg":
            print(
                f"WARN: SVG backgrounds are not embedded directly by python-pptx: {ref}. "
                "Use a raster preview or a template/SVG pipeline for exact SVG art.",
                file=sys.stderr,
            )
            return False
        size = _image_pixel_size(image_path)
        if not size:
            slide.shapes.add_picture(str(image_path), _inches(0), _inches(0), width=_inches(slide_w), height=_inches(slide_h))
            return True
        image_w, image_h = size
        image_ratio = image_w / image_h
        slide_ratio = slide_w / slide_h
        fit_mode = str(fit or "cover").lower()
        if fit_mode == "contain":
            if image_ratio > slide_ratio:
                pic_w = slide_w
                pic_h = slide_w / image_ratio
            else:
                pic_h = slide_h
                pic_w = slide_h * image_ratio
        else:
            if image_ratio > slide_ratio:
                pic_h = slide_h
                pic_w = slide_h * image_ratio
            else:
                pic_w = slide_w
                pic_h = slide_w / image_ratio
        left = (slide_w - pic_w) / 2
        top = (slide_h - pic_h) / 2
        slide.shapes.add_picture(str(image_path), _inches(left), _inches(top), width=_inches(pic_w), height=_inches(pic_h))
        return True
    except Exception as exc:
        print(f"WARN: could not add background image {ref}: {exc}", file=sys.stderr)
        return False


def _add_background_art(slide: Any, style: str, theme: dict[str, str], slide_w: float, slide_h: float) -> None:
    try:
        from pptx.enum.shapes import MSO_SHAPE  # type: ignore
    except ImportError:
        _missing("python-pptx")

    key = str(style or "").strip().lower()
    if key in {"", "none", "solid", "flat"}:
        return

    if key in {"blueprint_grid", "clinical_grid", "data_grid"}:
        spacing = 0.45 if key == "data_grid" else 0.5
        line_key = "primary_color" if key != "clinical_grid" else "accent_color"
        x = 0.0
        while x <= slide_w:
            line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, _inches(x), _inches(0), _inches(0.006), _inches(slide_h))
            _style_shape(line, theme, line_key, transparency=84)
            x += spacing
        y = 0.0
        while y <= slide_h:
            line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, _inches(0), _inches(y), _inches(slide_w), _inches(0.006))
            _style_shape(line, theme, line_key, transparency=86)
            y += spacing
        if key == "data_grid":
            for idx in range(8):
                dot = slide.shapes.add_shape(
                    MSO_SHAPE.OVAL,
                    _inches(0.7 + idx * 1.35),
                    _inches(slide_h - 0.95 - (idx % 3) * 0.22),
                    _inches(0.055),
                    _inches(0.055),
                )
                _style_shape(dot, theme, "accent_color", transparency=20)
        return

    if key == "paper_texture":
        for idx in range(14):
            left = (idx * 1.07) % slide_w
            top = (idx * 0.71) % slide_h
            width = 0.55 + (idx % 3) * 0.18
            height = 0.12 + (idx % 2) * 0.08
            patch = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, _inches(left), _inches(top), _inches(width), _inches(height))
            patch.rotation = -6 + (idx % 5) * 3
            _style_shape(patch, theme, "muted_surface_color", transparency=78)
        return

    if key == "spotlight":
        panel = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, _inches(-0.25), _inches(0.75), _inches(slide_w * 0.72), _inches(slide_h * 0.68))
        panel.rotation = -4
        _style_shape(panel, theme, "surface_color", transparency=16)
        edge = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, _inches(0.7), _inches(1.02), _inches(0.08), _inches(slide_h * 0.58))
        _style_shape(edge, theme, "accent_color")
        return

    if key in {"diagonal", "section", "editorial"}:
        band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, _inches(slide_w * 0.62), _inches(-0.45), _inches(slide_w * 0.5), _inches(slide_h + 0.9))
        band.rotation = -12
        _style_shape(band, theme, "accent_color", transparency=22)
        veil = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, _inches(-0.2), _inches(slide_h * 0.72), _inches(slide_w * 1.2), _inches(slide_h * 0.28))
        _style_shape(veil, theme, "muted_surface_color", transparency=55)
        return

    if key in {"mesh", "gradient_mesh", "soft_geometry", "ambient"}:
        orb = slide.shapes.add_shape(MSO_SHAPE.OVAL, _inches(slide_w * 0.68), _inches(-0.65), _inches(slide_w * 0.42), _inches(slide_h * 0.72))
        _style_shape(orb, theme, "primary_color", transparency=72)
        orb2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, _inches(-0.65), _inches(slide_h * 0.62), _inches(slide_w * 0.38), _inches(slide_h * 0.5))
        _style_shape(orb2, theme, "accent_color", transparency=76)
        rail = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, _inches(0), _inches(0), _inches(slide_w), _inches(0.08))
        _style_shape(rail, theme, "primary_color", transparency=20)
        return

    motif = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, _inches(slide_w * 0.74), _inches(0), _inches(slide_w * 0.26), _inches(slide_h))
    _style_shape(motif, theme, "muted_surface_color", transparency=38)


def _apply_slide_background(
    slide: Any,
    slide_spec: dict[str, Any],
    theme: dict[str, str],
    slide_w: float,
    slide_h: float,
    temp_paths: list[Path],
    workspace_root: Path,
    *,
    color_key: str = "background_color",
    default_style: str = "soft_geometry",
) -> None:
    try:
        from pptx.enum.shapes import MSO_SHAPE  # type: ignore
    except ImportError:
        _missing("python-pptx")

    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _rgb_from_theme_or_hex(
        theme,
        _background_value(slide_spec, "color", "background_color", "bg"),
        color_key,
    )

    image_added = False
    image_ref = _background_image_ref(slide_spec)
    if image_ref:
        fit = str(_background_value(slide_spec, "fit", "image_fit", "background_fit") or "cover")
        image_added = _add_background_picture(slide, image_ref, slide_w, slide_h, temp_paths, workspace_root, fit=fit)
        overlay_value = _background_value(slide_spec, "overlay_color", "background_overlay_color")
        if image_added and overlay_value is not False and str(overlay_value).lower() != "none":
            overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, _inches(0), _inches(0), _inches(slide_w), _inches(slide_h))
            _style_shape_rgb(
                overlay,
                _rgb_from_theme_or_hex(theme, overlay_value, "background_color"),
                transparency=_background_int(slide_spec, 24, "overlay_transparency", "background_overlay_transparency"),
            )

    style = str(_background_value(slide_spec, "style", "background_style") or theme.get("background_style") or default_style)
    if (not image_added) or _background_bool(slide_spec, False, "ornaments_over_image", "background_ornaments"):
        _add_background_art(slide, style, theme, slide_w, slide_h)


def _render_title_slide(slide: Any, slide_spec: dict[str, Any], theme: dict[str, str], slide_w: float, slide_h: float, temp_paths: list[Path], workspace_root: Path) -> None:
    try:
        from pptx.enum.shapes import MSO_SHAPE  # type: ignore
    except ImportError:
        _missing("python-pptx")
    _apply_slide_background(slide, slide_spec, theme, slide_w, slide_h, temp_paths, workspace_root, default_style="diagonal")
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, _inches(0), _inches(0), _inches(slide_w * 0.42), _inches(slide_h))
    _style_shape(band, theme, "primary_color")
    ref = _image_ref(slide_spec)
    if ref:
        _add_image_or_placeholder(slide, ref, slide_w * 0.42, 0, slide_w * 0.58, slide_h, theme, temp_paths, workspace_root)
    _add_text(slide, slide_spec.get("title") or "", 0.75, 1.4, slide_w * 0.36, 1.5, theme, size=38, color_key="inverse_text_color", bold=True, font_key="title_font")
    if slide_spec.get("subtitle"):
        _add_text(slide, slide_spec["subtitle"], 0.8, 3.05, slide_w * 0.34, 0.9, theme, size=17, color_key="inverse_text_color")
    if slide_spec.get("author"):
        _add_text(slide, slide_spec["author"], 0.8, slide_h - 1.0, slide_w * 0.34, 0.35, theme, size=11, color_key="inverse_text_color")


def _render_agenda_slide(slide: Any, slide_spec: dict[str, Any], theme: dict[str, str], slide_w: float, slide_h: float, temp_paths: list[Path], workspace_root: Path) -> None:
    try:
        from pptx.enum.shapes import MSO_SHAPE  # type: ignore
        from pptx.enum.text import MSO_ANCHOR, PP_ALIGN  # type: ignore
    except ImportError:
        _missing("python-pptx")
    _apply_slide_background(slide, slide_spec, theme, slide_w, slide_h, temp_paths, workspace_root)
    _add_title(slide, slide_spec.get("title") or "Agenda", theme, slide_w)
    items = slide_spec.get("items") or slide_spec.get("bullets") or []
    if not isinstance(items, list):
        _die("ERROR: agenda slide requires items array", 3)
    y = 1.45
    for idx, item in enumerate(items[:8], start=1):
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, _inches(0.9), _inches(y), _inches(0.42), _inches(0.42))
        _style_shape(dot, theme, "primary_color")
        tf = dot.text_frame
        _clear_text_frame(tf)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = str(idx)
        p.alignment = PP_ALIGN.CENTER
        _apply_font(p, theme, 12, "inverse_text_color", bold=True)
        _add_text(slide, item, 1.5, y - 0.04, slide_w - 2.3, 0.45, theme, size=18, color_key="text_color")
        y += 0.62


def _render_body_slide(slide: Any, slide_spec: dict[str, Any], theme: dict[str, str], slide_w: float, slide_h: float, temp_paths: list[Path], workspace_root: Path) -> None:
    _apply_slide_background(slide, slide_spec, theme, slide_w, slide_h, temp_paths, workspace_root)
    _add_title(slide, slide_spec.get("title"), theme, slide_w)
    ref = _image_ref(slide_spec)
    content_w = slide_w - 1.4
    if ref:
        content_w = slide_w * 0.54
        _add_image_or_placeholder(
            slide,
            ref,
            slide_w * 0.62,
            1.35,
            slide_w * 0.32,
            slide_h - 2.1,
            theme,
            temp_paths,
            workspace_root,
            caption=slide_spec.get("image_caption"),
            fit=_image_fit(slide_spec),
        )
    paragraph = slide_spec.get("paragraph") or slide_spec.get("body")
    bullets = slide_spec.get("bullets") or []
    if paragraph:
        _add_text(slide, paragraph, 0.75, 1.35, content_w, 1.45 if bullets else slide_h - 2.0, theme, size=17, color_key="text_color")
    if isinstance(bullets, list) and bullets:
        bullet_top = 2.75 if paragraph else 1.35
        _add_bullets(slide, bullets[:6], 0.78, bullet_top, content_w, slide_h - bullet_top - 0.75, theme)
    if slide_spec.get("table"):
        _add_table(slide, slide_spec["table"], 0.75, 2.0, content_w, slide_h - 2.8, theme)


def _column_payload(value: Any) -> dict[str, Any]:
    if isinstance(value, dict):
        return value
    if isinstance(value, list):
        return {"bullets": value}
    return {"paragraph": value}


def _render_column(slide: Any, payload: dict[str, Any], left: float, top: float, width: float, height: float, theme: dict[str, str], temp_paths: list[Path], workspace_root: Path) -> None:
    try:
        from pptx.enum.shapes import MSO_SHAPE  # type: ignore
    except ImportError:
        _missing("python-pptx")
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, _inches(left), _inches(top), _inches(width), _inches(height))
    _style_shape(panel, theme, "surface_color", "muted_surface_color")
    y = top + 0.25
    if payload.get("heading"):
        _add_text(slide, payload["heading"], left + 0.25, y, width - 0.5, 0.35, theme, size=17, color_key="title_color", bold=True, font_key="title_font")
        y += 0.55
    ref = _image_ref(payload)
    if ref:
        _add_image_or_placeholder(slide, ref, left + 0.25, y, width - 0.5, height * 0.38, theme, temp_paths, workspace_root, fit=_image_fit(payload))
        y += height * 0.43
    paragraph = payload.get("paragraph") or payload.get("body")
    if paragraph:
        _add_text(slide, paragraph, left + 0.25, y, width - 0.5, 0.85, theme, size=13.5, color_key="text_color")
        y += 1.0
    bullets = payload.get("bullets") or payload.get("items") or []
    if isinstance(bullets, list) and bullets:
        _add_bullets(slide, bullets[:5], left + 0.25, y, width - 0.5, top + height - y - 0.25, theme, size=13.5)


def _render_two_column_slide(slide: Any, slide_spec: dict[str, Any], theme: dict[str, str], slide_w: float, slide_h: float, temp_paths: list[Path], workspace_root: Path) -> None:
    _apply_slide_background(slide, slide_spec, theme, slide_w, slide_h, temp_paths, workspace_root)
    _add_title(slide, slide_spec.get("title"), theme, slide_w)
    gap = 0.35
    left = 0.75
    top = 1.35
    col_w = (slide_w - 1.5 - gap) / 2
    height = slide_h - 2.0
    _render_column(slide, _column_payload(slide_spec.get("left") or {}), left, top, col_w, height, theme, temp_paths, workspace_root)
    _render_column(slide, _column_payload(slide_spec.get("right") or {}), left + col_w + gap, top, col_w, height, theme, temp_paths, workspace_root)


def _render_stat_slide(slide: Any, slide_spec: dict[str, Any], theme: dict[str, str], slide_w: float, slide_h: float, temp_paths: list[Path], workspace_root: Path) -> None:
    try:
        from pptx.enum.shapes import MSO_SHAPE  # type: ignore
    except ImportError:
        _missing("python-pptx")
    _apply_slide_background(slide, slide_spec, theme, slide_w, slide_h, temp_paths, workspace_root)
    _add_title(slide, slide_spec.get("title"), theme, slide_w)
    stats = slide_spec.get("stats") or []
    if not isinstance(stats, list) or not stats:
        _die("ERROR: stat slide requires non-empty stats array", 3)
    stats = stats[:4]
    gap = 0.25
    card_w = (slide_w - 1.5 - gap * (len(stats) - 1)) / len(stats)
    top = 2.0
    for idx, item in enumerate(stats):
        if not isinstance(item, dict):
            item = {"value": item, "label": ""}
        left = 0.75 + idx * (card_w + gap)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, _inches(left), _inches(top), _inches(card_w), _inches(3.0))
        _style_shape(card, theme, "surface_color", "muted_surface_color")
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, _inches(left), _inches(top), _inches(card_w), _inches(0.08))
        _style_shape(accent, theme, "accent_color")
        _add_text(slide, item.get("value", ""), left + 0.2, top + 0.55, card_w - 0.4, 0.7, theme, size=32, color_key="primary_color", bold=True, align="center", font_key="title_font")
        _add_text(slide, item.get("label", ""), left + 0.25, top + 1.35, card_w - 0.5, 0.45, theme, size=15, color_key="title_color", bold=True, align="center")
        if item.get("caption"):
            _add_text(slide, item["caption"], left + 0.25, top + 1.95, card_w - 0.5, 0.7, theme, size=11, color_key="muted_text_color", align="center")


def _render_quote_slide(slide: Any, slide_spec: dict[str, Any], theme: dict[str, str], slide_w: float, slide_h: float, temp_paths: list[Path], workspace_root: Path) -> None:
    try:
        from pptx.enum.shapes import MSO_SHAPE  # type: ignore
    except ImportError:
        _missing("python-pptx")
    _apply_slide_background(slide, slide_spec, theme, slide_w, slide_h, temp_paths, workspace_root)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, _inches(0.75), _inches(1.45), _inches(0.12), _inches(slide_h - 2.8))
    _style_shape(bar, theme, "accent_color")
    quote = slide_spec.get("text") or slide_spec.get("quote") or slide_spec.get("body") or ""
    _add_text(slide, quote, 1.08, 1.35, slide_w - 2.0, 2.5, theme, size=28, color_key="title_color", bold=True, font_key="title_font")
    if slide_spec.get("attribution"):
        _add_text(slide, slide_spec["attribution"], 1.12, 4.25, slide_w - 2.2, 0.35, theme, size=14, color_key="muted_text_color")


def _render_section_slide(slide: Any, slide_spec: dict[str, Any], theme: dict[str, str], slide_w: float, slide_h: float, temp_paths: list[Path], workspace_root: Path) -> None:
    _apply_slide_background(slide, slide_spec, theme, slide_w, slide_h, temp_paths, workspace_root, color_key="primary_color", default_style="diagonal")
    _add_text(slide, slide_spec.get("title") or "", 1.05, 2.25, slide_w - 2.1, 0.95, theme, size=36, color_key="inverse_text_color", bold=True, font_key="title_font", align="center")
    if slide_spec.get("subtitle"):
        _add_text(slide, slide_spec["subtitle"], 1.5, 3.25, slide_w - 3.0, 0.55, theme, size=17, color_key="inverse_text_color", align="center")


def _render_image_full_slide(slide: Any, slide_spec: dict[str, Any], theme: dict[str, str], slide_w: float, slide_h: float, temp_paths: list[Path], workspace_root: Path) -> None:
    try:
        from pptx.enum.shapes import MSO_SHAPE  # type: ignore
    except ImportError:
        _missing("python-pptx")
    _apply_slide_background(slide, slide_spec, theme, slide_w, slide_h, temp_paths, workspace_root)
    _add_image_or_placeholder(slide, _image_ref(slide_spec), 0, 0, slide_w, slide_h, theme, temp_paths, workspace_root, fill_width_only=True)
    overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, _inches(0), _inches(slide_h * 0.58), _inches(slide_w), _inches(slide_h * 0.42))
    _style_shape(overlay, theme, "background_color", transparency=18)
    if slide_spec.get("title"):
        _add_text(slide, slide_spec["title"], 0.8, slide_h * 0.63, slide_w - 1.6, 0.75, theme, size=30, color_key="title_color", bold=True, font_key="title_font")
    if slide_spec.get("caption"):
        _add_text(slide, slide_spec["caption"], 0.85, slide_h * 0.75, slide_w - 1.7, 0.5, theme, size=13, color_key="text_color")


def _render_timeline_slide(slide: Any, slide_spec: dict[str, Any], theme: dict[str, str], slide_w: float, slide_h: float, temp_paths: list[Path], workspace_root: Path) -> None:
    try:
        from pptx.enum.shapes import MSO_SHAPE  # type: ignore
    except ImportError:
        _missing("python-pptx")
    _apply_slide_background(slide, slide_spec, theme, slide_w, slide_h, temp_paths, workspace_root)
    _add_title(slide, slide_spec.get("title"), theme, slide_w)
    events = _item_dicts(slide_spec.get("events") or slide_spec.get("items"))[:6]
    line_y = slide_h * 0.48
    left = 0.9
    usable = slide_w - 1.8
    rail = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, _inches(left), _inches(line_y), _inches(usable), _inches(0.035))
    _style_shape(rail, theme, "accent_color")
    step = usable / max(1, len(events) - 1)
    card_w = min(2.05, max(1.5, usable / max(1, len(events)) - 0.12))
    for idx, event in enumerate(events):
        x = left + step * idx
        node = slide.shapes.add_shape(MSO_SHAPE.OVAL, _inches(x - 0.12), _inches(line_y - 0.12), _inches(0.24), _inches(0.24))
        _style_shape(node, theme, "primary_color", "background_color")
        date = event.get("date") or event.get("label") or f"{idx + 1}"
        _add_text(slide, date, x - 0.55, line_y - 0.55, 1.1, 0.25, theme, size=10, color_key="primary_color", bold=True, align="center")
        card_top = 1.55 if idx % 2 == 0 else line_y + 0.55
        card_left = max(0.55, min(slide_w - card_w - 0.55, x - card_w / 2))
        connector_top = card_top + 1.08 if idx % 2 == 0 else line_y + 0.12
        connector_h = abs(line_y - connector_top)
        connector = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, _inches(x - 0.01), _inches(min(line_y, connector_top)), _inches(0.02), _inches(max(0.12, connector_h)))
        _style_shape(connector, theme, "muted_surface_color")
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, _inches(card_left), _inches(card_top), _inches(card_w), _inches(1.15))
        _style_shape(card, theme, "surface_color", "muted_surface_color")
        _add_text(slide, _item_title(event, str(date)), card_left + 0.12, card_top + 0.12, card_w - 0.24, 0.28, theme, size=12, color_key="title_color", bold=True)
        detail = _item_detail(event)
        if detail:
            _add_text(slide, detail, card_left + 0.12, card_top + 0.47, card_w - 0.24, 0.48, theme, size=9.5, color_key="muted_text_color")


def _render_process_slide(slide: Any, slide_spec: dict[str, Any], theme: dict[str, str], slide_w: float, slide_h: float, temp_paths: list[Path], workspace_root: Path) -> None:
    try:
        from pptx.enum.shapes import MSO_SHAPE  # type: ignore
    except ImportError:
        _missing("python-pptx")
    _apply_slide_background(slide, slide_spec, theme, slide_w, slide_h, temp_paths, workspace_root)
    _add_title(slide, slide_spec.get("title"), theme, slide_w)
    steps = _item_dicts(slide_spec.get("steps") or slide_spec.get("items"))[:6]
    gap = 0.18
    left = 0.65
    top = 2.05
    height = min(3.0, slide_h - 3.05)
    card_w = (slide_w - 1.3 - gap * (len(steps) - 1)) / len(steps)
    for idx, step in enumerate(steps):
        x = left + idx * (card_w + gap)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, _inches(x), _inches(top), _inches(card_w), _inches(height))
        _style_shape(card, theme, "surface_color", "muted_surface_color")
        badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, _inches(x + 0.18), _inches(top + 0.2), _inches(0.42), _inches(0.42))
        _style_shape(badge, theme, "primary_color")
        _add_text(slide, idx + 1, x + 0.18, top + 0.28, 0.42, 0.16, theme, size=10, color_key="inverse_text_color", bold=True, align="center", valign="middle")
        if step.get("icon"):
            _add_icon_or_image(
                slide,
                step.get("icon"),
                step.get("icon_label") or step.get("icon_id") or step.get("icon"),
                x + card_w - 0.68,
                top + 0.18,
                0.38,
                theme,
                temp_paths,
                workspace_root,
            )
        _add_text(slide, _item_title(step, f"Step {idx + 1}"), x + 0.2, top + 0.78, card_w - 0.4, 0.48, theme, size=14, color_key="title_color", bold=True)
        detail = _item_detail(step)
        if detail:
            _add_text(slide, detail, x + 0.2, top + 1.35, card_w - 0.4, 0.95, theme, size=10.5, color_key="muted_text_color")
        if idx < len(steps) - 1:
            connector = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, _inches(x + card_w - 0.04), _inches(top + height / 2 - 0.11), _inches(gap + 0.08), _inches(0.22))
            _style_shape(connector, theme, "accent_color")


def _comparison_columns(slide_spec: dict[str, Any]) -> list[dict[str, Any]]:
    columns = _item_dicts(slide_spec.get("columns"))
    if len(columns) >= 2:
        return columns[:3]
    left = slide_spec.get("left")
    right = slide_spec.get("right")
    if isinstance(left, dict) and isinstance(right, dict):
        return [left, right]
    return []


def _render_comparison_slide(slide: Any, slide_spec: dict[str, Any], theme: dict[str, str], slide_w: float, slide_h: float, temp_paths: list[Path], workspace_root: Path) -> None:
    try:
        from pptx.enum.shapes import MSO_SHAPE  # type: ignore
    except ImportError:
        _missing("python-pptx")
    _apply_slide_background(slide, slide_spec, theme, slide_w, slide_h, temp_paths, workspace_root)
    _add_title(slide, slide_spec.get("title"), theme, slide_w)
    columns = _comparison_columns(slide_spec)
    gap = 0.28
    left = 0.7
    top = 1.55
    height = slide_h - 2.35
    col_w = (slide_w - 1.4 - gap * (len(columns) - 1)) / len(columns)
    for idx, column in enumerate(columns):
        x = left + idx * (col_w + gap)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, _inches(x), _inches(top), _inches(col_w), _inches(height))
        _style_shape(card, theme, "surface_color", "muted_surface_color")
        accent_key = "primary_color" if idx == 0 else "accent_color"
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, _inches(x), _inches(top), _inches(col_w), _inches(0.09))
        _style_shape(accent, theme, accent_key)
        _render_column(slide, _column_payload(column), x + 0.22, top + 0.28, col_w - 0.44, height - 0.5, theme, temp_paths, workspace_root)


def _render_matrix_slide(slide: Any, slide_spec: dict[str, Any], theme: dict[str, str], slide_w: float, slide_h: float, temp_paths: list[Path], workspace_root: Path) -> None:
    try:
        from pptx.enum.shapes import MSO_SHAPE  # type: ignore
    except ImportError:
        _missing("python-pptx")
    _apply_slide_background(slide, slide_spec, theme, slide_w, slide_h, temp_paths, workspace_root)
    _add_title(slide, slide_spec.get("title"), theme, slide_w)
    quadrants = _item_dicts(slide_spec.get("quadrants") or slide_spec.get("items"))[:4]
    left = 1.0
    top = 1.45
    width = slide_w - 2.0
    height = slide_h - 2.3
    cell_w = width / 2
    cell_h = height / 2
    for idx in range(4):
        quadrant = quadrants[idx] if idx < len(quadrants) else {"title": ""}
        row = idx // 2
        col = idx % 2
        x = left + col * cell_w
        y = top + row * cell_h
        fill_key = "surface_color" if idx in {0, 3} else "muted_surface_color"
        cell = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, _inches(x), _inches(y), _inches(cell_w), _inches(cell_h))
        _style_shape(cell, theme, fill_key, "background_color")
        _add_text(slide, _item_title(quadrant), x + 0.22, y + 0.2, cell_w - 0.44, 0.36, theme, size=15, color_key="title_color", bold=True)
        detail = _item_detail(quadrant)
        if detail:
            _add_text(slide, detail, x + 0.22, y + 0.72, cell_w - 0.44, cell_h - 0.95, theme, size=11, color_key="muted_text_color")
    vline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, _inches(left + cell_w - 0.01), _inches(top), _inches(0.02), _inches(height))
    _style_shape(vline, theme, "primary_color")
    hline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, _inches(left), _inches(top + cell_h - 0.01), _inches(width), _inches(0.02))
    _style_shape(hline, theme, "primary_color")
    if slide_spec.get("x_axis"):
        _add_text(slide, slide_spec["x_axis"], left, top + height + 0.08, width, 0.25, theme, size=10, color_key="muted_text_color", align="center")
    if slide_spec.get("y_axis"):
        _add_text(slide, slide_spec["y_axis"], 0.25, top + height / 2 - 0.2, 0.55, 0.4, theme, size=10, color_key="muted_text_color", align="center")


def _chart_series(slide_spec: dict[str, Any]) -> list[dict[str, Any]]:
    result: list[dict[str, Any]] = []
    for idx, series in enumerate(slide_spec.get("series") or []):
        if isinstance(series, dict):
            result.append({"name": str(series.get("name") or f"Series {idx + 1}"), "values": series.get("values") or []})
        else:
            result.append({"name": f"Series {idx + 1}", "values": series})
    return result


def _chart_type(chart_type_key: str, xl_chart_type: Any):
    return {
        "bar": xl_chart_type.BAR_CLUSTERED,
        "column": xl_chart_type.COLUMN_CLUSTERED,
        "line": xl_chart_type.LINE_MARKERS,
        "area": xl_chart_type.AREA,
        "bar_stacked": xl_chart_type.BAR_STACKED,
        "stacked_bar": xl_chart_type.BAR_STACKED,
        "column_stacked": xl_chart_type.COLUMN_STACKED,
        "stacked_column": xl_chart_type.COLUMN_STACKED,
        "pie": xl_chart_type.PIE,
        "doughnut": getattr(xl_chart_type, "DOUGHNUT", xl_chart_type.PIE),
    }.get(chart_type_key, xl_chart_type.COLUMN_CLUSTERED)


def _add_chart(slide: Any, slide_spec: dict[str, Any], theme: dict[str, str], left: float, top: float, width: float, height: float) -> None:
    try:
        from pptx.chart.data import CategoryChartData  # type: ignore
        from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION  # type: ignore
    except ImportError:
        _missing("python-pptx")
    categories = [str(item) for item in (slide_spec.get("categories") or [])]
    chart_data = CategoryChartData()
    chart_data.categories = categories
    for series in _chart_series(slide_spec):
        chart_data.add_series(series["name"], [float(value) for value in series["values"]])
    chart_type_key = str(slide_spec.get("chart_type") or slide_spec.get("type") or "column").lower()
    chart_frame = slide.shapes.add_chart(
        _chart_type(chart_type_key, XL_CHART_TYPE),
        _inches(left),
        _inches(top),
        _inches(width),
        _inches(height),
        chart_data,
    )
    chart = chart_frame.chart
    chart.has_title = bool(slide_spec.get("chart_title"))
    if chart.has_title:
        chart.chart_title.text_frame.text = str(slide_spec["chart_title"])
    chart.has_legend = len(_chart_series(slide_spec)) > 1
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
    if hasattr(chart, "value_axis"):
        chart.value_axis.has_major_gridlines = True
    if slide_spec.get("data_labels"):
        try:
            plot = chart.plots[0]
            plot.has_data_labels = True
            plot.data_labels.show_value = True
        except Exception:
            pass


def _render_chart_slide(slide: Any, slide_spec: dict[str, Any], theme: dict[str, str], slide_w: float, slide_h: float, temp_paths: list[Path], workspace_root: Path) -> None:
    _apply_slide_background(slide, slide_spec, theme, slide_w, slide_h, temp_paths, workspace_root)
    _add_title(slide, slide_spec.get("title"), theme, slide_w)
    _add_chart(slide, slide_spec, theme, 0.8, 1.45, slide_w - 1.6, slide_h - 2.25)


def _text_lines_for_slide(slide_spec: dict[str, Any], layout_name: str) -> list[str]:
    if slide_spec.get("paragraph"):
        return [str(slide_spec["paragraph"])]
    if isinstance(slide_spec.get("bullets"), list):
        return [str(item) for item in slide_spec["bullets"]]
    if isinstance(slide_spec.get("items"), list):
        return [str(item) for item in slide_spec["items"]]
    if layout_name == "timeline":
        return [_item_title(item, f"Event {idx}") for idx, item in enumerate(_item_dicts(slide_spec.get("events") or slide_spec.get("items")), start=1)]
    if layout_name == "process":
        return [_item_title(item, f"Step {idx}") for idx, item in enumerate(_item_dicts(slide_spec.get("steps") or slide_spec.get("items")), start=1)]
    if layout_name == "stat":
        return [f"{_item_title(item)} {item.get('label', '')}".strip() for item in _item_dicts(slide_spec.get("stats"))]
    if layout_name == "comparison":
        return [_item_title(item, f"Option {idx}") for idx, item in enumerate(_comparison_columns(slide_spec), start=1)]
    if layout_name == "matrix":
        return [_item_title(item, f"Quadrant {idx}") for idx, item in enumerate(_item_dicts(slide_spec.get("quadrants") or slide_spec.get("items")), start=1)]
    if slide_spec.get("subtitle"):
        return [str(slide_spec["subtitle"])]
    return []


def _placeholder_kind(shape: Any) -> str:
    try:
        return str(shape.placeholder_format.type).lower()
    except Exception:
        return ""


def _set_placeholder_lines(shape: Any, lines: list[str]) -> None:
    tf = shape.text_frame
    tf.clear()
    for idx, line in enumerate(lines or [""]):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = str(line)


def _remove_empty_template_placeholders(slide: Any) -> None:
    for shape in list(slide.shapes):
        try:
            if not shape.is_placeholder or not getattr(shape, "has_text_frame", False):
                continue
            text = "\n".join(paragraph.text for paragraph in shape.text_frame.paragraphs).strip()
            if text:
                continue
            element = shape._element
            element.getparent().remove(element)
        except Exception:
            continue


def _render_template_bound_slide(slide: Any, slide_spec: dict[str, Any], layout_name: str, theme: dict[str, str], slide_w: float, slide_h: float) -> None:
    title = str(slide_spec.get("title") or "")
    subtitle = str(slide_spec.get("subtitle") or slide_spec.get("author") or "")
    body_lines = _text_lines_for_slide(slide_spec, layout_name)
    title_done = False
    subtitle_done = False
    body_done = False
    for shape in slide.shapes:
        try:
            if not shape.is_placeholder or not getattr(shape, "has_text_frame", False):
                continue
        except Exception:
            continue
        kind = _placeholder_kind(shape)
        if title and not title_done and ("title" in kind or "center_title" in kind or "ctrtitle" in kind):
            _set_placeholder_lines(shape, [title])
            title_done = True
        elif subtitle and not subtitle_done and "subtitle" in kind:
            _set_placeholder_lines(shape, [subtitle])
            subtitle_done = True
        elif body_lines and not body_done and any(token in kind for token in ["body", "object", "content", "placeholder"]):
            _set_placeholder_lines(shape, body_lines[:6])
            body_done = True
    if title and not title_done:
        _add_title(slide, title, theme, slide_w)
    if subtitle and not subtitle_done and layout_name in {"title", "section"}:
        _add_text(slide, subtitle, 0.75, 1.55, slide_w - 1.5, 0.5, theme, size=18, color_key="muted_text_color", align="center")
    if body_lines and not body_done and layout_name not in {"table", "chart"}:
        _add_bullets(slide, body_lines[:6], 0.85, 1.55, slide_w - 1.7, slide_h - 2.35, theme)
    if layout_name == "table":
        _add_table(slide, slide_spec.get("table"), 0.85, 1.65, slide_w - 1.7, slide_h - 2.45, theme)
    elif layout_name == "chart":
        _add_chart(slide, slide_spec, theme, 0.85, 1.65, slide_w - 1.7, slide_h - 2.45)
    _remove_empty_template_placeholders(slide)


def _link_items(slide_spec: dict[str, Any]) -> list[dict[str, str]]:
    result: list[dict[str, str]] = []
    for key in ("links", "citations"):
        value = slide_spec.get(key)
        if not isinstance(value, list):
            continue
        for item in value:
            if isinstance(item, dict):
                url = str(item.get("url") or item.get("href") or "").strip()
                label = str(item.get("label") or item.get("title") or url).strip()
            else:
                url = str(item).strip()
                label = url
            if url.startswith(("http://", "https://")):
                result.append({"label": label[:80], "url": url})
    return result[:4]


def _add_slide_links(slide: Any, slide_spec: dict[str, Any], theme: dict[str, str], slide_w: float, slide_h: float) -> None:
    links = _link_items(slide_spec)
    if not links:
        return
    try:
        from pptx.enum.text import PP_ALIGN  # type: ignore
    except ImportError:
        _missing("python-pptx")
    box = slide.shapes.add_textbox(_inches(0.65), _inches(slide_h - 0.55), _inches(slide_w - 2.1), _inches(0.35))
    tf = box.text_frame
    _clear_text_frame(tf)
    for idx, link in enumerate(links):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = link["label"]
        run.hyperlink.address = link["url"]
        run.font.name = theme.get("body_font") or theme["body_font"]
        run.font.size = _pt(7.5)
        run.font.color.rgb = _rgb(theme, "muted_text_color")


def _apply_notes(slide: Any, notes: Any) -> None:
    if not notes:
        return
    try:
        slide.notes_slide.notes_text_frame.text = str(notes)
    except Exception:
        pass


def _transition_dict(value: Any) -> dict[str, Any] | None:
    if not value:
        return None
    if isinstance(value, str):
        return {"type": value}
    if isinstance(value, dict):
        return value
    return None


def _transition_bool(value: Any, default: bool) -> str:
    if value is None:
        return "1" if default else "0"
    if isinstance(value, bool):
        return "1" if value else "0"
    if isinstance(value, str):
        return "0" if value.strip().lower() in {"0", "false", "no", "off"} else "1"
    return "1" if bool(value) else "0"


def _transition_speed(value: Any) -> str:
    speed = str(value or "med").strip().lower()
    if speed in {"slow", "slw"}:
        return "slow"
    if speed in {"fast", "fst"}:
        return "fast"
    return "med"


def _transition_dir(value: Any, default: str = "l") -> str:
    direction = str(value or default).strip().lower()
    return {
        "left": "l",
        "right": "r",
        "up": "u",
        "down": "d",
        "l": "l",
        "r": "r",
        "u": "u",
        "d": "d",
    }.get(direction, default)


def _transition_xml(spec: dict[str, Any]) -> str:
    transition_type = str(spec.get("type") or spec.get("effect") or "fade").strip().lower()
    if transition_type in {"none", "off"}:
        return ""
    if transition_type not in PPTX_TRANSITION_TYPES:
        _die(
            f"ERROR: unsupported slide transition '{transition_type}'. "
            f"Supported transitions: {', '.join(sorted(PPTX_TRANSITION_TYPES))}",
            3,
        )
    attrs = [
        f'spd="{_transition_speed(spec.get("speed"))}"',
        f'advClick="{_transition_bool(spec.get("advance_click"), True)}"',
    ]
    advance_ms = spec.get("advance_ms") or spec.get("advanceTime") or spec.get("advance_time")
    if advance_ms not in (None, ""):
        try:
            attrs.append(f'advTm="{max(0, int(advance_ms))}"')
        except (TypeError, ValueError):
            _die("ERROR: slide transition advance_ms must be an integer", 3)

    if transition_type == "fade":
        child = "<p:fade/>"
    elif transition_type in {"push", "wipe", "cover", "pull"}:
        child = f'<p:{transition_type} dir="{_transition_dir(spec.get("dir") or spec.get("direction"))}"/>'
    elif transition_type == "split":
        orient = str(spec.get("orient") or spec.get("orientation") or "vert").strip().lower()
        if orient not in {"horz", "vert"}:
            orient = "vert"
        split_dir = str(spec.get("dir") or spec.get("direction") or "out").strip().lower()
        if split_dir not in {"in", "out"}:
            split_dir = "out"
        child = f'<p:split orient="{orient}" dir="{split_dir}"/>'
    else:
        child = "<p:cut/>"

    return f'<p:transition {" ".join(attrs)} xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">{child}</p:transition>'


def _apply_slide_transition(slide: Any, transition: Any) -> None:
    spec = _transition_dict(transition)
    if not spec:
        return
    xml = _transition_xml(spec)
    if not xml:
        return
    try:
        from pptx.oxml import parse_xml  # type: ignore
        from pptx.oxml.ns import qn  # type: ignore
    except ImportError:
        _missing("python-pptx")

    slide_element = slide._element
    transition_tag = qn("p:transition")
    timing_tag = qn("p:timing")
    ext_tag = qn("p:extLst")
    for child in list(slide_element):
        if child.tag == transition_tag:
            slide_element.remove(child)

    insert_at = len(slide_element)
    for idx, child in enumerate(slide_element):
        if child.tag in {timing_tag, ext_tag}:
            insert_at = idx
            break
    slide_element.insert(insert_at, parse_xml(xml))


def _require_items(idx: int, slide_spec: dict[str, Any], layout_name: str, *keys: str) -> None:
    for key in keys:
        if _item_dicts(slide_spec.get(key)):
            return
    _die(f"ERROR: slide {idx} layout '{layout_name}' requires a non-empty {keys[0]} array", 3)


def _validate_chart_payload(idx: int, slide_spec: dict[str, Any]) -> None:
    categories = slide_spec.get("categories")
    if not isinstance(categories, list) or not categories:
        _die(f"ERROR: slide {idx} layout 'chart' requires non-empty categories array", 3)
    chart_type = str(slide_spec.get("chart_type") or slide_spec.get("type") or "column").lower()
    if chart_type not in PPTX_SUPPORTED_CHART_TYPES:
        _die(f"ERROR: slide {idx} uses unsupported chart_type '{chart_type}'. Supported chart types: {', '.join(sorted(PPTX_SUPPORTED_CHART_TYPES))}", 3)
    series = _chart_series(slide_spec)
    if not series:
        _die(f"ERROR: slide {idx} layout 'chart' requires non-empty series array", 3)
    for series_idx, item in enumerate(series, start=1):
        values = item.get("values")
        if not isinstance(values, list) or len(values) != len(categories):
            _die(f"ERROR: slide {idx} chart series {series_idx} values must match categories length", 3)
        for value in values:
            try:
                float(value)
            except (TypeError, ValueError):
                _die(f"ERROR: slide {idx} chart series {series_idx} contains a non-numeric value", 3)


def _validate_layout_payload(idx: int, layout_name: str, slide_spec: dict[str, Any]) -> None:
    if layout_name == "timeline":
        _require_items(idx, slide_spec, layout_name, "events", "items")
    elif layout_name == "process":
        _require_items(idx, slide_spec, layout_name, "steps", "items")
    elif layout_name == "comparison":
        if len(_comparison_columns(slide_spec)) < 2:
            _die(f"ERROR: slide {idx} layout 'comparison' requires left/right objects or at least two columns", 3)
    elif layout_name == "matrix":
        quadrants = _item_dicts(slide_spec.get("quadrants") or slide_spec.get("items"))
        if len(quadrants) != 4:
            _die(f"ERROR: slide {idx} layout 'matrix' requires exactly four quadrants", 3)
    elif layout_name == "chart":
        _validate_chart_payload(idx, slide_spec)


def _validate_spec(spec: dict[str, Any]) -> tuple[list[dict[str, Any]], list[Any]]:
    slides = spec.get("slides") or []
    if not isinstance(slides, list) or not slides:
        _die("ERROR: create_pptx spec requires non-empty 'slides' array", 3)
    for idx, slide_spec in enumerate(slides, start=1):
        if not isinstance(slide_spec, dict):
            _die(f"ERROR: slide {idx} must be an object", 3)
        layout_name = str(slide_spec.get("layout") or "body").lower()
        if layout_name not in PPTX_SUPPORTED_LAYOUTS:
            _die(f"ERROR: slide {idx} uses unsupported layout '{layout_name}'. Supported layouts: {', '.join(sorted(PPTX_SUPPORTED_LAYOUTS))}", 3)
        _validate_layout_payload(idx, layout_name, slide_spec)
    notes = spec.get("notes_per_slide") or []
    if notes and not isinstance(notes, list):
        _die("ERROR: notes_per_slide must be an array when provided", 3)
    if notes and len(notes) != len(slides):
        _die("ERROR: notes_per_slide length must match slides length", 3)
    return slides, notes


def create_pptx_from_spec(path: str, spec_path: str, template: str | None = None, workspace_root: str | Path | None = None) -> Path:
    try:
        from pptx import Presentation  # type: ignore
    except ImportError:
        _missing("python-pptx")

    root = _workspace_root(workspace_root)
    output_path = _validate_output_path(path, workspace_root=root)
    spec = _read_json(spec_path, workspace_root=root)
    slides, notes_per_slide = _validate_spec(spec)
    image_catalog = _normalize_image_catalog(spec.get("images"))
    icon_catalog = _normalize_icon_catalog(spec.get("icons"))
    slides = [_apply_icon_catalog_to_slide(_apply_image_catalog_to_slide(slide_spec, image_catalog), icon_catalog) for slide_spec in slides]
    theme = _normalize_theme(spec.get("theme"))

    prs = Presentation(str(_validate_path(template, workspace_root=root))) if template else Presentation()
    if template and not bool(spec.get("preserve_template_slides")):
        _remove_existing_slides(prs)
    _set_slide_size(prs, spec, bool(template))
    blank = _blank_layout(prs)
    slide_w, slide_h = _slide_size_inches(prs)
    temp_paths: list[Path] = []

    try:
        for index, slide_spec in enumerate(slides, start=1):
            layout_name = str(slide_spec.get("layout") or "body").lower()
            template_layout = _bound_template_layout(prs, slide_spec) if template else None
            slide = prs.slides.add_slide(template_layout or blank)
            if template_layout is not None:
                _render_template_bound_slide(slide, slide_spec, layout_name, theme, slide_w, slide_h)
            elif layout_name == "title":
                _render_title_slide(slide, slide_spec, theme, slide_w, slide_h, temp_paths, root)
            elif layout_name == "agenda":
                _render_agenda_slide(slide, slide_spec, theme, slide_w, slide_h, temp_paths, root)
            elif layout_name == "body":
                _render_body_slide(slide, slide_spec, theme, slide_w, slide_h, temp_paths, root)
            elif layout_name == "two_column":
                _render_two_column_slide(slide, slide_spec, theme, slide_w, slide_h, temp_paths, root)
            elif layout_name == "stat":
                _render_stat_slide(slide, slide_spec, theme, slide_w, slide_h, temp_paths, root)
            elif layout_name == "quote":
                _render_quote_slide(slide, slide_spec, theme, slide_w, slide_h, temp_paths, root)
            elif layout_name == "section":
                _render_section_slide(slide, slide_spec, theme, slide_w, slide_h, temp_paths, root)
            elif layout_name == "image_full":
                _render_image_full_slide(slide, slide_spec, theme, slide_w, slide_h, temp_paths, root)
            elif layout_name == "table":
                _apply_slide_background(slide, slide_spec, theme, slide_w, slide_h, temp_paths, root)
                _add_title(slide, slide_spec.get("title"), theme, slide_w)
                _add_table(slide, slide_spec.get("table"), 0.7, 1.35, slide_w - 1.4, slide_h - 2.0, theme)
            elif layout_name == "timeline":
                _render_timeline_slide(slide, slide_spec, theme, slide_w, slide_h, temp_paths, root)
            elif layout_name == "process":
                _render_process_slide(slide, slide_spec, theme, slide_w, slide_h, temp_paths, root)
            elif layout_name == "comparison":
                _render_comparison_slide(slide, slide_spec, theme, slide_w, slide_h, temp_paths, root)
            elif layout_name == "matrix":
                _render_matrix_slide(slide, slide_spec, theme, slide_w, slide_h, temp_paths, root)
            elif layout_name == "chart":
                _render_chart_slide(slide, slide_spec, theme, slide_w, slide_h, temp_paths, root)
            _add_slide_icon(slide, slide_spec, theme, slide_w, temp_paths, root)
            _add_slide_links(slide, slide_spec, theme, slide_w, slide_h)
            _add_footer(slide, slide_spec, theme, slide_w, slide_h, index, len(slides))
            _apply_notes(slide, slide_spec.get("notes") or (notes_per_slide[index - 1] if notes_per_slide else None))
            _apply_slide_transition(slide, slide_spec.get("transition"))
    finally:
        for temp_path in temp_paths:
            try:
                temp_path.unlink(missing_ok=True)
            except Exception:
                pass

    prs.save(str(output_path))
    return output_path


def main() -> int:
    parser = argparse.ArgumentParser(description="Create an editable PPTX from a JSON deck spec.")
    parser.add_argument("--path", required=True, help="Absolute output .pptx path")
    parser.add_argument("--spec", required=True, help="Absolute JSON deck spec path, or '-' to read JSON from stdin")
    parser.add_argument("--template", default=None, help="Optional absolute .pptx template path")
    args = parser.parse_args()

    output = create_pptx_from_spec(args.path, args.spec, args.template, Path.cwd())
    print(f"created PPTX: {output}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
