#!/usr/bin/env python3
"""HTML-first PPTX renderer.

This pipeline is intentionally separate from pptx_renderer.py. It creates a
reviewable HTML deck project, optionally captures browser screenshots, exports
a hybrid PPTX, maps simple transitions/entrance animations, and writes a QA
manifest. The goal is a high-fidelity visual workflow with an explicit
editability tradeoff.
"""

from __future__ import annotations

import argparse
import html
import json
import re
import shutil
import sys
from pathlib import Path
from typing import Any


DEFAULT_WIDTH_PX = 1280
DEFAULT_HEIGHT_PX = 720
DEFAULT_WIDTH_IN = 13.333333
DEFAULT_HEIGHT_IN = 7.5

SUPPORTED_MODES = {"hybrid", "raster", "native"}
SCREENSHOT_MODES = {"auto", "require", "skip"}
ANIMATION_EFFECTS = {
    "appear": {"presetID": 1, "filter": None},
    "fade": {"presetID": 10, "filter": "fade"},
    "wipe": {"presetID": 22, "filter": "wipe(up)"},
    "fly": {"presetID": 2, "filter": "fly"},
    "zoom": {"presetID": 16, "filter": "zoom"},
    "split": {"presetID": 25, "filter": "split(in)"},
}


def _die(message: str, code: int = 1) -> None:
    print(message, file=sys.stderr)
    raise SystemExit(code)


def _missing(pkg: str) -> None:
    print(f"MISSING_DEP: {pkg}", file=sys.stderr)
    print(f"Install with: python -m pip install {pkg}", file=sys.stderr)
    raise SystemExit(2)


def _read_json(path: str) -> dict[str, Any]:
    if path == "-":
        data = json.load(sys.stdin)
    else:
        with Path(path).expanduser().resolve().open("r", encoding="utf-8") as handle:
            data = json.load(handle)
    if not isinstance(data, dict):
        _die("ERROR: spec root must be a JSON object", 3)
    return data


def _write_json(path: Path, data: Any, *, pretty: bool = True) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(
        json.dumps(data, ensure_ascii=False, indent=2 if pretty else None),
        encoding="utf-8",
    )


def _as_list(value: Any) -> list[Any]:
    return value if isinstance(value, list) else []


def _as_dict(value: Any) -> dict[str, Any]:
    return value if isinstance(value, dict) else {}


def _as_str(value: Any, default: str = "") -> str:
    return value if isinstance(value, str) else default


def _num(value: Any, default: float = 0.0) -> float:
    try:
        if value is None or value == "":
            return default
        return float(value)
    except (TypeError, ValueError):
        return default


def _bool(value: Any, default: bool = False) -> bool:
    if value is None:
        return default
    if isinstance(value, bool):
        return value
    if isinstance(value, str):
        return value.strip().lower() not in {"0", "false", "no", "off"}
    return bool(value)


def _clean_hex(value: Any, default: str = "FFFFFF") -> str:
    raw = str(value or default).strip().lstrip("#")
    if len(raw) == 3:
        raw = "".join(ch * 2 for ch in raw)
    if re.fullmatch(r"[0-9a-fA-F]{6}", raw):
        return raw.upper()
    return default.upper()


def _rgb_tuple(value: Any, default: str = "FFFFFF") -> tuple[int, int, int]:
    clean = _clean_hex(value, default)
    return int(clean[0:2], 16), int(clean[2:4], 16), int(clean[4:6], 16)


def _css_color(value: Any, default: str = "FFFFFF") -> str:
    return f"#{_clean_hex(value, default)}"


def _slide_size(spec: dict[str, Any]) -> dict[str, float]:
    raw = spec.get("slide_size") or spec.get("size") or "ppt169"
    if isinstance(raw, dict):
        width_in = _num(raw.get("width_in") or raw.get("width"), DEFAULT_WIDTH_IN)
        height_in = _num(raw.get("height_in") or raw.get("height"), DEFAULT_HEIGHT_IN)
        width_px = int(_num(raw.get("width_px"), DEFAULT_WIDTH_PX))
        height_px = int(_num(raw.get("height_px"), DEFAULT_HEIGHT_PX))
    elif str(raw).lower() in {"ppt43", "4:3"}:
        width_in, height_in = 10.0, 7.5
        width_px, height_px = 1024, 768
    else:
        width_in, height_in = DEFAULT_WIDTH_IN, DEFAULT_HEIGHT_IN
        width_px, height_px = DEFAULT_WIDTH_PX, DEFAULT_HEIGHT_PX
    return {
        "width_in": width_in,
        "height_in": height_in,
        "width_px": width_px,
        "height_px": height_px,
    }


def _validate_spec(spec: dict[str, Any]) -> list[dict[str, Any]]:
    slides = spec.get("slides")
    if not isinstance(slides, list) or not slides:
        _die("ERROR: HTML deck spec requires a non-empty slides array", 3)
    normalized: list[dict[str, Any]] = []
    for index, slide in enumerate(slides, start=1):
        if not isinstance(slide, dict):
            _die(f"ERROR: slide {index} must be an object", 3)
        html_source = _as_str(slide.get("html"))
        elements = _as_list(slide.get("elements") or slide.get("native"))
        if not html_source and not elements:
            _die(f"ERROR: slide {index} needs html or elements/native content", 3)
        copy = dict(slide)
        copy.setdefault("id", f"slide_{index:02d}")
        copy.setdefault("title", f"Slide {index}")
        normalized.append(copy)
    return normalized


def _theme(spec: dict[str, Any]) -> dict[str, str]:
    theme = _as_dict(spec.get("theme"))
    return {
        "background": _clean_hex(theme.get("background") or theme.get("background_color"), "F8FAFC"),
        "surface": _clean_hex(theme.get("surface") or theme.get("surface_color"), "FFFFFF"),
        "primary": _clean_hex(theme.get("primary") or theme.get("primary_color"), "2563EB"),
        "accent": _clean_hex(theme.get("accent") or theme.get("accent_color"), "F97316"),
        "text": _clean_hex(theme.get("text") or theme.get("text_color"), "111827"),
        "muted_text": _clean_hex(theme.get("muted_text") or theme.get("muted_text_color"), "64748B"),
        "title_font": _as_str(theme.get("title_font"), "Aptos Display"),
        "body_font": _as_str(theme.get("body_font"), "Aptos"),
    }


def _base_css(size: dict[str, float], theme: dict[str, str]) -> str:
    return f"""
:root {{
  --bg: #{theme["background"]};
  --surface: #{theme["surface"]};
  --primary: #{theme["primary"]};
  --accent: #{theme["accent"]};
  --text: #{theme["text"]};
  --muted: #{theme["muted_text"]};
}}
* {{ box-sizing: border-box; }}
body {{
  margin: 0;
  background: #111827;
  font-family: "{theme["body_font"]}", "Aptos", "Inter", Arial, sans-serif;
  color: var(--text);
}}
.deck-stage {{
  width: {int(size["width_px"])}px;
  min-height: {int(size["height_px"])}px;
  margin: 0;
}}
.slide {{
  position: relative;
  overflow: hidden;
  width: {int(size["width_px"])}px;
  height: {int(size["height_px"])}px;
  background: var(--bg);
}}
.slide h1, .slide h2, .slide h3 {{
  font-family: "{theme["title_font"]}", "{theme["body_font"]}", Arial, sans-serif;
  margin: 0;
}}
.slide p {{ margin: 0; }}
"""


def _default_slide_html(slide: dict[str, Any], theme: dict[str, str]) -> str:
    title = html.escape(_as_str(slide.get("title"), "Untitled"))
    subtitle = html.escape(_as_str(slide.get("subtitle") or slide.get("summary")))
    bullets = [html.escape(str(item)) for item in _as_list(slide.get("bullets"))]
    bullet_html = "".join(f"<li>{item}</li>" for item in bullets)
    subtitle_html = f"<p class=\"subtitle\">{subtitle}</p>" if subtitle else ""
    bullets_html = f"<ul>{bullet_html}</ul>" if bullet_html else ""
    return f"""
<div class="fallback-layout">
  <div class="accent"></div>
  <h1>{title}</h1>
  {subtitle_html}
  {bullets_html}
</div>
<style>
.fallback-layout {{
  position: absolute;
  inset: 72px;
  display: grid;
  align-content: center;
  gap: 28px;
}}
.fallback-layout .accent {{
  width: 96px;
  height: 8px;
  border-radius: 999px;
  background: #{theme["accent"]};
}}
.fallback-layout h1 {{
  max-width: 860px;
  color: #{theme["text"]};
  font-size: 64px;
  line-height: 0.96;
  letter-spacing: 0;
}}
.fallback-layout .subtitle {{
  max-width: 720px;
  color: #{theme["muted_text"]};
  font-size: 28px;
  line-height: 1.25;
}}
.fallback-layout ul {{
  margin: 0;
  padding-left: 28px;
  color: #{theme["text"]};
  font-size: 25px;
  line-height: 1.3;
}}
</style>
"""


def _html_document(
    *,
    title: str,
    body: str,
    css: str,
    size: dict[str, float],
    theme: dict[str, str],
) -> str:
    return f"""<!doctype html>
<html lang="en">
<head>
  <meta charset="utf-8">
  <meta name="viewport" content="width={int(size["width_px"])}, initial-scale=1">
  <title>{html.escape(title)}</title>
  <style>
{_base_css(size, theme)}
{css}
  </style>
</head>
<body>
  <main class="deck-stage">
{body}
  </main>
</body>
</html>
"""


def _write_html_project(
    spec: dict[str, Any],
    slides: list[dict[str, Any]],
    out_dir: Path,
    size: dict[str, float],
    theme: dict[str, str],
) -> list[dict[str, Any]]:
    source_dir = out_dir / "source"
    slide_dir = source_dir / "slides"
    slide_dir.mkdir(parents=True, exist_ok=True)
    css = _as_str(spec.get("css"))
    records: list[dict[str, Any]] = []
    deck_sections: list[str] = []

    for index, slide in enumerate(slides, start=1):
        body = _as_str(slide.get("html")) or _default_slide_html(slide, theme)
        slide_css = css + "\n" + _as_str(slide.get("css"))
        wrapped = f'    <section class="slide" data-slide="{index}" id="{html.escape(str(slide["id"]))}">\n{body}\n    </section>'
        deck_sections.append(wrapped)
        slide_html = _html_document(
            title=f"{index:02d} {slide.get('title')}",
            body=wrapped,
            css=slide_css,
            size=size,
            theme=theme,
        )
        slide_path = slide_dir / f"slide_{index:02d}.html"
        slide_path.write_text(slide_html, encoding="utf-8")
        records.append(
            {
                "index": index,
                "id": slide["id"],
                "title": slide.get("title"),
                "htmlPath": str(slide_path),
            }
        )

    deck_html = _html_document(
        title=_as_str(spec.get("title"), "HTML Deck"),
        body="\n".join(deck_sections),
        css=css,
        size=size,
        theme=theme,
    )
    deck_path = source_dir / "deck.html"
    deck_path.write_text(deck_html, encoding="utf-8")
    for record in records:
        record["deckPath"] = str(deck_path)
    return records


def _capture_screenshots(
    html_records: list[dict[str, Any]],
    out_dir: Path,
    size: dict[str, float],
    screenshot_mode: str,
) -> dict[str, Any]:
    rendered_dir = out_dir / "rendered"
    rendered_dir.mkdir(parents=True, exist_ok=True)
    result: dict[str, Any] = {
        "mode": screenshot_mode,
        "available": False,
        "engine": None,
        "slides": [],
        "warnings": [],
    }
    if screenshot_mode == "skip":
        result["warnings"].append("browser screenshots skipped")
        return result

    try:
        from playwright.sync_api import sync_playwright  # type: ignore
    except ImportError:
        message = "python package playwright is not installed; screenshot stage skipped"
        if screenshot_mode == "require":
            _die(f"ERROR: {message}", 2)
        result["warnings"].append(message)
        return result

    try:
        with sync_playwright() as playwright:
            browser = playwright.chromium.launch()
            page = browser.new_page(
                viewport={"width": int(size["width_px"]), "height": int(size["height_px"])},
                device_scale_factor=1,
            )
            for record in html_records:
                html_path = Path(str(record["htmlPath"]))
                png_path = rendered_dir / f"slide_{int(record['index']):02d}.png"
                page.goto(html_path.as_uri(), wait_until="networkidle")
                page.screenshot(path=str(png_path), full_page=False)
                record["screenshotPath"] = str(png_path)
                result["slides"].append({"index": record["index"], "path": str(png_path)})
            browser.close()
    except Exception as exc:  # noqa: BLE001
        message = f"browser screenshot stage failed: {type(exc).__name__}: {exc}"
        if screenshot_mode == "require":
            _die(f"ERROR: {message}", 2)
        result["warnings"].append(message)
        return result

    result["available"] = True
    result["engine"] = "playwright.chromium"
    return result


def _inches(value: float):
    try:
        from pptx.util import Inches  # type: ignore
    except ImportError:
        _missing("python-pptx")
    return Inches(float(value))


def _pt(value: float):
    try:
        from pptx.util import Pt  # type: ignore
    except ImportError:
        _missing("python-pptx")
    return Pt(float(value))


def _rgb(value: Any, default: str = "FFFFFF"):
    try:
        from pptx.dml.color import RGBColor  # type: ignore
    except ImportError:
        _missing("python-pptx")
    return RGBColor(*_rgb_tuple(value, default))


def _blank_layout(prs: Any):
    for layout in prs.slide_layouts:
        if len(layout.placeholders) == 0:
            return layout
    return prs.slide_layouts[-1]


def _resolve_asset(raw: Any, workspace_root: Path) -> Path:
    value = _as_str(raw)
    if not value:
        _die("ERROR: image src/path is required", 3)
    path = Path(value)
    if not path.is_absolute():
        path = workspace_root / path
    resolved = path.resolve()
    try:
        resolved.relative_to(workspace_root)
    except ValueError:
        _die(f"ERROR: asset escapes workspace: {value}", 3)
    if not resolved.exists():
        _die(f"ERROR: asset not found: {resolved}", 3)
    return resolved


def _paragraph_alignment(value: Any):
    try:
        from pptx.enum.text import PP_ALIGN  # type: ignore
    except ImportError:
        _missing("python-pptx")
    align = str(value or "left").strip().lower()
    return {
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
        "justify": PP_ALIGN.JUSTIFY,
    }.get(align, PP_ALIGN.LEFT)


def _set_text_frame(shape: Any, element: dict[str, Any], theme: dict[str, str]) -> None:
    try:
        from pptx.enum.text import MSO_ANCHOR  # type: ignore
    except ImportError:
        _missing("python-pptx")
    text = str(element.get("text") or "")
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.margin_left = _inches(_num(element.get("padding_left"), 0.06))
    text_frame.margin_right = _inches(_num(element.get("padding_right"), 0.06))
    text_frame.margin_top = _inches(_num(element.get("padding_top"), 0.03))
    text_frame.margin_bottom = _inches(_num(element.get("padding_bottom"), 0.03))
    valign = str(element.get("valign") or "top").strip().lower()
    text_frame.vertical_anchor = {
        "middle": MSO_ANCHOR.MIDDLE,
        "center": MSO_ANCHOR.MIDDLE,
        "bottom": MSO_ANCHOR.BOTTOM,
    }.get(valign, MSO_ANCHOR.TOP)

    paragraphs = text.split("\n") or [""]
    for index, paragraph_text in enumerate(paragraphs):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.alignment = _paragraph_alignment(element.get("align"))
        paragraph.space_after = _pt(_num(element.get("space_after"), 0))
        run = paragraph.add_run()
        run.text = paragraph_text
        run.font.name = _as_str(element.get("font") or element.get("font_family"), theme["body_font"])
        run.font.size = _pt(_num(element.get("font_size") or element.get("fontSize"), 20))
        run.font.bold = _bool(element.get("bold"), False)
        run.font.italic = _bool(element.get("italic"), False)
        run.font.color.rgb = _rgb(element.get("color"), theme["text"])


def _apply_fill(shape: Any, fill: Any, default: str = "FFFFFF") -> None:
    if fill in (None, "", "none", "transparent"):
        shape.fill.background()
        return
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill, default)


def _apply_line(shape: Any, element: dict[str, Any]) -> None:
    line_color = element.get("line") or element.get("stroke")
    if line_color in (None, "", "none", "transparent"):
        shape.line.fill.background()
        return
    shape.line.color.rgb = _rgb(line_color, "111827")
    shape.line.width = _pt(_num(element.get("line_width") or element.get("stroke_width"), 1.0))


def _add_native_element(
    slide: Any,
    element: dict[str, Any],
    *,
    theme: dict[str, str],
    workspace_root: Path,
) -> tuple[str | None, int | None]:
    try:
        from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE  # type: ignore
    except ImportError:
        _missing("python-pptx")
    kind = str(element.get("type") or element.get("kind") or "text").strip().lower()
    element_id = _as_str(element.get("id") or element.get("name"))
    x = _num(element.get("x"), 0.0)
    y = _num(element.get("y"), 0.0)
    w = _num(element.get("w") or element.get("width"), 1.0)
    h = _num(element.get("h") or element.get("height"), 1.0)

    if kind in {"text", "textbox", "title"}:
        shape = slide.shapes.add_textbox(_inches(x), _inches(y), _inches(w), _inches(h))
        _set_text_frame(shape, element, theme)
    elif kind in {"rect", "rectangle", "roundrect", "card"}:
        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if kind in {"roundrect", "card"} or _num(element.get("radius"), 0) > 0 else MSO_SHAPE.RECTANGLE
        shape = slide.shapes.add_shape(shape_type, _inches(x), _inches(y), _inches(w), _inches(h))
        _apply_fill(shape, element.get("fill"), theme["surface"])
        _apply_line(shape, element)
    elif kind in {"ellipse", "oval", "circle"}:
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, _inches(x), _inches(y), _inches(w), _inches(h))
        _apply_fill(shape, element.get("fill"), theme["accent"])
        _apply_line(shape, element)
    elif kind in {"image", "picture"}:
        src = element.get("src") or element.get("path")
        path = _resolve_asset(src, workspace_root)
        shape = slide.shapes.add_picture(str(path), _inches(x), _inches(y), _inches(w), _inches(h))
    elif kind == "line":
        x2 = _num(element.get("x2"), x + w)
        y2 = _num(element.get("y2"), y + h)
        shape = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            _inches(x),
            _inches(y),
            _inches(x2),
            _inches(y2),
        )
        _apply_line(shape, {"line": element.get("color") or element.get("stroke"), "line_width": element.get("width") or element.get("line_width")})
    else:
        _die(f"ERROR: unsupported native element type: {kind}", 3)

    if element_id:
        try:
            shape.name = element_id
        except Exception:
            pass
    return element_id or None, int(getattr(shape, "shape_id", 0) or 0)


def _apply_notes(slide: Any, notes: Any) -> None:
    if not notes:
        return
    try:
        slide.notes_slide.notes_text_frame.text = str(notes)
    except Exception:
        pass


def _apply_transition(slide: Any, transition: Any) -> None:
    if not transition:
        return
    try:
        import pptx_renderer  # type: ignore
    except Exception:
        return
    try:
        pptx_renderer._apply_slide_transition(slide, transition)
    except Exception:
        return


def _animation_xml(targets: list[dict[str, Any]]) -> str:
    if not targets:
        return ""
    rows: list[str] = []
    next_id = 3
    for target in targets:
        shape_id = int(target["shape_id"])
        effect = str(target.get("effect") or "fade").lower()
        if effect not in ANIMATION_EFFECTS:
            effect = "fade"
        delay = max(0, int(_num(target.get("delay_ms"), 0)))
        duration = max(1, int(_num(target.get("duration_ms"), 450)))
        preset_id = ANIMATION_EFFECTS[effect]["presetID"]
        filter_name = ANIMATION_EFFECTS[effect]["filter"]
        wrapper_id = next_id
        set_id = next_id + 1
        effect_id = next_id + 2
        next_id += 3
        delay_attr = f' delay="{delay}"' if delay else ""
        effect_xml = ""
        if filter_name:
            effect_xml = f"""
                          <p:animEffect transition="in" filter="{filter_name}">
                            <p:cBhvr>
                              <p:cTn id="{effect_id}" dur="{duration}"/>
                              <p:tgtEl><p:spTgt spid="{shape_id}"/></p:tgtEl>
                            </p:cBhvr>
                          </p:animEffect>"""
        rows.append(
            f"""
                <p:par>
                  <p:cTn id="{wrapper_id}" fill="hold"{delay_attr}>
                    <p:childTnLst>
                      <p:set>
                        <p:cBhvr>
                          <p:cTn id="{set_id}" dur="1" fill="hold"/>
                          <p:tgtEl><p:spTgt spid="{shape_id}"/></p:tgtEl>
                          <p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>
                        </p:cBhvr>
                        <p:to><p:strVal val="visible"/></p:to>
                      </p:set>{effect_xml}
                    </p:childTnLst>
                  </p:cTn>
                </p:par>"""
        )
    return f"""<p:timing xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
  <p:tnLst>
    <p:par>
      <p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">
        <p:childTnLst>
          <p:seq concurrent="1" nextAc="seek">
            <p:cTn id="2" dur="indefinite" nodeType="mainSeq">
              <p:childTnLst>{''.join(rows)}
              </p:childTnLst>
            </p:cTn>
          </p:seq>
        </p:childTnLst>
      </p:cTn>
    </p:par>
  </p:tnLst>
</p:timing>"""


def _apply_animations(slide: Any, targets: list[dict[str, Any]]) -> None:
    xml = _animation_xml(targets)
    if not xml:
        return
    try:
        from pptx.oxml import parse_xml  # type: ignore
        from pptx.oxml.ns import qn  # type: ignore
    except ImportError:
        _missing("python-pptx")
    slide_element = slide._element
    timing_tag = qn("p:timing")
    ext_tag = qn("p:extLst")
    for child in list(slide_element):
        if child.tag == timing_tag:
            slide_element.remove(child)
    insert_at = len(slide_element)
    for idx, child in enumerate(slide_element):
        if child.tag == ext_tag:
            insert_at = idx
            break
    slide_element.insert(insert_at, parse_xml(xml))


def _collect_animation_targets(slide_spec: dict[str, Any], shape_ids: dict[str, int]) -> list[dict[str, Any]]:
    targets: list[dict[str, Any]] = []
    seen: set[str] = set()
    for element in _as_list(slide_spec.get("elements") or slide_spec.get("native")):
        if not isinstance(element, dict):
            continue
        element_id = _as_str(element.get("id") or element.get("name"))
        animation = element.get("animation")
        if element_id and animation and element_id in shape_ids:
            item = dict(animation if isinstance(animation, dict) else {"effect": animation})
            item["target"] = element_id
            item["shape_id"] = shape_ids[element_id]
            targets.append(item)
            seen.add(element_id)
    for animation in _as_list(slide_spec.get("animations")):
        if not isinstance(animation, dict):
            continue
        target = _as_str(animation.get("target") or animation.get("id"))
        if not target or target in seen or target not in shape_ids:
            continue
        item = dict(animation)
        item["shape_id"] = shape_ids[target]
        targets.append(item)
        seen.add(target)
    return targets


def _export_pptx(
    *,
    spec: dict[str, Any],
    slides: list[dict[str, Any]],
    html_records: list[dict[str, Any]],
    out_path: Path,
    mode: str,
    size: dict[str, float],
    theme: dict[str, str],
    workspace_root: Path,
) -> dict[str, Any]:
    try:
        from pptx import Presentation  # type: ignore
        from pptx.enum.shapes import MSO_SHAPE  # type: ignore
    except ImportError:
        _missing("python-pptx")

    prs = Presentation()
    prs.slide_width = _inches(size["width_in"])
    prs.slide_height = _inches(size["height_in"])
    blank = _blank_layout(prs)
    raster_slides = 0
    hybrid_slides = 0
    editable_slides = 0
    native_elements = 0
    animation_count = 0
    slide_outputs: list[dict[str, Any]] = []

    for index, slide_spec in enumerate(slides, start=1):
        record = html_records[index - 1]
        slide = prs.slides.add_slide(blank)
        bg = slide_spec.get("background") or theme["background"]
        bg_color = bg.get("color") if isinstance(bg, dict) else bg
        background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            _inches(0),
            _inches(0),
            _inches(size["width_in"]),
            _inches(size["height_in"]),
        )
        _apply_fill(background, bg_color, theme["background"])
        background.line.fill.background()

        screenshot_path = record.get("screenshotPath")
        elements = [item for item in _as_list(slide_spec.get("elements") or slide_spec.get("native")) if isinstance(item, dict)]
        use_raster = mode == "raster" or (mode == "hybrid" and screenshot_path and (not elements or _bool(slide_spec.get("raster_background"), False)))
        if use_raster:
            if not screenshot_path:
                _die(f"ERROR: slide {index} needs screenshot for raster export", 3)
            slide.shapes.add_picture(
                str(Path(str(screenshot_path)).resolve()),
                _inches(0),
                _inches(0),
                _inches(size["width_in"]),
                _inches(size["height_in"]),
            )
            raster_slides += 1

        shape_ids: dict[str, int] = {}
        if mode != "raster":
            for element in elements:
                element_id, shape_id = _add_native_element(
                    slide,
                    element,
                    theme=theme,
                    workspace_root=workspace_root,
                )
                native_elements += 1
                if element_id and shape_id:
                    shape_ids[element_id] = shape_id

        if use_raster and elements and mode != "raster":
            editability_class = "partial-raster-backplate"
            hybrid_slides += 1
        elif use_raster:
            editability_class = "raster"
        elif elements:
            editability_class = "native-editable"
            editable_slides += 1
        else:
            editability_class = "background-only"

        _apply_notes(slide, slide_spec.get("notes"))
        _apply_transition(slide, slide_spec.get("transition") or spec.get("transition"))
        animation_targets = _collect_animation_targets(slide_spec, shape_ids)
        if animation_targets:
            _apply_animations(slide, animation_targets)
            animation_count += len(animation_targets)

        slide_outputs.append(
            {
                "index": index,
                "mode": "raster" if use_raster else "native",
                "editabilityClass": editability_class,
                "nativeElements": len(elements) if mode != "raster" else 0,
                "animationTargets": len(animation_targets),
                "screenshotPath": screenshot_path,
            }
        )

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    total_slides = max(1, len(slides))
    conservative_editability = (
        editable_slides + hybrid_slides * 0.35
    ) / total_slides
    if raster_slides == 0 and editable_slides == len(slides):
        editability_level = "fully_editable"
    elif editable_slides or hybrid_slides:
        editability_level = "partially_editable"
    else:
        editability_level = "raster"
    return {
        "path": str(out_path),
        "mode": mode,
        "slides": slide_outputs,
        "metrics": {
            "slides": len(slides),
            "rasterSlides": raster_slides,
            "hybridSlides": hybrid_slides,
            "editableSlides": editable_slides,
            "nativeElements": native_elements,
            "animationTargets": animation_count,
            "editabilityScore": round(max(0.0, min(1.0, conservative_editability)), 3),
            "editabilityLevel": editability_level,
            "finalPptxRenderVerified": False,
        },
    }


def _luminance(hex_color: str) -> float:
    channels = [int(hex_color[i : i + 2], 16) / 255.0 for i in (0, 2, 4)]
    linear = [v / 12.92 if v <= 0.03928 else ((v + 0.055) / 1.055) ** 2.4 for v in channels]
    return 0.2126 * linear[0] + 0.7152 * linear[1] + 0.0722 * linear[2]


def _contrast(fg: str, bg: str) -> float:
    hi = max(_luminance(fg), _luminance(bg))
    lo = min(_luminance(fg), _luminance(bg))
    return (hi + 0.05) / (lo + 0.05)


def _evaluate_qa(
    *,
    spec: dict[str, Any],
    slides: list[dict[str, Any]],
    screenshot_result: dict[str, Any],
    pptx_result: dict[str, Any] | None,
    theme: dict[str, str],
) -> dict[str, Any]:
    checks: list[dict[str, Any]] = []
    issues: list[str] = []
    warnings = list(screenshot_result.get("warnings") or [])
    html_count = sum(1 for slide in slides if _as_str(slide.get("html")))
    native_count = sum(len(_as_list(slide.get("elements") or slide.get("native"))) for slide in slides)
    visual_anchor_count = 0
    low_contrast: list[int] = []
    dense: list[int] = []
    for index, slide in enumerate(slides, start=1):
        elements = [item for item in _as_list(slide.get("elements") or slide.get("native")) if isinstance(item, dict)]
        if _as_str(slide.get("html")) or any(str(el.get("type") or "").lower() in {"image", "picture", "rect", "rectangle", "ellipse", "circle", "card"} for el in elements):
            visual_anchor_count += 1
        text_chars = len(_as_str(slide.get("html"))) + sum(len(str(el.get("text") or "")) for el in elements)
        if text_chars > 1400:
            dense.append(index)
        background = slide.get("background") or theme["background"]
        bg_color = _clean_hex(background.get("color") if isinstance(background, dict) else background, theme["background"])
        for el in elements:
            if str(el.get("type") or "text").lower() not in {"text", "textbox", "title"}:
                continue
            color = _clean_hex(el.get("color"), theme["text"])
            if _contrast(color, bg_color) < 3.0:
                low_contrast.append(index)
                break

    checks.append({"name": "html_sources", "status": "pass" if html_count == len(slides) else "warn", "metric": html_count, "threshold": len(slides)})
    checks.append({"name": "native_manifest", "status": "pass" if native_count else "warn", "metric": native_count, "threshold": 1})
    checks.append({"name": "visual_anchors", "status": "pass" if visual_anchor_count == len(slides) else "fail", "metric": visual_anchor_count, "threshold": len(slides)})
    checks.append({"name": "browser_screenshots", "status": "pass" if screenshot_result.get("available") else "warn", "metric": len(screenshot_result.get("slides") or []), "threshold": len(slides)})
    if low_contrast:
        issues.append(f"low contrast native text on slides: {', '.join(map(str, low_contrast))}")
        checks.append({"name": "contrast", "status": "fail", "metric": len(low_contrast), "threshold": 0})
    else:
        checks.append({"name": "contrast", "status": "pass", "metric": 0, "threshold": 0})
    if dense:
        warnings.append(f"text-dense slides: {', '.join(map(str, dense))}")
        checks.append({"name": "text_density", "status": "warn", "metric": len(dense), "threshold": 0})
    else:
        checks.append({"name": "text_density", "status": "pass", "metric": 0, "threshold": 0})

    if pptx_result:
        editability = float(_as_dict(pptx_result.get("metrics")).get("editabilityScore") or 0)
        editability_level = str(_as_dict(pptx_result.get("metrics")).get("editabilityLevel") or "unknown")
        checks.append({"name": "editability_score", "status": "pass" if editability >= 0.8 else "warn", "metric": editability, "threshold": 0.8, "level": editability_level})
        checks.append({"name": "final_pptx_render", "status": "warn", "metric": False, "threshold": True, "detail": "HTML/browser screenshots are source previews, not final-PPTX render evidence."})
        warnings.append("final exported PPTX has not been rendered; browser screenshots do not prove PowerPoint layout fidelity")
        checks.append({"name": "animation_mapping", "status": "pass" if int(_as_dict(pptx_result.get("metrics")).get("animationTargets") or 0) else "warn", "metric": _as_dict(pptx_result.get("metrics")).get("animationTargets") or 0, "threshold": 1})

    status = "fail" if any(check["status"] == "fail" for check in checks) else "warn" if any(check["status"] == "warn" for check in checks) else "pass"
    return {
        "kind": "htmlDeckQa",
        "status": status,
        "checks": checks,
        "issues": issues,
        "warnings": warnings,
        "metrics": {
            "slides": len(slides),
            "htmlSlides": html_count,
            "nativeElements": native_count,
            "screenshots": len(screenshot_result.get("slides") or []),
        },
        "pptx": pptx_result,
        "screenshot": screenshot_result,
        "source": {
            "title": spec.get("title"),
            "format": spec.get("slide_size") or spec.get("size") or "ppt169",
        },
    }


def render_html_deck(
    *,
    spec_path: str,
    out_dir: str,
    pptx_path: str | None = None,
    mode: str = "hybrid",
    screenshot: str = "auto",
    pretty: bool = True,
    workspace_root: str | Path | None = None,
) -> dict[str, Any]:
    if mode not in SUPPORTED_MODES:
        _die(f"ERROR: mode must be one of {', '.join(sorted(SUPPORTED_MODES))}", 3)
    if screenshot not in SCREENSHOT_MODES:
        _die(f"ERROR: screenshot must be one of {', '.join(sorted(SCREENSHOT_MODES))}", 3)
    root = Path(workspace_root).resolve() if workspace_root else Path.cwd().resolve()
    output_dir = Path(out_dir).expanduser().resolve()
    try:
        output_dir.relative_to(root)
    except ValueError:
        _die(f"ERROR: out-dir escapes workspace: {output_dir}", 3)
    if output_dir == root:
        _die("ERROR: out-dir must be a managed subdirectory, not the workspace root", 3)
    if output_dir.exists():
        marker = output_dir / ".html-deck-project"
        managed_manifest = False
        manifest_path = output_dir / "manifest.json"
        if manifest_path.exists():
            try:
                managed_manifest = json.loads(manifest_path.read_text(encoding="utf-8")).get("kind") == "htmlDeckProject"
            except Exception:
                managed_manifest = False
        if any(output_dir.iterdir()) and not marker.exists() and not managed_manifest:
            _die(
                f"ERROR: out-dir already exists and is not an HTML deck project: {output_dir}",
                3,
            )
        shutil.rmtree(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)
    (output_dir / ".html-deck-project").write_text("managed by html_deck_renderer\n", encoding="utf-8")

    spec = _read_json(spec_path)
    slides = _validate_spec(spec)
    size = _slide_size(spec)
    theme = _theme(spec)
    html_records = _write_html_project(spec, slides, output_dir, size, theme)
    screenshot_result = _capture_screenshots(html_records, output_dir, size, screenshot)

    pptx_result = None
    if pptx_path:
        pptx_out = Path(pptx_path).expanduser().resolve()
        try:
            pptx_out.relative_to(root)
        except ValueError:
            _die(f"ERROR: pptx path escapes workspace: {pptx_out}", 3)
        pptx_result = _export_pptx(
            spec=spec,
            slides=slides,
            html_records=html_records,
            out_path=pptx_out,
            mode=mode,
            size=size,
            theme=theme,
            workspace_root=root,
        )

    manifest = {
        "kind": "htmlDeckProject",
        "version": 1,
        "title": spec.get("title"),
        "mode": mode,
        "size": size,
        "source": {
            "deckHtml": str(output_dir / "source" / "deck.html"),
            "slides": html_records,
        },
        "pptx": pptx_result,
        "screenshot": screenshot_result,
    }
    qa = _evaluate_qa(
        spec=spec,
        slides=slides,
        screenshot_result=screenshot_result,
        pptx_result=pptx_result,
        theme=theme,
    )
    _write_json(output_dir / "manifest.json", manifest, pretty=pretty)
    _write_json(output_dir / "qa.json", qa, pretty=pretty)
    return {
        "manifest": manifest,
        "qa": qa,
        "manifestPath": str(output_dir / "manifest.json"),
        "qaPath": str(output_dir / "qa.json"),
    }


def main() -> int:
    parser = argparse.ArgumentParser(description="Render an HTML-first deck project and optional PPTX.")
    parser.add_argument("--spec", required=True, help="JSON HTML deck spec path, or '-' for stdin")
    parser.add_argument("--out-dir", required=True, help="Output project directory under the workspace")
    parser.add_argument("--pptx", default=None, help="Optional output .pptx path under the workspace")
    parser.add_argument("--mode", choices=sorted(SUPPORTED_MODES), default="hybrid")
    parser.add_argument("--screenshot", choices=sorted(SCREENSHOT_MODES), default="auto")
    parser.add_argument("--compact", action="store_true", help="Write compact JSON output")
    args = parser.parse_args()

    result = render_html_deck(
        spec_path=args.spec,
        out_dir=args.out_dir,
        pptx_path=args.pptx,
        mode=args.mode,
        screenshot=args.screenshot,
        pretty=not args.compact,
        workspace_root=Path.cwd(),
    )
    print(json.dumps(result, ensure_ascii=False, indent=None if args.compact else 2))
    return 0 if result["qa"]["status"] != "fail" else 4


if __name__ == "__main__":
    raise SystemExit(main())
